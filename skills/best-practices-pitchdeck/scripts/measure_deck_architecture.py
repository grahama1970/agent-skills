"""Measure DECK-level architecture from a corpus of .pptx decks.

Page-level design is measured elsewhere (`$best-practices-slide-design`).
This module measures the layer above: where sections start, how long they run,
how decks open and close. It exists because that layer was once invented in
code instead of measured — the numbers here are reproducible so no future
agent has to guess.

Inputs: a directory of .pptx decks. Outputs: one DeckArchitecture per deck
(JSON or table). Failure modes: a deck that cannot be opened is reported as an
error row, never skipped silently; archetype classification is heuristic and
labelled as such in the output.
"""

from __future__ import annotations

import json
import re
from dataclasses import asdict, dataclass, field
from pathlib import Path

import typer
from loguru import logger

EMU_PER_INCH = 914400.0
app = typer.Typer(help="Measure deck-level architecture from real .pptx decks.")


@dataclass(frozen=True, slots=True)
class SlideFeatures:
    """Measured features of one slide (the basis for archetype heuristics)."""

    index: int
    words: int
    pictures: int
    text_blocks: int
    big_centred_text: bool
    title: str


@dataclass
class DeckArchitecture:
    """Deck-level shape: opening, closing, divider cadence, section lengths."""

    deck: str
    slides: int
    opening: list[str] = field(default_factory=list)
    closing: list[str] = field(default_factory=list)
    divider_positions: list[int] = field(default_factory=list)
    section_lengths: list[int] = field(default_factory=list)
    divider_titles: list[str] = field(default_factory=list)
    classification: str = "heuristic"


def _features(slide, width_in: float, height_in: float, index: int) -> SlideFeatures:
    words = pictures = text_blocks = 0
    big_centred = False
    title = ""
    for shape in slide.shapes:
        shape_type = str(shape.shape_type or "")
        if "PICTURE" in shape_type or "GROUP" in shape_type:
            pictures += 1
        if not (getattr(shape, "has_text_frame", False) and shape.has_text_frame):
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        text_blocks += 1
        words += len(text.split())
        try:
            y_fraction = shape.top / EMU_PER_INCH / height_in
            sizes = [run.font.size.pt for para in shape.text_frame.paragraphs
                     for run in para.runs if run.font.size]
            if 0.2 < y_fraction < 0.7 and max(sizes or [0]) >= 32:
                big_centred = True
            if not title and y_fraction < 0.2:
                title = " ".join(text.split())[:60]
        except (TypeError, AttributeError, ZeroDivisionError):
            logger.debug("unreadable geometry on slide {} of a deck", index)
    return SlideFeatures(index, words, pictures, text_blocks, big_centred, title)


def classify(features: SlideFeatures) -> str:
    """Heuristic archetype label. Deliberately coarse: deck SHAPE needs only
    cover / toc / divider / close / body distinctions."""
    lowered = features.title.lower()
    if features.index == 1:
        return "cover"
    if re.search(r"table of contents|agenda|outline", lowered):
        return "toc"
    if re.search(r"thank you|open discussion|questions|contact", lowered):
        return "close"
    if features.words <= 12 and features.big_centred_text:
        return "section-divider"
    if features.words <= 12:
        return "art-only"
    if features.words >= 100:
        return "dense-reference"
    if features.pictures >= 6:
        return "art-rich"
    return "content"


def measure_deck(path: Path) -> DeckArchitecture:
    from pptx import Presentation

    presentation = Presentation(str(path))
    width_in = presentation.slide_width / EMU_PER_INCH
    height_in = presentation.slide_height / EMU_PER_INCH
    slides = list(presentation.slides)
    labelled = [(classify(_features(s, width_in, height_in, i + 1)),
                 _features(s, width_in, height_in, i + 1).title)
                for i, s in enumerate(slides)]
    kinds = [kind for kind, _ in labelled]
    dividers = [i + 1 for i, kind in enumerate(kinds) if kind == "section-divider"]
    lengths = [dividers[i + 1] - dividers[i] for i in range(len(dividers) - 1)]
    return DeckArchitecture(
        deck=path.stem,
        slides=len(slides),
        opening=kinds[:4],
        closing=kinds[-3:],
        divider_positions=dividers,
        section_lengths=lengths,
        divider_titles=[title for kind, title in labelled if kind == "section-divider"][:8],
    )


@app.command()
def measure(
    corpus: Path = typer.Option(..., help="Directory containing .pptx decks."),
    json_out: bool = typer.Option(False, "--json", help="Emit JSON instead of a table."),
) -> None:
    """Measure every deck in a corpus directory."""
    decks = sorted(corpus.glob("*.pptx"))
    if not decks:
        raise typer.BadParameter(f"no .pptx decks under {corpus}")
    results: list[DeckArchitecture] = []
    for deck in decks:
        try:
            results.append(measure_deck(deck))
        except Exception as exc:  # a deck that cannot be read is reported, never skipped
            logger.error("unreadable deck {}: {}", deck.name, exc)
            results.append(DeckArchitecture(deck=deck.stem, slides=-1, classification=f"error: {exc}"))
    if json_out:
        typer.echo(json.dumps([asdict(r) for r in results], indent=1))
        return
    for r in results:
        median = sorted(r.section_lengths)[len(r.section_lengths) // 2] if r.section_lengths else None
        typer.echo(f"{r.deck}: {r.slides} slides | open {r.opening} | close {r.closing} "
                   f"| {len(r.divider_positions)} dividers | median section {median}")


if __name__ == "__main__":
    app()
