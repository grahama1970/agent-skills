from __future__ import annotations

import math
import tempfile
from pathlib import Path

import importlib.util
import shutil
import subprocess
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from .io import SkillError, dump_json, expand_path, sha256_file
from .models import (
    AssetManifest,
    AssetSpec,
    DeckManifest,
    OperationClaims,
    OperationReceipt,
    Readiness,
    SeamValidation,
    SlideLayout,
    SourceManifest,
    ClaimLedger,
    ValidationReport,
    Visibility,
    VisualType,
)
from .validation import validate_bundle, validate_pptx


class Theme:
    background = "08131F"
    panel = "102437"
    panel_alt = "132D43"
    text = "F4F8FB"
    muted = "AFC2CF"
    cyan = "36D6E7"
    teal = "2DD4BF"
    amber = "F5B942"
    red = "EF5B5B"
    violet = "8B5CF6"
    line = "2A455A"
    font = "Arial"
    mono = "Courier New"


SLIDE_W = 13.333
SLIDE_H = 7.5


def _rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor.from_string(value)


def _set_background(slide, color: str = Theme.background) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


def _add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 20,
    bold: bool = False,
    color: str = Theme.text,
    font: str = Theme.font,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.0,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    paragraph.space_before = Pt(0)
    paragraph.line_spacing = 1.05
    for run in paragraph.runs:
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = _rgb(color)
    return box


def _add_title(slide, title: str, *, section_label: str | None = None) -> None:
    if section_label:
        _add_text(
            slide,
            section_label.upper(),
            0.55,
            0.26,
            4.8,
            0.28,
            size=9,
            bold=True,
            color=Theme.cyan,
        )
    _add_text(slide, title, 0.55, 0.55, 12.1, 0.68, size=28, bold=True)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.31), Inches(1.05), Inches(0.035)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = _rgb(Theme.cyan)
    line.line.fill.background()


def _add_footer(slide, deck_title: str, slide_no: int, visibility: Visibility) -> None:
    label = f"{visibility.value.upper()}  ·  {deck_title}"
    _add_text(slide, label, 0.55, 7.12, 8.9, 0.18, size=7.5, color=Theme.muted)
    _add_text(
        slide,
        f"{slide_no:02d}",
        12.05,
        7.08,
        0.7,
        0.23,
        size=8.5,
        bold=True,
        color=Theme.cyan,
        align=PP_ALIGN.RIGHT,
    )


def _add_panel(slide, x: float, y: float, w: float, h: float, *, color: str = Theme.panel):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    shape.line.color.rgb = _rgb(Theme.line)
    shape.line.width = Pt(0.8)
    return shape


def _add_body_rows(slide, items: list[str], x: float, y: float, w: float, h: float) -> None:
    if not items:
        return
    row_h = h / len(items)
    body_size = (
        10.8
        if row_h < 0.72
        else 12.0
        if row_h < 0.9 or max(len(item) for item in items) > 110 or len(items) > 4
        else 15.5
    )
    for index, item in enumerate(items):
        row_y = y + index * row_h
        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(row_y + 0.09),
            Inches(0.08),
            Inches(max(0.28, row_h - 0.18)),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = _rgb(Theme.teal if index % 2 == 0 else Theme.cyan)
        bar.line.fill.background()
        _add_text(
            slide,
            item,
            x + 0.22,
            row_y,
            w - 0.22,
            row_h,
            size=body_size,
            color=Theme.text,
            valign=MSO_ANCHOR.MIDDLE,
        )


def _split_card_item(item: str) -> tuple[str, str]:
    for separator in (" — ", ": ", " - "):
        if separator in item:
            title, body = item.split(separator, 1)
            if 2 <= len(title) <= 52:
                return title.strip(), body.strip()
    words = item.split()
    title = " ".join(words[: min(4, len(words))])
    return title, item


def _materialize_image(path: Path, temp_dir: Path) -> Path:
    suffix = path.suffix.lower()
    if suffix == ".svg":
        target = temp_dir / f"{path.stem}.png"
        if importlib.util.find_spec("cairosvg") is not None:
            import cairosvg  # type: ignore

            cairosvg.svg2png(url=str(path), write_to=str(target), output_width=2400)
            return target
        converter = shutil.which("rsvg-convert")
        if converter:
            result = subprocess.run(
                [converter, "-w", "2400", "-o", str(target), str(path)],
                check=False,
                capture_output=True,
                text=True,
            )
            if result.returncode == 0 and target.exists():
                return target
            raise SkillError(f"rsvg-convert failed for {path}: {result.stderr}")
        raise SkillError(
            "SVG conversion requires the optional cairosvg Python package or rsvg-convert binary."
        )
    if suffix in {".webp", ".gif", ".tif", ".tiff", ".bmp"}:
        target = temp_dir / f"{path.stem}.png"
        with Image.open(path) as image:
            if getattr(image, "n_frames", 1) > 1:
                image.seek(0)
            converted = image.convert("RGBA")
            converted.save(target, format="PNG")
        return target
    return path


def _add_missing_asset(slide, x: float, y: float, w: float, h: float, asset: AssetSpec) -> None:
    panel = _add_panel(slide, x, y, w, h, color="0D1D2B")
    panel.line.color.rgb = _rgb(Theme.amber)
    panel.line.width = Pt(1.4)
    _add_text(
        slide,
        "MISSING ASSET",
        x + 0.35,
        y + 0.35,
        w - 0.7,
        0.35,
        size=11,
        bold=True,
        color=Theme.amber,
    )
    _add_text(
        slide,
        asset.alt_text,
        x + 0.35,
        y + 0.92,
        w - 0.7,
        min(1.3, h - 1.2),
        size=20,
        bold=True,
        color=Theme.text,
        valign=MSO_ANCHOR.MIDDLE,
    )
    if asset.generation_brief and h > 2.8:
        _add_text(
            slide,
            asset.generation_brief,
            x + 0.35,
            y + h - 1.15,
            w - 0.7,
            0.78,
            size=10,
            color=Theme.muted,
        )


def _add_image_fit(
    slide,
    asset: AssetSpec,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    asset_base_dir: Path,
    temp_dir: Path,
) -> bool:
    if not asset.local_path:
        _add_missing_asset(slide, x, y, w, h, asset)
        return False
    raw_path = expand_path(asset.local_path, base_dir=asset_base_dir)
    if not raw_path.exists():
        _add_missing_asset(slide, x, y, w, h, asset)
        return False
    try:
        path = _materialize_image(raw_path, temp_dir)
        with Image.open(path) as image:
            image_w, image_h = image.size
        image_ratio = image_w / image_h
        frame_ratio = w / h
        if image_ratio >= frame_ratio:
            draw_w = w
            draw_h = w / image_ratio
        else:
            draw_h = h
            draw_w = h * image_ratio
        draw_x = x + (w - draw_w) / 2
        draw_y = y + (h - draw_h) / 2
        _add_panel(slide, x, y, w, h, color="071019")
        picture = slide.shapes.add_picture(
            str(path), Inches(draw_x), Inches(draw_y), width=Inches(draw_w), height=Inches(draw_h)
        )
        picture.name = f"asset:{asset.id}"
        return True
    except Exception as exc:
        if asset.required:
            raise SkillError(f"Required asset '{asset.id}' could not be materialized: {exc}") from exc
        _add_missing_asset(slide, x, y, w, h, asset)
        return False


def _render_cover(slide, spec, asset_map, asset_base_dir: Path, temp_dir: Path) -> None:
    # Quiet visual depth without rasterizing the narrative.
    for x, y, w, h, color in (
        (8.7, -0.35, 5.1, 3.1, Theme.panel_alt),
        (9.65, 2.35, 4.1, 4.4, Theme.panel),
        (7.85, 5.4, 3.7, 2.5, "0D2436"),
    ):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(color)
        shape.line.fill.background()
        shape.rotation = -6
    _add_text(slide, "SPARTA EXPLORER", 0.7, 0.62, 5.1, 0.35, size=10, bold=True, color=Theme.cyan)
    _add_text(slide, spec.title, 0.7, 1.32, 7.5, 1.35, size=38, bold=True)
    _add_text(slide, spec.message, 0.72, 3.05, 7.0, 1.6, size=22, color=Theme.muted)
    _add_text(
        slide,
        "Inspectable evidence · explicit uncertainty · human authority",
        0.72,
        5.65,
        7.0,
        0.46,
        size=12,
        bold=True,
        color=Theme.teal,
    )
    if spec.visual.asset_id and spec.visual.asset_id in asset_map:
        _add_image_fit(
            slide,
            asset_map[spec.visual.asset_id],
            8.25,
            1.0,
            4.25,
            5.45,
            asset_base_dir=asset_base_dir,
            temp_dir=temp_dir,
        )


def _render_statement(slide, spec) -> None:
    _add_title(slide, spec.title, section_label=spec.role)
    _add_text(slide, spec.message, 0.75, 1.75, 11.6, 1.55, size=26, bold=True)
    if spec.body:
        _add_panel(slide, 0.7, 3.65, 11.95, 2.65)
        _add_body_rows(slide, spec.body[:4], 1.0, 3.9, 11.35, 2.15)
    else:
        _add_text(
            slide,
            "Source-bound draft. Add one decisive visual or three supporting points before external use.",
            0.78,
            4.4,
            10.4,
            0.8,
            size=15,
            color=Theme.muted,
        )


def _render_split(slide, spec, asset_map, asset_base_dir: Path, temp_dir: Path) -> None:
    _add_title(slide, spec.title, section_label=spec.role)
    very_long_message = len(spec.message) > 180
    long_message = len(spec.message) > 120
    message_h = 2.5 if very_long_message else 2.0 if long_message else 1.55
    message_size = 15.5 if very_long_message else 18.0 if long_message else 22
    _add_text(slide, spec.message, 0.7, 1.7, 5.55, message_h, size=message_size, bold=True)
    body_y = 1.7 + message_h + 0.18
    body_h = max(1.45, 6.35 - body_y)
    _add_body_rows(slide, spec.body[:4], 0.75, body_y, 5.45, body_h)
    if spec.visual.asset_id and spec.visual.asset_id in asset_map:
        _add_image_fit(
            slide,
            asset_map[spec.visual.asset_id],
            6.65,
            1.65,
            5.95,
            4.9,
            asset_base_dir=asset_base_dir,
            temp_dir=temp_dir,
        )
    else:
        _add_panel(slide, 6.65, 1.65, 5.95, 4.9)
        _add_text(slide, "ONE INSPECTABLE\nDECISION THREAD", 7.18, 2.25, 4.85, 1.1, size=24, bold=True)
        stages = spec.body[:4] or ["Source", "Relationship", "Evidence", "Human decision"]
        for index, stage in enumerate(stages[:4]):
            y = 3.65 + index * 0.62
            _add_text(slide, f"{index + 1:02d}", 7.18, y, 0.45, 0.28, size=9, bold=True, color=Theme.cyan)
            _add_text(slide, stage, 7.75, y - 0.05, 4.05, 0.42, size=13, color=Theme.text)


def _render_screenshot(slide, spec, asset_map, asset_base_dir: Path, temp_dir: Path) -> None:
    _add_title(slide, spec.title, section_label=spec.role)
    _add_text(slide, spec.message, 0.7, 1.45, 12.0, 0.62, size=16.5, color=Theme.muted)
    asset = asset_map.get(spec.visual.asset_id or "")
    if asset:
        _add_image_fit(
            slide,
            asset,
            0.72,
            2.12,
            11.9,
            4.35,
            asset_base_dir=asset_base_dir,
            temp_dir=temp_dir,
        )
    if spec.visual.callouts:
        count = min(len(spec.visual.callouts), 3)
        chip_w = 11.9 / count - 0.08
        for index, callout in enumerate(spec.visual.callouts[:count]):
            x = 0.72 + index * (chip_w + 0.12)
            _add_panel(slide, x, 6.57, chip_w, 0.42, color="0D2232")
            _add_text(slide, callout, x + 0.1, 6.62, chip_w - 0.2, 0.29, size=8.5, color=Theme.text)
    elif spec.visual.caption:
        _add_text(slide, spec.visual.caption, 0.75, 6.62, 11.75, 0.28, size=8.5, color=Theme.muted)


def _render_flow(slide, spec) -> None:
    _add_title(slide, spec.title, section_label=spec.role)
    _add_text(slide, spec.message, 0.72, 1.45, 11.95, 0.72, size=18.5, color=Theme.muted)
    items = spec.visual.items or spec.body
    items = items[:6]
    if len(items) < 2:
        _render_statement(slide, spec)
        return
    available_w = 11.9
    gap = 0.18
    box_w = (available_w - gap * (len(items) - 1)) / len(items)
    y = 3.0
    for index, item in enumerate(items):
        x = 0.72 + index * (box_w + gap)
        panel = _add_panel(slide, x, y, box_w, 1.95, color=Theme.panel_alt if index % 2 else Theme.panel)
        panel.line.color.rgb = _rgb(Theme.cyan if index == len(items) - 1 else Theme.line)
        _add_text(slide, f"{index + 1:02d}", x + 0.18, y + 0.16, 0.42, 0.25, size=9, bold=True, color=Theme.cyan)
        _add_text(
            slide,
            item,
            x + 0.18,
            y + 0.58,
            box_w - 0.36,
            1.08,
            size=12.5 if len(items) > 4 else 14,
            bold=True,
            valign=MSO_ANCHOR.MIDDLE,
        )
        if index < len(items) - 1:
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(x + box_w),
                Inches(y + 0.98),
                Inches(x + box_w + gap),
                Inches(y + 0.98),
            )
            connector.line.color.rgb = _rgb(Theme.teal)
            connector.line.width = Pt(1.4)
    _add_text(
        slide,
        "Each transition preserves source, scope, review state, and decision authority.",
        1.05,
        5.55,
        11.1,
        0.62,
        size=14,
        bold=True,
        color=Theme.teal,
        align=PP_ALIGN.CENTER,
    )


def _render_cards(slide, spec, *, proof: bool = False) -> None:
    _add_title(slide, spec.title, section_label=spec.role)
    _add_text(slide, spec.message, 0.72, 1.45, 11.95, 0.8, size=18.5, color=Theme.muted)
    items = spec.body[:4] or [spec.message]
    columns = 4 if proof and len(items) >= 4 else min(3, len(items))
    rows = math.ceil(len(items) / columns)
    gap = 0.18
    card_w = (11.9 - gap * (columns - 1)) / columns
    card_h = 3.7 / rows - 0.12
    for index, item in enumerate(items):
        row = index // columns
        col = index % columns
        x = 0.72 + col * (card_w + gap)
        y = 2.65 + row * (card_h + 0.18)
        panel = _add_panel(slide, x, y, card_w, card_h, color=Theme.panel_alt if index % 2 else Theme.panel)
        panel.line.color.rgb = _rgb(Theme.teal if proof else Theme.line)
        title, body = _split_card_item(item)
        _add_text(slide, f"{index + 1:02d}", x + 0.22, y + 0.18, 0.46, 0.25, size=8.5, bold=True, color=Theme.cyan)
        _add_text(slide, title, x + 0.22, y + 0.62, card_w - 0.44, 0.6, size=16, bold=True)
        if body != title:
            _add_text(slide, body, x + 0.22, y + 1.34, card_w - 0.44, card_h - 1.55, size=11.5, color=Theme.muted)


def _render_roadmap(slide, spec) -> None:
    _add_title(slide, spec.title, section_label=spec.role)
    _add_text(slide, spec.message, 0.72, 1.45, 11.95, 0.72, size=18.5, color=Theme.muted)
    labels = [
        ("WORKING / DEMONSTRATED", Theme.teal),
        ("IN INTEGRATION", Theme.amber),
        ("NOT YET CLOSED", Theme.red),
    ]
    items = spec.body or ["Working surface", "Integration boundary", "Open closure gate"]
    buckets = [[], [], []]
    for index, item in enumerate(items):
        buckets[index % 3].append(item)
    for column, (label, color) in enumerate(labels):
        x = 0.72 + column * 4.03
        panel = _add_panel(slide, x, 2.45, 3.75, 3.85)
        panel.line.color.rgb = _rgb(color)
        _add_text(slide, label, x + 0.24, 2.68, 3.25, 0.35, size=9.5, bold=True, color=color)
        y = 3.25
        for item in buckets[column][:3]:
            _add_text(slide, item, x + 0.24, y, 3.25, 0.78, size=12.5, color=Theme.text)
            y += 0.9


def _render_collaboration(slide, spec) -> None:
    _add_title(slide, spec.title, section_label=spec.role)
    _add_text(slide, spec.message, 0.72, 1.45, 11.95, 0.9, size=21, bold=True)
    items = spec.body[:3] or ["Workflow co-design", "Integration design", "Deployment planning"]
    card_w = 3.75
    for index, item in enumerate(items):
        x = 0.72 + index * 4.03
        panel = _add_panel(slide, x, 2.75, card_w, 2.45, color=Theme.panel_alt if index % 2 else Theme.panel)
        panel.line.color.rgb = _rgb(Theme.teal)
        title, body = _split_card_item(item)
        _add_text(slide, f"0{index + 1}", x + 0.24, 2.98, 0.5, 0.26, size=9, bold=True, color=Theme.cyan)
        _add_text(slide, title, x + 0.24, 3.45, 3.2, 0.62, size=18, bold=True)
        if body != title:
            _add_text(slide, body, x + 0.24, 4.15, 3.2, 0.75, size=11.5, color=Theme.muted)
    callout = _add_panel(slide, 0.72, 5.62, 11.84, 0.76, color="0C2B35")
    callout.line.color.rgb = _rgb(Theme.cyan)
    _add_text(
        slide,
        "Start with one representative claim. Trace it from source to relationship, evidence, uncertainty, and human decision.",
        1.0,
        5.82,
        11.3,
        0.38,
        size=13,
        bold=True,
        color=Theme.text,
        align=PP_ALIGN.CENTER,
    )


def _render_appendix(slide, spec) -> None:
    _add_title(slide, spec.title, section_label="PRIVATE APPENDIX")
    _add_text(slide, spec.message, 0.72, 1.45, 11.95, 0.72, size=17, color=Theme.muted)
    _add_panel(slide, 0.7, 2.35, 11.95, 4.25)
    _add_body_rows(slide, spec.body[:6], 1.0, 2.6, 11.35, 3.75)


def _render_slide(slide, spec, asset_map, asset_base_dir: Path, temp_dir: Path) -> None:
    if spec.layout == SlideLayout.COVER:
        _render_cover(slide, spec, asset_map, asset_base_dir, temp_dir)
    elif spec.layout == SlideLayout.STATEMENT:
        _render_statement(slide, spec)
    elif spec.layout == SlideLayout.SPLIT:
        _render_split(slide, spec, asset_map, asset_base_dir, temp_dir)
    elif spec.layout == SlideLayout.SCREENSHOT:
        _render_screenshot(slide, spec, asset_map, asset_base_dir, temp_dir)
    elif spec.layout == SlideLayout.FLOW:
        _render_flow(slide, spec)
    elif spec.layout == SlideLayout.THREE_CARDS:
        _render_cards(slide, spec, proof=False)
    elif spec.layout == SlideLayout.PROOF_CARDS:
        _render_cards(slide, spec, proof=True)
    elif spec.layout == SlideLayout.ROADMAP:
        _render_roadmap(slide, spec)
    elif spec.layout == SlideLayout.COLLABORATION:
        _render_collaboration(slide, spec)
    else:
        _render_appendix(slide, spec)


def _notes_text(spec) -> str:
    lines = [spec.notes.strip()] if spec.notes.strip() else []
    if spec.claim_ids:
        lines.append("Claim IDs: " + ", ".join(spec.claim_ids))
    if spec.source_refs:
        refs = []
        for ref in spec.source_refs:
            value = ref.source_id
            if ref.section:
                value += f" / {ref.section}"
            if ref.locator:
                value += f" / {ref.locator}"
            refs.append(value)
        lines.append("Source refs: " + " | ".join(refs))
    return "\n\n".join(lines)


def build_pptx(
    deck: DeckManifest,
    ledger: ClaimLedger,
    sources: SourceManifest,
    assets: AssetManifest,
    *,
    source_manifest_dir: Path,
    asset_manifest_dir: Path,
    output_path: Path,
    require_approved_claims: bool = False,
) -> tuple[OperationReceipt, ValidationReport]:
    report = validate_bundle(
        deck,
        ledger,
        sources,
        assets,
        source_manifest_dir=source_manifest_dir,
        asset_manifest_dir=asset_manifest_dir,
        require_approved_claims=require_approved_claims,
    )
    if report.errors:
        preview = "\n".join(
            f"- {issue.code}: {issue.message}" for issue in report.issues if issue.severity == "error"
        )
        raise SkillError(f"Deck validation failed closed:\n{preview}")

    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_W)
    presentation.slide_height = Inches(SLIDE_H)
    blank = presentation.slide_layouts[6]
    asset_map = {asset.id: asset for asset in assets.assets}
    output_path.parent.mkdir(parents=True, exist_ok=True)

    with tempfile.TemporaryDirectory(prefix="readme-to-pitchdeck-images-") as temp:
        temp_dir = Path(temp)
        for spec in sorted(deck.slides, key=lambda value: value.order):
            slide = presentation.slides.add_slide(blank)
            _set_background(slide)
            _render_slide(slide, spec, asset_map, asset_manifest_dir, temp_dir)
            _add_footer(slide, deck.deck.title, spec.order, spec.visibility)
            slide.notes_slide.notes_text_frame.text = _notes_text(spec)

    presentation.save(output_path)
    pptx_issues = validate_pptx(output_path, len(deck.slides))
    structural_errors = [issue for issue in pptx_issues if issue.severity == "error"]
    if structural_errors:
        preview = "\n".join(f"- {issue.code}: {issue.message}" for issue in structural_errors)
        raise SkillError(f"Generated PPTX failed structural verification:\n{preview}")

    all_warnings = [issue.message for issue in report.issues if issue.severity == "warning"]
    all_warnings.extend(issue.message for issue in pptx_issues if issue.severity == "warning")
    readiness = Readiness.USABLE_WITH_GAPS if all_warnings else Readiness.READY
    receipt = OperationReceipt(
        schema="readme_to_pitchdeck.build_receipt.v1",
        operation="build",
        readiness=readiness,
        mocked=False,
        live=False,
        inputs={
            "deck_id": deck.deck.id,
            "source_policy": deck.deck.source_policy.value,
        },
        outputs={
            "pptx": str(output_path.resolve()),
            "pptx_sha256": sha256_file(output_path),
        },
        counts={
            "slides": len(deck.slides),
            "claims_referenced": len({claim for slide in deck.slides for claim in slide.claim_ids}),
            "sources_referenced": len(
                {ref.source_id for slide in deck.slides for ref in slide.source_refs}
            ),
            "validation_warnings": len(all_warnings),
        },
        gaps=all_warnings,
        claims=OperationClaims(
            proves=[
                "The PPTX was generated from typed, producer-validated manifests.",
                "Narrative text and diagrams remain editable PowerPoint objects.",
                "The generated PPTX was reopened and structurally verified.",
                "Encoded public/private source and claim boundaries passed validation.",
            ],
            does_not_prove=[
                "The source README claims match the current implementation or runtime.",
                "Candidate claims have human approval unless require-approved-claims was used.",
                "The deck is visually approved after import into Google Slides.",
                "Screenshots or dated proof remain current at presentation time.",
                "The product is production-ready, certified, accredited, or deployed.",
            ],
        ),
        seam_validation=SeamValidation(kind="build_receipt"),
    )
    receipt_path = output_path.with_suffix(".build-receipt.json")
    dump_json(receipt, receipt_path)
    return receipt, report
