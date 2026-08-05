"""The reviewer's 14-case proof bundle (WebGPT review item 15) as the release gate.

Each case maps to the reviewer's executable list. Implemented protections
assert; not-yet-implemented ones are strict xfails so the gate documents the
gap honestly instead of green-washing it. This module is the CI release bar:
all-pass (with zero unexpected passes) == the boundary holds as specified.
"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from readme_to_pitchdeck.io import load_yaml
from readme_to_pitchdeck.models import (
    AssetManifest,
    BindingKind,
    Claim,
    ClaimApproval,
    ClaimGuard,
    ClaimKind,
    ClaimLedger,
    ClaimRisk,
    ClaimStatus,
    DeckManifest,
    DeckMeta,
    SlideSpec,
    SourceRef,
    TextBinding,
    Visibility,
    VisualSpec,
)
from readme_to_pitchdeck.planner import plan_bundle
from readme_to_pitchdeck.pptx_builder import build_pptx
from readme_to_pitchdeck.validation import validate_bundle

from test_skill import FIXTURE, _base_models  # shared fixtures


@pytest.fixture()
def planned(tmp_path: Path) -> Path:
    from readme_to_pitchdeck.models import SourceManifest

    source_path = FIXTURE / "source_manifest.yaml"
    src = load_yaml(source_path, SourceManifest)
    out = tmp_path / "planned"
    plan_bundle(src, source_manifest_path=source_path, output_dir=out, max_slides=10)
    return out


def _run(deck, ledger, sources, assets, tmp_path, publish=False):
    return validate_bundle(
        deck, ledger, sources, assets,
        source_manifest_dir=tmp_path, asset_manifest_dir=tmp_path,
        require_approved_claims=publish,
    )


def _slide(**kw):
    base = dict(
        id="s1", order=1, role="content", layout="statement",
        visibility=Visibility.PUBLIC, title="T", message="M",
        claim_ids=["public-claim"], visual=VisualSpec(), claim_guard=ClaimGuard(),
        source_refs=[SourceRef(source_id="public-source")],
    )
    base.update(kw)
    return SlideSpec(**base)


def _deck(slide, visibility=Visibility.PUBLIC, policy="public_only"):
    return DeckManifest(
        deck=DeckMeta(id="d", title="D", audience="R", visibility=visibility, source_policy=policy),
        slides=[slide],
    )


def test_case1_visual_items_forbidden_phrase_rejected(tmp_path: Path) -> None:
    sources, ledger, assets = _base_models(tmp_path)
    slide = _slide(visual=VisualSpec(type="native_diagram", items=["production-ready", "Ship it"]))
    report = _run(_deck(slide), ledger, sources, assets, tmp_path)
    assert any(i.code == "FORBIDDEN_UNQUALIFIED_CLAIM" and i.severity == "error" for i in report.issues)


def test_case2_new_slide_unbound_text_blocks_publish(tmp_path: Path) -> None:
    sources, ledger, assets = _base_models(tmp_path)
    slide = _slide(title="A substantive assertion with no binding")
    report = _run(_deck(slide), ledger, sources, assets, tmp_path, publish=True)
    assert any(i.code == "UNBOUND_TEXT" and i.severity == "error" for i in report.issues)


def test_case3_qualifier_in_body_five_cannot_disappear(tmp_path: Path) -> None:
    sources, ledger, assets = _base_models(tmp_path)
    slide = _slide(body=[f"line {i}" for i in range(5)])  # statement capacity is 4
    report = _run(_deck(slide), ledger, sources, assets, tmp_path)
    assert any(i.code == "TARGET_CONTENT_TRUNCATED" and i.severity == "error" for i in report.issues)


def test_case4_stale_freeform_qualifier_cleared_on_layout_switch(planned: Path, tmp_path: Path) -> None:
    from readme_to_pitchdeck.slide_edit import apply_slide_edit

    deck = load_yaml(planned / "deck.public.yaml", DeckManifest)
    sid = sorted(deck.slides, key=lambda s: s.order)[1].id
    apply_slide_edit(planned, tmp_path / "ui", slide_id=sid, field="layout", value="freeform")
    apply_slide_edit(planned, tmp_path / "ui", slide_id=sid, field="layout", value="statement")
    after = load_yaml(planned / "deck.public.yaml", DeckManifest)
    assert [s for s in after.slides if s.id == sid][0].elements == []


def test_case5_deck_classification_immutable(planned: Path, tmp_path: Path) -> None:
    from readme_to_pitchdeck.source_edit import apply_deck_source

    src = (planned / "deck.public.yaml").read_text()
    with pytest.raises(ValueError, match="immutable"):
        apply_deck_source(planned, tmp_path / "ui", source_yaml=src.replace("visibility: public", "visibility: private", 1))


def test_case6_publish_blocks_candidates_draft_watermarks(planned: Path) -> None:
    deck = load_yaml(planned / "deck.public.yaml", DeckManifest)
    ledger = load_yaml(planned / "claim_ledger.yaml", ClaimLedger)
    from readme_to_pitchdeck.models import SourceManifest

    sources = load_yaml(planned / "source_manifest.resolved.yaml", SourceManifest)
    assets = load_yaml(planned / "asset_manifest.yaml", AssetManifest)
    with pytest.raises(Exception):
        build_pptx(deck, ledger, sources, assets, source_manifest_dir=planned,
                   asset_manifest_dir=planned, output_path=planned / "p.pptx", require_approved_claims=True)
    build_pptx(deck, ledger, sources, assets, source_manifest_dir=planned,
               asset_manifest_dir=planned, output_path=planned / "d.pptx", draft_watermark=True)
    texts = [sh.text_frame.text for s in Presentation(planned / "d.pptx").slides for sh in s.shapes if sh.has_text_frame]
    assert any("DRAFT — UNAPPROVED CLAIMS" in t for t in texts)


def test_case7_every_rendered_field_scanned(tmp_path: Path) -> None:
    sources, ledger, assets = _base_models(tmp_path)
    for field, kw in [
        ("footer", {"footer": "totally production-ready"}),
        ("caption", {"visual": VisualSpec(type="cards", items=["x", "y"], caption="production-ready caption")}),
    ]:
        report = _run(_deck(_slide(**kw)), ledger, sources, assets, tmp_path)
        assert any(i.code == "FORBIDDEN_UNQUALIFIED_CLAIM" for i in report.issues), field


def test_case8_unrelated_negation_cannot_publish_unbound_assertion(tmp_path: Path) -> None:
    # The negation-window heuristic would accept this sentence; the publish
    # coverage gate still blocks it because the string is unbound.
    sources, ledger, assets = _base_models(tmp_path)
    slide = _slide(message="No blockers are discussed. The system is production-ready.")
    report = _run(_deck(slide), ledger, sources, assets, tmp_path, publish=True)
    assert any(i.code == "UNBOUND_TEXT" and i.severity == "error" for i in report.issues)


def test_case9_claim_level_source_authority(tmp_path: Path) -> None:
    sources, ledger, assets = _base_models(tmp_path)
    bad_claim = Claim(
        id="bad-src", text="Claim with unknown source", kind=ClaimKind.PRODUCT,
        visibility=Visibility.PUBLIC, source_refs=[SourceRef(source_id="does-not-exist")],
        risk=ClaimRisk.MEDIUM, status=ClaimStatus.APPROVED,
        approval=ClaimApproval(approved_by="t", approved_at="2026-08-05"),
    )
    ledger2 = ledger.model_copy(update={"claims": [*ledger.claims, bad_claim]})
    slide = _slide(claim_ids=["bad-src"])
    report = _run(_deck(slide), ledger2, sources, assets, tmp_path)
    assert any(i.code == "CLAIM_UNKNOWN_SOURCE" for i in report.issues)
    # Private-source claim on a public deck.
    priv_claim = Claim(
        id="priv-src", text="Backed by a private source", kind=ClaimKind.PRODUCT,
        visibility=Visibility.PUBLIC, source_refs=[SourceRef(source_id="private-source")],
        risk=ClaimRisk.MEDIUM, status=ClaimStatus.APPROVED,
        approval=ClaimApproval(approved_by="t", approved_at="2026-08-05"),
    )
    ledger3 = ledger.model_copy(update={"claims": [*ledger.claims, priv_claim]})
    report2 = _run(_deck(_slide(claim_ids=["priv-src"])), ledger3, sources, assets, tmp_path)
    assert any(i.code == "CLAIM_PRIVATE_SOURCE" for i in report2.issues)
    # Expired approval blocks.
    stale = priv_claim.model_copy(update={
        "id": "stale", "source_refs": [SourceRef(source_id="public-source")],
        "approval": ClaimApproval(approved_by="t", approved_at="2025-01-01", expires_at="2025-06-01"),
    })
    ledger4 = ledger.model_copy(update={"claims": [*ledger.claims, stale]})
    report3 = _run(_deck(_slide(claim_ids=["stale"])), ledger4, sources, assets, tmp_path)
    assert any(i.code == "APPROVAL_EXPIRED" for i in report3.issues)


def test_case10_cleared_asset_not_served(planned: Path, tmp_path: Path) -> None:
    from PIL import Image

    from readme_to_pitchdeck.asset_ops import add_asset_to_slide, clear_slide_visual

    deck = load_yaml(planned / "deck.public.yaml", DeckManifest)
    sid = sorted(deck.slides, key=lambda s: s.order)[1].id
    img = tmp_path / "x.png"
    Image.new("RGB", (64, 64), "red").save(img)
    out = tmp_path / "ui"
    add_asset_to_slide(planned, out, slide_id=sid, file_path=img, alt_text="x")
    served = list((out / "assets").glob("*"))
    assert served
    clear_slide_visual(planned, out, slide_id=sid)
    assert not list((out / "assets").glob("*"))  # stale copies purged on re-emit


def test_case11_concurrent_edit_conflict(planned: Path, tmp_path: Path) -> None:
    from readme_to_pitchdeck.revisions import RevisionConflict, current_revision
    from readme_to_pitchdeck.slide_edit import apply_slide_edit

    deck = load_yaml(planned / "deck.public.yaml", DeckManifest)
    sid = sorted(deck.slides, key=lambda s: s.order)[1].id
    base = current_revision(planned)
    apply_slide_edit(planned, tmp_path / "ui", slide_id=sid, field="footer", value="a", expected_revision=base)
    with pytest.raises(RevisionConflict):
        apply_slide_edit(planned, tmp_path / "ui", slide_id=sid, field="footer", value="b", expected_revision=base)


def test_case12_pptx_text_inventory_matches_plan(planned: Path) -> None:
    from readme_to_pitchdeck.artifact_scan import scan_artifact
    from readme_to_pitchdeck.models import SourceManifest

    deck = load_yaml(planned / "deck.public.yaml", DeckManifest)
    ledger = load_yaml(planned / "claim_ledger.yaml", ClaimLedger)
    sources = load_yaml(planned / "source_manifest.resolved.yaml", SourceManifest)
    assets = load_yaml(planned / "asset_manifest.yaml", AssetManifest)
    out = planned / "inv.pptx"
    build_pptx(deck, ledger, sources, assets, source_manifest_dir=planned,
               asset_manifest_dir=planned, output_path=out)
    counts = scan_artifact(out, deck, ledger, sources)
    assert counts["strings_verified"] > 10  # full whole-string inventory ran


@pytest.mark.xfail(strict=True, reason="visual-diff threshold between browser and LibreOffice render not implemented (roundtable backlog)")
def test_case13_browser_vs_libreoffice_visual_diff() -> None:
    raise NotImplementedError


@pytest.mark.xfail(strict=True, reason="magic-byte verification and size limits on asset intake not implemented (P1-7 full)")
def test_case14_malformed_asset_upload_magic_bytes(tmp_path: Path) -> None:
    from readme_to_pitchdeck.asset_ops import add_asset_to_slide

    fake = tmp_path / "fake.png"
    fake.write_bytes(b"#!/bin/sh\necho not-an-image\n")  # wrong magic bytes, right suffix
    # Should be rejected on content, not suffix; today it passes intake.
    with pytest.raises(ValueError, match="magic"):
        add_asset_to_slide(tmp_path, tmp_path / "ui", slide_id="s1", file_path=fake, alt_text="x")


def test_undo_restores_previous_state_and_redo(planned: Path, tmp_path: Path) -> None:
    from readme_to_pitchdeck.revisions import NoHistory, undo_last_write
    from readme_to_pitchdeck.slide_edit import apply_slide_edit

    deck = load_yaml(planned / "deck.public.yaml", DeckManifest)
    sid = sorted(deck.slides, key=lambda s: s.order)[1].id
    original = [s for s in deck.slides if s.id == sid][0].footer
    apply_slide_edit(planned, tmp_path / "ui", slide_id=sid, field="footer", value="edited footer")
    edited = load_yaml(planned / "deck.public.yaml", DeckManifest)
    assert [s for s in edited.slides if s.id == sid][0].footer == "edited footer"

    undo_last_write(planned)
    restored = load_yaml(planned / "deck.public.yaml", DeckManifest)
    assert [s for s in restored.slides if s.id == sid][0].footer == original

    undo_last_write(planned)  # undo of the undo = redo
    redone = load_yaml(planned / "deck.public.yaml", DeckManifest)
    assert [s for s in redone.slides if s.id == sid][0].footer == "edited footer"

    # History can never be exhausted by undoing (each undo archives the
    # pre-undo state, so undo/redo ping-pongs by design). NoHistory applies
    # only to a bundle that has never been edited.
    pristine = tmp_path / "never-edited"
    pristine.mkdir()
    with pytest.raises(NoHistory):
        undo_last_write(pristine)
