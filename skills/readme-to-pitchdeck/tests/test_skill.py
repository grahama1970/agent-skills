from __future__ import annotations

from pathlib import Path

import pytest
import yaml
from pptx import Presentation

from readme_to_pitchdeck.io import SkillError, load_yaml
from readme_to_pitchdeck.models import (
    AssetKind,
    AssetManifest,
    AssetSpec,
    AssetStatus,
    Claim,
    ClaimKind,
    ClaimLedger,
    ClaimRisk,
    ClaimStatus,
    DeckManifest,
    DeckMeta,
    DeckSourcePolicy,
    SlideLayout,
    SlideSpec,
    SourceManifest,
    SourceRef,
    Visibility,
    VisualSpec,
    VisualType,
)
from readme_to_pitchdeck.planner import plan_bundle
from readme_to_pitchdeck.pptx_builder import build_pptx
from readme_to_pitchdeck.validation import validate_bundle


ROOT = Path(__file__).resolve().parents[1]
FIXTURE = ROOT / "fixtures" / "minimal"


def test_positive_plan_build_and_reopen(tmp_path: Path) -> None:
    source_path = FIXTURE / "source_manifest.yaml"
    source = load_yaml(source_path, SourceManifest)
    planned = tmp_path / "planned"
    receipt = plan_bundle(
        source,
        source_manifest_path=source_path,
        output_dir=planned,
        max_slides=10,
    )
    assert receipt.counts["slides"] >= 6
    assert receipt.counts["candidate_claims"] > 0

    deck = load_yaml(planned / "deck.public.yaml", DeckManifest)
    ledger = load_yaml(planned / "claim_ledger.yaml", ClaimLedger)
    sources = load_yaml(planned / "source_manifest.resolved.yaml", SourceManifest)
    assets = load_yaml(planned / "asset_manifest.yaml", AssetManifest)
    output = planned / "atlas.pptx"
    build_receipt, validation = build_pptx(
        deck,
        ledger,
        sources,
        assets,
        source_manifest_dir=planned,
        asset_manifest_dir=planned,
        output_path=output,
        require_approved_claims=False,
    )
    assert validation.errors == 0
    assert output.exists()
    assert build_receipt.counts["slides"] == len(deck.slides)
    reopened = Presentation(output)
    assert len(reopened.slides) == len(deck.slides)
    assert all(any(shape.has_text_frame and shape.text.strip() for shape in slide.shapes) for slide in reopened.slides)


def _base_models(tmp_path: Path):
    readme = tmp_path / "README.md"
    readme.write_text("# Demo\n\nSource.\n", encoding="utf-8")
    source_data = {
        "schema": "readme_to_pitchdeck.source_manifest.v1",
        "project_name": "Demo",
        "deck_title": "Demo",
        "audience": "Reviewers",
        "sources": [
            {
                "id": "public-source",
                "title": "Public",
                "path": str(readme),
                "visibility": "public",
                "role": "primary",
                "required": True,
            },
            {
                "id": "private-source",
                "title": "Private",
                "path": str(readme),
                "visibility": "private",
                "role": "evidence",
                "required": True,
            },
        ],
        "policy": {
            "public_deck_source_ids": ["public-source"],
            "forbidden_unqualified_claims": ["production-ready"],
            "mandatory_non_claims": [],
        },
        "seam_validation": {"kind": "source_manifest", "status": "PASS"},
    }
    sources = SourceManifest.model_validate(source_data)
    ledger = ClaimLedger(
        project_name="Demo",
        claims=[
            Claim(
                id="public-claim",
                text="A bounded public claim.",
                kind=ClaimKind.PRODUCT,
                visibility=Visibility.PUBLIC,
                source_refs=[SourceRef(source_id="public-source", section="Overview")],
                risk=ClaimRisk.MEDIUM,
                status=ClaimStatus.APPROVED,
            ),
            Claim(
                id="private-claim",
                text="A private implementation claim.",
                kind=ClaimKind.STATUS,
                visibility=Visibility.PRIVATE,
                source_refs=[SourceRef(source_id="private-source", section="Status")],
                risk=ClaimRisk.HIGH,
                status=ClaimStatus.APPROVED,
                required_qualifier="Verify current private evidence.",
            ),
        ],
    )
    assets = AssetManifest(project_name="Demo", assets=[])
    return sources, ledger, assets


def test_public_private_source_and_claim_leak_fails(tmp_path: Path) -> None:
    sources, ledger, assets = _base_models(tmp_path)
    deck = DeckManifest(
        deck=DeckMeta(
            id="demo-public",
            title="Demo",
            audience="Reviewers",
            visibility=Visibility.PUBLIC,
            source_policy=DeckSourcePolicy.PUBLIC_ONLY,
        ),
        slides=[
            SlideSpec(
                id="01-leak",
                order=1,
                role="test",
                layout=SlideLayout.STATEMENT,
                visibility=Visibility.PUBLIC,
                title="Leak",
                message="Private implementation claim.",
                source_refs=[SourceRef(source_id="private-source", section="Status")],
                claim_ids=["private-claim"],
                notes="Verify current private evidence.",
            )
        ],
    )
    report = validate_bundle(
        deck,
        ledger,
        sources,
        assets,
        source_manifest_dir=tmp_path,
        asset_manifest_dir=tmp_path,
    )
    codes = {issue.code for issue in report.issues}
    assert "PUBLIC_DECK_PRIVATE_SOURCE" in codes
    assert "PUBLIC_DECK_PRIVATE_CLAIM" in codes
    assert report.errors >= 2


def test_forbidden_unqualified_claim_fails(tmp_path: Path) -> None:
    sources, ledger, assets = _base_models(tmp_path)
    deck = DeckManifest(
        deck=DeckMeta(
            id="demo-public",
            title="Demo",
            audience="Reviewers",
            visibility=Visibility.PUBLIC,
            source_policy=DeckSourcePolicy.PUBLIC_ONLY,
        ),
        slides=[
            SlideSpec(
                id="01-overclaim",
                order=1,
                role="test",
                layout=SlideLayout.STATEMENT,
                visibility=Visibility.PUBLIC,
                title="Status",
                message="The system is production-ready.",
                source_refs=[SourceRef(source_id="public-source", section="Overview")],
                claim_ids=["public-claim"],
            )
        ],
    )
    report = validate_bundle(
        deck,
        ledger,
        sources,
        assets,
        source_manifest_dir=tmp_path,
        asset_manifest_dir=tmp_path,
    )
    assert any(issue.code == "FORBIDDEN_UNQUALIFIED_CLAIM" for issue in report.issues)


def test_negated_forbidden_phrase_is_allowed(tmp_path: Path) -> None:
    sources, ledger, assets = _base_models(tmp_path)
    deck = DeckManifest(
        deck=DeckMeta(
            id="demo-public",
            title="Demo",
            audience="Reviewers",
            visibility=Visibility.PUBLIC,
            source_policy=DeckSourcePolicy.PUBLIC_ONLY,
        ),
        slides=[
            SlideSpec(
                id="01-honest",
                order=1,
                role="test",
                layout=SlideLayout.STATEMENT,
                visibility=Visibility.PUBLIC,
                title="Status",
                message="The system is not production-ready.",
                source_refs=[SourceRef(source_id="public-source", section="Overview")],
                claim_ids=["public-claim"],
            )
        ],
    )
    report = validate_bundle(
        deck,
        ledger,
        sources,
        assets,
        source_manifest_dir=tmp_path,
        asset_manifest_dir=tmp_path,
    )
    assert not any(issue.code == "FORBIDDEN_UNQUALIFIED_CLAIM" for issue in report.issues)


def test_required_missing_asset_fails(tmp_path: Path) -> None:
    sources, ledger, _ = _base_models(tmp_path)
    assets = AssetManifest(
        project_name="Demo",
        assets=[
            AssetSpec(
                id="missing-shot",
                kind=AssetKind.SCREENSHOT,
                visibility=Visibility.PUBLIC,
                local_path=str(tmp_path / "missing.png"),
                required=True,
                alt_text="Required screenshot",
                status=AssetStatus.MISSING,
            )
        ],
    )
    deck = DeckManifest(
        deck=DeckMeta(
            id="demo-public",
            title="Demo",
            audience="Reviewers",
            visibility=Visibility.PUBLIC,
            source_policy=DeckSourcePolicy.PUBLIC_ONLY,
        ),
        slides=[
            SlideSpec(
                id="01-shot",
                order=1,
                role="test",
                layout=SlideLayout.SCREENSHOT,
                visibility=Visibility.PUBLIC,
                title="Screenshot",
                message="Inspect the real surface.",
                source_refs=[SourceRef(source_id="public-source", section="Overview")],
                claim_ids=["public-claim"],
                visual=VisualSpec(type=VisualType.SCREENSHOT, asset_id="missing-shot"),
            )
        ],
    )
    report = validate_bundle(
        deck,
        ledger,
        sources,
        assets,
        source_manifest_dir=tmp_path,
        asset_manifest_dir=tmp_path,
    )
    assert any(issue.code == "REQUIRED_ASSET_MISSING" for issue in report.issues)


def test_high_risk_qualifier_must_be_present(tmp_path: Path) -> None:
    sources, ledger, assets = _base_models(tmp_path)
    deck = DeckManifest(
        deck=DeckMeta(
            id="demo-private",
            title="Demo",
            audience="Reviewers",
            visibility=Visibility.PRIVATE,
            source_policy=DeckSourcePolicy.PUBLIC_AND_PRIVATE,
        ),
        slides=[
            SlideSpec(
                id="01-status",
                order=1,
                role="test",
                layout=SlideLayout.APPENDIX,
                visibility=Visibility.PRIVATE,
                title="Status",
                message="A private implementation claim.",
                source_refs=[SourceRef(source_id="private-source", section="Status")],
                claim_ids=["private-claim"],
                notes="Qualifier intentionally omitted.",
            )
        ],
    )
    report = validate_bundle(
        deck,
        ledger,
        sources,
        assets,
        source_manifest_dir=tmp_path,
        asset_manifest_dir=tmp_path,
    )
    assert any(issue.code == "CLAIM_QUALIFIER_MISSING" for issue in report.issues)


def test_sparta_example_manifests_validate_with_resolved_sources(tmp_path: Path) -> None:
    profile = ROOT / "examples" / "sparta-explorer"
    sources = load_yaml(profile / "source_manifest.yaml", SourceManifest)
    ledger = load_yaml(profile / "claim_ledger.yaml", ClaimLedger)
    assets = load_yaml(profile / "asset_manifest.yaml", AssetManifest)
    public = load_yaml(profile / "deck.public.yaml", DeckManifest)
    private = load_yaml(profile / "deck.private-appendix.yaml", DeckManifest)
    assert len(public.slides) == 12
    assert len(private.slides) == 7

    resolved_sources = []
    for source in sources.sources:
        path = tmp_path / f"{source.id}.md"
        path.write_text("# Fixture source\n", encoding="utf-8")
        resolved_sources.append(source.model_copy(update={"path": str(path)}))
    sources = sources.model_copy(update={"sources": resolved_sources})

    image_path = tmp_path / "asset.png"
    from PIL import Image

    Image.new("RGB", (1600, 900), "white").save(image_path)
    assets = assets.model_copy(
        update={
            "assets": [
                asset.model_copy(
                    update={"local_path": str(image_path), "status": AssetStatus.PRESENT}
                )
                for asset in assets.assets
            ]
        }
    )

    for deck in (public, private):
        report = validate_bundle(
            deck,
            ledger,
            sources,
            assets,
            source_manifest_dir=tmp_path,
            asset_manifest_dir=tmp_path,
            require_approved_claims=True,
        )
        assert report.errors == 0, [issue.model_dump() for issue in report.issues]


def test_emit_ui_bundle_positive_and_seam_stamp(tmp_path: Path) -> None:
    source_path = FIXTURE / "source_manifest.yaml"
    source = load_yaml(source_path, SourceManifest)
    planned = tmp_path / "planned"
    plan_bundle(source, source_manifest_path=source_path, output_dir=planned, max_slides=10)

    from readme_to_pitchdeck.ui_emitter import emit_ui_bundle

    deck = load_yaml(planned / "deck.public.yaml", DeckManifest)
    ledger = load_yaml(planned / "claim_ledger.yaml", ClaimLedger)
    sources = load_yaml(planned / "source_manifest.resolved.yaml", SourceManifest)
    assets = load_yaml(planned / "asset_manifest.yaml", AssetManifest)
    out = tmp_path / "ui"
    receipt, bundle = emit_ui_bundle(
        deck,
        ledger,
        sources,
        assets,
        source_manifest_dir=planned,
        asset_manifest_dir=planned,
        output_dir=out,
    )
    assert (out / "deck.data.json").exists()
    assert (out / "emit_ui_receipt.json").exists()
    assert bundle.seam_validation.kind == "ui_deck_bundle"
    assert bundle.seam_validation.status == "PASS"
    assert len(bundle.slides) == len(deck.slides)
    assert [slide.order for slide in bundle.slides] == sorted(s.order for s in deck.slides)
    # Every slide claim badge must resolve to a real ledger claim.
    ledger_ids = {claim.id for claim in ledger.claims}
    assert all(badge.id in ledger_ids for slide in bundle.slides for badge in slide.claims)


def test_emit_ui_bundle_fails_closed_on_validation_errors(tmp_path: Path) -> None:
    sources, ledger, assets = _base_models(tmp_path)
    from readme_to_pitchdeck.models import ClaimGuard, DeckManifest as DM, DeckMeta, SlideSpec, VisualSpec
    from readme_to_pitchdeck.ui_emitter import emit_ui_bundle

    deck = DM(
        deck=DeckMeta(
            id="demo-public",
            title="Demo",
            audience="Reviewers",
            visibility=Visibility.PUBLIC,
            source_policy="public_only",
        ),
        slides=[
            SlideSpec(
                id="s1",
                order=1,
                role="cover",
                layout="cover",
                visibility=Visibility.PUBLIC,
                title="Demo",
                message="Leaks a private claim.",
                claim_ids=["private-claim"],
                visual=VisualSpec(),
                claim_guard=ClaimGuard(),
                source_refs=[SourceRef(source_id="public-source")],
            )
        ],
    )
    with pytest.raises(ValueError, match="validation"):
        emit_ui_bundle(
            deck,
            ledger,
            sources,
            assets,
            source_manifest_dir=tmp_path,
            asset_manifest_dir=tmp_path,
            output_dir=tmp_path / "ui",
        )
    assert not (tmp_path / "ui" / "deck.data.json").exists()


def test_apply_slide_edit_positive_and_fail_closed(tmp_path: Path) -> None:
    source_path = FIXTURE / "source_manifest.yaml"
    source = load_yaml(source_path, SourceManifest)
    planned = tmp_path / "planned"
    plan_bundle(source, source_manifest_path=source_path, output_dir=planned, max_slides=10)

    from readme_to_pitchdeck.slide_edit import apply_slide_edit

    deck = load_yaml(planned / "deck.public.yaml", DeckManifest)
    slide = deck.slides[1]
    out = tmp_path / "ui"

    receipt = apply_slide_edit(
        planned, out, slide_id=slide.id, field="title", value="An edited, bounded title"
    )
    assert receipt.seam_validation.kind == "slide_edit_receipt"
    edited = load_yaml(planned / "deck.public.yaml", DeckManifest)
    assert edited.slides[1].title == "An edited, bounded title"
    assert (out / "deck.data.json").exists()

    # Fail closed: an over-length title is rejected and nothing changes.
    before = (planned / "deck.public.yaml").read_text()
    with pytest.raises(Exception):
        apply_slide_edit(planned, out, slide_id=slide.id, field="title", value="x" * 200)
    assert (planned / "deck.public.yaml").read_text() == before

    # Fail closed: forbidden unqualified phrase in visible text is rejected.
    forbidden = source.policy.forbidden_unqualified_claims
    if forbidden:
        with pytest.raises(Exception):
            apply_slide_edit(
                planned, out, slide_id=slide.id, field="message", value=f"We are {forbidden[0]}."
            )
        assert (planned / "deck.public.yaml").read_text().count("An edited, bounded title") == 1

    # Non-editable fields are refused.
    with pytest.raises(ValueError, match="not editable"):
        apply_slide_edit(planned, out, slide_id=slide.id, field="claim_ids", value="x")


def _make_test_mp4(path: Path) -> None:
    import subprocess

    subprocess.run(
        [
            "ffmpeg", "-y", "-f", "lavfi", "-i", "testsrc=duration=1:size=320x180:rate=10",
            "-pix_fmt", "yuv420p", str(path),
        ],
        check=True,
        capture_output=True,
        timeout=60,
    )


def test_video_asset_builds_and_validates(tmp_path: Path) -> None:
    import shutil as _shutil

    if _shutil.which("ffmpeg") is None:
        pytest.skip("ffmpeg unavailable")
    sources, ledger, assets = _base_models(tmp_path)
    video_path = tmp_path / "demo.mp4"
    _make_test_mp4(video_path)

    from readme_to_pitchdeck.models import (
        AssetKind,
        AssetSpec,
        ClaimGuard,
        DeckManifest as DM,
        DeckMeta,
        SlideSpec,
        VisualSpec,
        VisualType,
    )

    assets = AssetManifest(
        project_name="Demo",
        assets=[
            AssetSpec(
                id="demo-video",
                kind=AssetKind.VIDEO,
                visibility=Visibility.PUBLIC,
                local_path=str(video_path),
                alt_text="Product demo clip",
                required=True,
                status=AssetStatus.PRESENT,
            )
        ],
    )
    deck = DM(
        deck=DeckMeta(
            id="demo-public", title="Demo", audience="Reviewers",
            visibility=Visibility.PUBLIC, source_policy="public_only",
        ),
        slides=[
            SlideSpec(
                id="s1", order=1, role="demo", layout="screenshot",
                visibility=Visibility.PUBLIC, title="Demo video",
                message="A one-second generated clip.",
                claim_ids=["public-claim"],
                visual=VisualSpec(type=VisualType.SCREENSHOT, asset_id="demo-video"),
                claim_guard=ClaimGuard(),
                source_refs=[SourceRef(source_id="public-source")],
            )
        ],
    )
    report = validate_bundle(
        deck, ledger, sources, assets,
        source_manifest_dir=tmp_path, asset_manifest_dir=tmp_path,
    )
    assert report.errors == 0, [i.model_dump() for i in report.issues]

    output = tmp_path / "video.pptx"
    _, validation = build_pptx(
        deck, ledger, sources, assets,
        source_manifest_dir=tmp_path, asset_manifest_dir=tmp_path,
        output_path=output, require_approved_claims=False,
    )
    assert validation.errors == 0
    reopened = Presentation(output)
    movie_shapes = [
        shape for slide in reopened.slides for shape in slide.shapes
        if shape.shape_type == 16  # MSO_SHAPE_TYPE.MEDIA
    ]
    assert movie_shapes, "no embedded media shape found in built PPTX"

    # Wrong container for a video asset fails validation.
    png = tmp_path / "not-video.png"
    from PIL import Image

    Image.new("RGB", (100, 100), "white").save(png)
    bad_assets = AssetManifest(
        project_name="Demo",
        assets=[assets.assets[0].model_copy(update={"local_path": str(png)})],
    )
    bad_report = validate_bundle(
        deck, ledger, sources, bad_assets,
        source_manifest_dir=tmp_path, asset_manifest_dir=tmp_path,
    )
    assert any(i.code == "ASSET_UNSUPPORTED_FORMAT" and i.severity == "error" for i in bad_report.issues)


def test_emit_markdown_one_way_export(tmp_path: Path) -> None:
    source_path = FIXTURE / "source_manifest.yaml"
    source = load_yaml(source_path, SourceManifest)
    planned = tmp_path / "planned"
    plan_bundle(source, source_manifest_path=source_path, output_dir=planned, max_slides=10)

    from readme_to_pitchdeck.md_emitter import emit_markdown

    deck = load_yaml(planned / "deck.public.yaml", DeckManifest)
    receipt = emit_markdown(
        deck,
        load_yaml(planned / "claim_ledger.yaml", ClaimLedger),
        load_yaml(planned / "source_manifest.resolved.yaml", SourceManifest),
        load_yaml(planned / "asset_manifest.yaml", AssetManifest),
        source_manifest_dir=planned,
        asset_manifest_dir=planned,
        output_dir=tmp_path / "md",
    )
    text = (tmp_path / "md" / "deck.md").read_text()
    assert text.startswith("---\nmarp: true")
    assert text.count("\n---\n") >= len(deck.slides)
    assert receipt.seam_validation.kind == "emit_md_receipt"
    for slide in deck.slides:
        assert f"## {slide.title}" in text


def test_deck_ops_add_move_delete(tmp_path: Path) -> None:
    source_path = FIXTURE / "source_manifest.yaml"
    source = load_yaml(source_path, SourceManifest)
    planned = tmp_path / "planned"
    plan_bundle(source, source_manifest_path=source_path, output_dir=planned, max_slides=10)

    from readme_to_pitchdeck.slide_edit import apply_deck_op

    deck = load_yaml(planned / "deck.public.yaml", DeckManifest)
    first = sorted(deck.slides, key=lambda s: s.order)[0]
    count = len(deck.slides)
    out = tmp_path / "ui"

    apply_deck_op(planned, out, op="add_after", slide_id=first.id)
    deck2 = load_yaml(planned / "deck.public.yaml", DeckManifest)
    assert len(deck2.slides) == count + 1
    ordered = sorted(deck2.slides, key=lambda s: s.order)
    assert ordered[1].id == "new-slide"
    assert [s.order for s in ordered] == list(range(1, count + 2))

    apply_deck_op(planned, out, op="move_right", slide_id="new-slide")
    deck3 = load_yaml(planned / "deck.public.yaml", DeckManifest)
    assert sorted(deck3.slides, key=lambda s: s.order)[2].id == "new-slide"

    apply_deck_op(planned, out, op="delete", slide_id="new-slide")
    deck4 = load_yaml(planned / "deck.public.yaml", DeckManifest)
    assert len(deck4.slides) == count
    assert [s.order for s in sorted(deck4.slides, key=lambda s: s.order)] == list(range(1, count + 1))

    with pytest.raises(ValueError, match="already first"):
        apply_deck_op(planned, out, op="move_left", slide_id=first.id)


def test_transitions_reveals_and_asset_ops(tmp_path: Path) -> None:
    source_path = FIXTURE / "source_manifest.yaml"
    source = load_yaml(source_path, SourceManifest)
    planned = tmp_path / "planned"
    plan_bundle(source, source_manifest_path=source_path, output_dir=planned, max_slides=10)
    out = tmp_path / "ui"

    from readme_to_pitchdeck.asset_ops import add_asset_to_slide, clear_slide_visual
    from readme_to_pitchdeck.slide_edit import apply_slide_edit

    deck = load_yaml(planned / "deck.public.yaml", DeckManifest)
    slide = sorted(deck.slides, key=lambda s: s.order)[1]

    # Transition + reveal flow through the whole pipeline as typed data.
    apply_slide_edit(planned, out, slide_id=slide.id, field="transition", value="zoom")
    apply_slide_edit(planned, out, slide_id=slide.id, field="reveal", value="stagger_fade")
    import json as _json

    bundle = _json.loads((out / "deck.data.json").read_text())
    ui_slide = [s for s in bundle["slides"] if s["id"] == slide.id][0]
    assert ui_slide["transition"] == "zoom"
    assert ui_slide["reveal"] == "stagger_fade"
    with pytest.raises(Exception):
        apply_slide_edit(planned, out, slide_id=slide.id, field="transition", value="wobble")

    # Asset drop: image binds, appears in manifest + emitted bundle.
    from PIL import Image

    image_path = tmp_path / "shot.png"
    Image.new("RGB", (1280, 720), "navy").save(image_path)
    receipt = add_asset_to_slide(planned, out, slide_id=slide.id, file_path=image_path, alt_text="A navy screenshot")
    assert receipt.operation == "asset-add"
    assets = load_yaml(planned / "asset_manifest.yaml", AssetManifest)
    added = [a for a in assets.assets if a.id.startswith(slide.id)]
    assert added and added[0].alt_text == "A navy screenshot"
    bundle = _json.loads((out / "deck.data.json").read_text())
    ui_slide = [s for s in bundle["slides"] if s["id"] == slide.id][0]
    assert ui_slide["visual"]["asset"]["file"].startswith("assets/")

    # Position variant round-trips.
    apply_slide_edit(planned, out, slide_id=slide.id, field="visual:position", value="left")
    deck2 = load_yaml(planned / "deck.public.yaml", DeckManifest)
    assert [s for s in deck2.slides if s.id == slide.id][0].visual.position.value == "left"

    # Unsupported format fails closed with no orphan file.
    bad = tmp_path / "notes.txt"
    bad.write_text("nope")
    with pytest.raises(ValueError, match="unsupported asset format"):
        add_asset_to_slide(planned, out, slide_id=slide.id, file_path=bad, alt_text="x")

    # Clear visual detaches.
    clear_slide_visual(planned, out, slide_id=slide.id)
    deck3 = load_yaml(planned / "deck.public.yaml", DeckManifest)
    assert [s for s in deck3.slides if s.id == slide.id][0].visual.asset_id is None


def test_freeform_layout_round_trip(tmp_path: Path) -> None:
    source_path = FIXTURE / "source_manifest.yaml"
    source = load_yaml(source_path, SourceManifest)
    planned = tmp_path / "planned"
    plan_bundle(source, source_manifest_path=source_path, output_dir=planned, max_slides=10)
    out = tmp_path / "ui"

    from readme_to_pitchdeck.slide_edit import apply_slide_edit

    deck = load_yaml(planned / "deck.public.yaml", DeckManifest)
    slide = sorted(deck.slides, key=lambda s: s.order)[1]

    # Switching to freeform synthesizes elements from existing content.
    apply_slide_edit(planned, out, slide_id=slide.id, field="layout", value="freeform")
    deck2 = load_yaml(planned / "deck.public.yaml", DeckManifest)
    ff = [s for s in deck2.slides if s.id == slide.id][0]
    assert ff.layout.value == "freeform"
    assert any(e.id == "title" for e in ff.elements)

    # Frame move persists exact fractions; PPTX places the shape at them.
    apply_slide_edit(planned, out, slide_id=slide.id, field="element:title:frame", value="0.25,0.1,0.5,0.12")
    ledger = load_yaml(planned / "claim_ledger.yaml", ClaimLedger)
    sources = load_yaml(planned / "source_manifest.resolved.yaml", SourceManifest)
    assets = load_yaml(planned / "asset_manifest.yaml", AssetManifest)
    deck3 = load_yaml(planned / "deck.public.yaml", DeckManifest)
    output = planned / "freeform.pptx"
    build_pptx(
        deck3, ledger, sources, assets,
        source_manifest_dir=planned, asset_manifest_dir=planned,
        output_path=output, require_approved_claims=False,
    )
    from pptx.util import Emu

    prs = Presentation(output)
    target = list(prs.slides)[1]
    boxes = [
        (Emu(sh.left).inches / 13.333, Emu(sh.top).inches / 7.5)
        for sh in target.shapes
        if sh.has_text_frame and ff.title[:20] in sh.text_frame.text
    ]
    assert boxes and abs(boxes[0][0] - 0.25) < 0.001 and abs(boxes[0][1] - 0.1) < 0.001

    # Out-of-bounds frames fail closed.
    with pytest.raises(Exception):
        apply_slide_edit(planned, out, slide_id=slide.id, field="element:title:frame", value="0.8,0.8,0.5,0.5")

    # Element text joins the forbidden-phrase scan.
    forbidden = sources.policy.forbidden_unqualified_claims
    if forbidden:
        with pytest.raises(Exception):
            apply_slide_edit(
                planned, out, slide_id=slide.id, field="element:title:text", value=f"We are {forbidden[0]}!"
            )
