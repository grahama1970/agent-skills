"""#1380: matched adversarial style negatives for the house gate.

Each mutant preserves the nuisance variables the gate measures (words, colors,
object counts, rough areas) while breaking exactly one style dimension a real
looks-like-Graham classifier must catch. A mutant that PASSES the gate is a
documented false-pass — the hole, reproduced.

Inputs: a passing deck.pptx. Outputs: mutant .pptx files + cases.json.
Failure modes: an unreadable deck raises; every mutation is seeded (7) so
mutants are deterministic.
"""

from __future__ import annotations

import json
import random
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

SEED = 7


def _content_shapes(slide):
    """Authored, movable content: skip inherited chrome, band boxes, footers."""
    out = []
    for shape in slide.shapes:
        name = shape.name or ""
        if name.startswith("chrome:"):
            continue
        try:
            if shape.top is None or shape.left is None:
                continue
            # keep the band title in place: it lives in the top strip
            if shape.top < 0.14 * 5143500:
                continue
        except (TypeError, AttributeError):
            continue
        out.append(shape)
    return out


def bbox_shuffle(deck_path: Path, out_path: Path) -> None:
    """Permute content positions. Pixel marginals barely move; layout dies."""
    rng = random.Random(SEED)
    pres = Presentation(str(deck_path))
    for slide in pres.slides:
        shapes = _content_shapes(slide)
        if len(shapes) < 2:
            continue
        spots = [(s.left, s.top) for s in shapes]
        rng.shuffle(spots)
        for shape, (left, top) in zip(shapes, spots):
            shape.left, shape.top = left, top
    pres.save(str(out_path))


def typography_swap(deck_path: Path, out_path: Path) -> None:
    """Same words and boxes; ransom-note type: mixed fonts/sizes, no hierarchy."""
    rng = random.Random(SEED)
    fonts = ["Comic Sans MS", "Courier New", "Impact", "Times New Roman"]
    pres = Presentation(str(deck_path))
    for slide in pres.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
                continue
            if (shape.name or "").startswith("chrome:"):
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = rng.choice(fonts)
                    if run.font.size is not None:
                        run.font.size = Pt(max(8, int(run.font.size.pt * rng.choice([0.55, 0.8, 1.5, 2.1]))))
                    run.font.underline = rng.random() < 0.3
                    run.font.bold = rng.random() < 0.5
    pres.save(str(out_path))


def two_tiny_visuals(deck_path: Path, out_path: Path) -> None:
    """Strip every real visual, satisfy density with two meaningless 20px groups."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Emu

    pres = Presentation(str(deck_path))
    for slide in pres.slides:
        doomed = [s for s in slide.shapes
                  if s.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.GROUP)
                  and not (s.name or "").startswith("chrome:")]
        for shape in doomed:
            shape._element.getparent().remove(shape._element)
        for i in (0, 1):
            group = slide.shapes.add_group_shape()
            group.left, group.top = Emu(457200 + i * 300000), Emu(4000000)
            group.width, group.height = Emu(190500), Emu(190500)  # ~20px each
    pres.save(str(out_path))


def layout_mirror(deck_path: Path, out_path: Path) -> None:
    """Mirror every content shape horizontally. Ink, palette, words, and object
    counts are EXACTLY preserved (no overlap is introduced), but every
    composition rule — chevrons left, art placement, mark positions, reading
    order — is broken. A spatially blind gate cannot see this."""
    pres = Presentation(str(deck_path))
    width = pres.slide_width
    for slide in pres.slides:
        for shape in _content_shapes(slide):
            shape.left = width - shape.left - shape.width
    pres.save(str(out_path))


def arc_shuffle(deck_path: Path, out_path: Path) -> None:
    """Individually valid pages in an implausible order: close first, cover
    mid-deck, dividers orphaned. Per-slide channels see nothing."""
    import copy

    rng = random.Random(SEED)
    pres = Presentation(str(deck_path))
    xml_slides = pres.slides._sldIdLst
    slides = list(xml_slides)
    rng.shuffle(slides)
    for sld in slides:
        xml_slides.remove(sld)
    for sld in slides:
        xml_slides.append(sld)
    pres.save(str(out_path))


def art_register_swap(deck_path: Path, out_path: Path, replacement: Path | None = None) -> None:
    """Same text, layout, and boxes — the drawn-scene art replaced by generic
    glossy 3D stock art (the register Graham never uses)."""
    replacement = replacement or Path(__file__).parent.parent / "fixtures" / "house-gate" / "generic-3d-robot.png"
    payload = replacement.read_bytes()
    pres = Presentation(str(deck_path))
    for slide in pres.slides:
        for shape in slide.shapes:
            if (shape.name or "") in {"el:scene-art", "el:proof-scene"}:
                ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
                ns_r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
                blip = shape._element.find(f".//{{{ns_a}}}blip")
                if blip is None:
                    continue
                image_part = shape.part.related_part(blip.get(f"{{{ns_r}}}embed"))
                image_part._blob = payload
    pres.save(str(out_path))


def text_dump_divider(deck_path: Path, out_path: Path) -> None:
    """Archetype mismatch: section dividers carry a dense left-aligned prose
    dump instead of one huge centered heading."""
    from pptx.enum.text import PP_ALIGN

    pres = Presentation(str(deck_path))
    for slide in pres.slides:
        for shape in slide.shapes:
            if (shape.name or "") == "el:divider-heading":
                tf = shape.text_frame
                text = tf.text
                tf.clear()
                run = tf.paragraphs[0].add_run()
                run.text = (text + ". ") * 12
                run.font.size = Pt(12)
                tf.paragraphs[0].alignment = PP_ALIGN.LEFT
                shape.left = int(0.05 * 9144000)
    pres.save(str(out_path))


def card_grid(deck_path: Path, out_path: Path) -> None:
    """The review's most-likely false PASS: a generic startup card-grid deck
    skinned with the house chrome, words, and palette."""
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
    from pptx.util import Emu

    pres = Presentation(str(deck_path))
    slide_w, slide_h = pres.slide_width, pres.slide_height
    for slide in pres.slides:
        texts = []
        doomed = []
        for shape in slide.shapes:
            name = shape.name or ""
            if name.startswith("chrome:") or not name.startswith("el:"):
                continue
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame and shape.text_frame.text.strip():
                if shape.top is not None and shape.top > 0.14 * slide_h:
                    texts.append(shape.text_frame.text.strip())
            doomed.append(shape)
        for shape in doomed:
            shape._element.getparent().remove(shape._element)
        # 2x2 mechanically symmetric cards, house palette, same words
        chunks = texts + [""] * 4
        for i in range(4):
            gx, gy = i % 2, i // 2
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Emu(int((0.07 + gx * 0.45) * slide_w)), Emu(int((0.2 + gy * 0.36) * slide_h)),
                Emu(int(0.41 * slide_w)), Emu(int(0.3 * slide_h)))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(0x07, 0x68, 0x89) if i % 2 == 0 else RGBColor(0xD6, 0xA3, 0x00)
            card.line.color.rgb = RGBColor(0x06, 0x5E, 0x7C)
            tf = card.text_frame
            tf.word_wrap = True
            run = tf.paragraphs[0].add_run()
            run.text = (chunks[i] or "SPARTA Explorer")[:180]
            run.font.size = Pt(13)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    pres.save(str(out_path))


MUTANTS = {"bbox-shuffle": bbox_shuffle, "typography-swap": typography_swap,
           "two-tiny-visuals": two_tiny_visuals, "layout-mirror": layout_mirror,
           "arc-shuffle": arc_shuffle, "art-register-swap": art_register_swap,
           "text-dump-divider": text_dump_divider, "card-grid": card_grid}


def main() -> None:
    source = Path(sys.argv[sys.argv.index("--source-pptx") + 1])
    out_dir = Path(sys.argv[sys.argv.index("--output-dir") + 1])
    out_dir.mkdir(parents=True, exist_ok=True)
    cases = []
    for name, fn in MUTANTS.items():
        target = out_dir / f"{name}.pptx"
        fn(source, target)
        cases.append({"mutant": name, "pptx": str(target),
                      "breaks": {"bbox-shuffle": "layout/composition",
                                 "typography-swap": "typography/hierarchy",
                                 "two-tiny-visuals": "visual substance",
                                 "layout-mirror": "composition/reading order",
                                 "arc-shuffle": "narrative arc",
                                 "art-register-swap": "illustration register",
                                 "text-dump-divider": "archetype anatomy",
                                 "card-grid": "composition (generic grid, house skin)"}[name],
                      "expected_from_real_classifier": "FAIL"})
    (out_dir / "cases.json").write_text(json.dumps({"source": str(source), "seed": SEED,
                                                     "cases": cases}, indent=1))
    print(json.dumps({"built": [c["mutant"] for c in cases], "out_dir": str(out_dir)}))


if __name__ == "__main__":
    main()
