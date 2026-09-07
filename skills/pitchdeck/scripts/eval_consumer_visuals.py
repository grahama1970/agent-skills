#!/usr/bin/env python3
"""Live OAI consumer regression: positional header titles and themed picture frames.

Requires the emitted OAI deck and its dedicated Vite instance. Uses real Surf
and real PPTX output; never edits the source deck or recolors image pixels.
"""
from __future__ import annotations

import argparse
import hashlib
import json
import os
from pathlib import Path
import re
import subprocess
import time
from zipfile import ZipFile

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[3]
SKILL = ROOT / 'skills/pitchdeck'
SURF = ROOT / 'skills/surf/run.sh'
OAI = Path(os.environ.get('OAI_TRIAL_ROOT', str(ROOT.parent / 'oai-trial')))


def run(*args):
    print('STEP', str(args[0]), str(args[1]) if len(args) > 1 else '', flush=True)
    p = subprocess.run(list(map(str, args)), cwd=ROOT, capture_output=True,
                       text=True, timeout=120)
    if p.returncode:
        raise RuntimeError(f'{args[0]} failed ({p.returncode}): {p.stdout}\n{p.stderr}')
    return p.stdout


def media(path):
    with ZipFile(path) as z:
        return {n: hashlib.sha256(z.read(n)).hexdigest()
                for n in z.namelist() if n.startswith('ppt/media/')}


def geometry(path):
    p = Presentation(path)
    return [[(s.name, s.left, s.top, s.width, s.height, s.text if s.has_text_frame else '')
             for s in slide.shapes] for slide in p.slides]


def frame_check(out):
    source = OAI / 'docs/pitch/oai-trial/reorganized'
    baseline = Path('/mnt/storage12tb/oai-trial/visual-fixes/before/oai-trial-current.pptx')
    target = out.with_suffix('.pptx')
    run(SKILL / 'run.sh', 'build', '--deck', source / 'deck.public.yaml',
        '--claim-ledger', source / 'claim_ledger.yaml', '--source-manifest',
        source / 'source_manifest.yaml', '--asset-manifest', source / 'asset_manifest.yaml',
        '--output', target, '--draft-watermark', '--allow-candidate-claims', '--json')
    ns = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    counts = []
    for path in [baseline, target]:
        deck = Presentation(path)
        colors = [n.get('val', '').upper() for s in deck.slides
                  for n in s._element.findall(f'.//{ns}srgbClr')]
        counts.append(sum(c in {'071019', '2A455A'} for c in colors))
    assert counts[0] > 0, 'baseline must demonstrate the old blue frame colors'
    assert counts[1] == 0, f'legacy blue frame paint remains: {counts[1]}'
    assert media(baseline) == media(target), 'raster/SVG materialized media changed'
    assert geometry(baseline) == geometry(target), 'fixed geometry or text changed'
    # Theme selection stays opt-in: no-canvas means no presentation mutation.
    from pitchdeck.models import ThemeTokens
    from pitchdeck.theme_style import apply_presentation_theme
    deck = Presentation(baseline)
    before = [s._element.xml for s in deck.slides]
    apply_presentation_theme(deck, ThemeTokens())
    assert before == [s._element.xml for s in deck.slides], 'default theme changed'
    return {'before_blue_nodes': counts[0], 'after_blue_nodes': counts[1],
            'media_unchanged': True, 'geometry_unchanged': True,
            'default_theme_unchanged': True, 'pptx': str(target)}


def header_check(out, url):
    run(SURF, 'tab.list', '--json')
    created = run(SURF, 'window.new', url, '--unfocused', '--width', '1400', '--height', '1100')
    tab = re.search(r'\(tab (\d+)\)', created).group(1)

    def surf(*args):
        flags = [] if args[0] == 'snap' else ['--no-screenshot']
        return run(SURF, *args, *flags, '--tab-id', tab, '--no-activate')

    def js(code):
        raw = json.loads(surf('js', 'return (async () => {' + code + '})()', '--no-screenshot'))
        return json.loads(raw) if isinstance(raw, str) else raw

    result = {'tab_id': tab, 'window_creation': created, 'checks': []}
    out.with_name(out.stem + '-ownership.json').write_text(json.dumps(result, indent=2) + '\n')
    try:
        deadline = time.monotonic() + 30
        while not js('return !!document.querySelector(".slide-viewport");'):
            if time.monotonic() > deadline:
                raise RuntimeError('actual OAI viewer did not mount')
            time.sleep(.2)
        if 'Zoom: 100%' not in surf('zoom'):
            surf('emulate.device', 'reset')
            surf('zoom', '1')
        for width in [1920, 960]:
            surf('emulate.viewport', '--width', str(width), '--height', '1080')
            for slide_id, title, expected in [
                ('r01-toc', 'Table of Contents', 1),
                ('r30-thank-you', 'Thank you', 2),
            ]:
                js('location.hash=' + json.dumps('/slide/' + slide_id) + ';return true')
                deadline = time.monotonic() + 30
                while True:
                    settled = js('const p=document.querySelector(".slide-viewport");return {id:p?.dataset.slideId,layout:p?.dataset.layout,width:innerWidth}')
                    if settled == {'id': slide_id, 'layout': 'responsive' if width < 1100 else 'canvas', 'width': width}:
                        break
                    if time.monotonic() > deadline:
                        raise RuntimeError(f'viewport did not settle: {settled}')
                    time.sleep(.1)
                code = '''
const id=__ID__, title=__TITLE__, width=__WIDTH__;
await document.fonts.ready;
const p=document.querySelector('.slide-viewport');
if(p?.dataset.slideId!==id) throw Error('wrong slide');
if(innerWidth!==width) throw Error('viewport not applied');
const norm=s=>s.replace(/\\s+/g,' ').trim();
const matches=[...p.querySelectorAll('.freeform-band span,.freeform-element p')]
 .filter(e=>norm(e.innerText)===title && getComputedStyle(e).display!=='none');
const targets=[...p.querySelectorAll('[data-animation-target="title"]')];
const band=p.querySelector('.freeform-band');
return {band_position:getComputedStyle(band).position, band_height:band.getBoundingClientRect().height,
 texture_height:parseFloat(getComputedStyle(band,'::before').height),slide_id:id, viewport:[innerWidth,innerHeight], pane_width:p.getBoundingClientRect().width,
 scroll_width:document.documentElement.scrollWidth, hidden:document.hidden, layout:p.dataset.layout,
 title_count:matches.length, title_targets:targets.length,
 target_in_band:!!targets[0]?.closest('.freeform-band'),
 source_elements:targets.map(e=>e.dataset.elementId)};
'''.replace('__ID__', json.dumps(slide_id)).replace('__TITLE__', json.dumps(title)).replace('__WIDTH__', str(width))
                observed = js(code)
                if width < 1100:
                    assert observed['band_position'] == 'relative', observed
                    assert observed['texture_height'] <= observed['band_height'] * 1.2 + 1, observed
                assert observed['title_count'] == expected, observed
                assert observed['title_targets'] == 1, 'animation identity lost or duplicated'
                if slide_id == 'r01-toc':
                    assert observed['target_in_band'] == (width < 1100), observed
                else:
                    assert not observed['target_in_band'], 'centered hero was suppressed'
                shot = out.with_name(f'{out.stem}-{width}-{slide_id}.png')
                surf('snap', '--output', str(shot))
                observed['screenshot'] = str(shot)
                result['checks'].append(observed)
        result['passed'] = True
        return result
    except Exception:
        result['failure_state'] = js('return {url:location.href,ready:document.readyState,text:document.body.innerText.slice(0,2000)};')
        out.with_name(out.stem + '-diagnostic.json').write_text(json.dumps(result, indent=2) + '\n')
        surf('snap', '--output', str(out.with_name(out.stem + '-failure.png')))
        raise
    finally:
        # The window/tab was created by this test, never the human's project tab.
        if result.get('passed'):
            run(SURF, 'tab.close', tab, '--no-screenshot')


def main():
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument('--case', choices=['header', 'frame'], required=True)
    ap.add_argument('--out', type=Path, required=True)
    ap.add_argument('--url', default='http://127.0.0.1:3016/?deck=./oai-trial-current/deck.data.json&rehearse=1#/slide/r01-toc')
    args = ap.parse_args()
    args.out.parent.mkdir(parents=True, exist_ok=True)
    receipt = {'schema': 'pitchdeck.consumer_visuals.v1', 'live': True,
               'mocked': False, 'case': args.case,
               'does_not_prove': 'Human approval, universal export parity, or complete animation coverage.'}
    try:
        receipt['evidence'] = header_check(args.out, args.url) if args.case == 'header' else frame_check(args.out)
        receipt['status'] = 'PASS'
    except Exception as exc:
        receipt.update(status='FAIL', error=str(exc))
    args.out.write_text(json.dumps(receipt, indent=2) + '\n')
    print(json.dumps({'status': receipt['status'], 'receipt': str(args.out), 'error': receipt.get('error')}))
    return 0 if receipt['status'] == 'PASS' else 1


if __name__ == '__main__':
    raise SystemExit(main())
