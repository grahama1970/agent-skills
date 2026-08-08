"""Voice-profile and transform-policy compilers (#1311, voice layer slice 1).

Compiles the author's REAL headline corpus (best-practices-slide-design
references/exemplars.yaml) into a content-addressed
``pitchdeck.voice_profile.v1`` artifact, and provides
``pitchdeck.transform_policy.v1``: the guard taxonomy plus a registry of
DIRECTIONAL, versioned aggregation/generalization mappings. The policy is the
LAW half of the voice layer — the proposal stage (slice 2, #1312) may only
emit candidates this policy can prove or explicitly refuse.

Inputs: exemplars.yaml (headlines quoted inside ``why`` fields — a documented
stable grammar), transform mapping instances under design/transform_mappings/.
Outputs: VoiceProfile (deterministic; sha256-stable across runs),
PolicyCheck violations with typed codes. Failure modes: exemplars without an
extractable headline are recorded as coverage gaps, never silently dropped;
unapproved mappings never license a rewrite; every check is fail-closed —
pass, refuse with a code, never warn-and-continue.
"""

from __future__ import annotations

import hashlib
import json
import re
from pathlib import Path
from typing import Literal

import yaml
from pydantic import Field

from .models import StrictModel

# Grammar note (correctness-regex-only-known-grammar): exemplars.yaml `why`
# fields quote headlines in straight or curly double quotes. This is a
# committed file with fixture coverage, not free text.
_QUOTED = re.compile(r"[\"“]([^\"”]{4,80})[\"”]")

# Guard taxonomy (transform_policy_v1). Modality guards are DIRECTIONAL:
# strengthening `can` into `does/will/always` asserts more than the ledger.
POLARITY_GUARDS = ("not", "no", "without", "never", "unless", "except")
MODALITY_GUARDS = ("can", "may", "must", "could", "should")
SCOPE_GUARDS = ("only", "some", "every", "all", "each")
TIME_GUARDS = ("today", "now", "current", "planned", "before", "after", "remain", "remains")
STRENGTHENERS = ("does", "will", "always", "guarantees", "proves", "ensures")


class VoiceExemplar(StrictModel):
    id: str = Field(min_length=1)
    role: str = Field(min_length=1, description="Slide role from the exemplar's cited rule.")
    headline: str = Field(min_length=1)
    word_count: int = Field(ge=1)
    syntax: Literal["copular_assertion", "sv_assertion", "question", "fragment"]
    devices: list[str] = Field(default_factory=list)
    source_image: str = Field(min_length=1)


class VoiceStats(StrictModel):
    """Quantified house voice, measured over the author's real deck titles."""

    title_count: int = Field(ge=1)
    median_words: int = Field(ge=1)
    p10_words: int = Field(ge=1)
    p90_words: int = Field(ge=1)
    title_case_ratio: float = Field(ge=0.0, le=1.0)
    question_ratio: float = Field(ge=0.0, le=1.0)
    parenthetical_ratio: float = Field(ge=0.0, le=1.0)
    copular_ratio: float = Field(ge=0.0, le=1.0)


class VoiceProfile(StrictModel):
    schema_: Literal["pitchdeck.voice_profile.v1"] = Field(
        default="pitchdeck.voice_profile.v1", alias="schema"
    )
    corpus_path: str
    exemplars: list[VoiceExemplar]
    coverage_gaps: list[str] = Field(
        default_factory=list, description="Exemplar ids with no extractable headline — visible, not dropped."
    )
    word_count_range: tuple[int, int]
    stats: VoiceStats | None = Field(
        default=None, description="Measured over the real deck corpus when decks_dir is supplied."
    )

    def content_sha256(self) -> str:
        payload = self.model_dump(by_alias=True, mode="json")
        return hashlib.sha256(json.dumps(payload, sort_keys=True).encode()).hexdigest()


def _classify_syntax(headline: str) -> str:
    lower = f" {headline.lower()} "
    if headline.rstrip().endswith("?"):
        return "question"
    if any(f" {v} " in lower for v in ("is", "are", "was", "were")):
        return "copular_assertion"
    words = headline.split()
    if len(words) >= 2:
        return "sv_assertion"
    return "fragment"


def _detect_devices(headline: str) -> list[str]:
    devices: list[str] = []
    if "(" in headline:
        devices.append("parenthetical_aside")
    if any(w in headline.lower() for w in ("horrible", "expensive", "saves", "hate", "love")):
        devices.append("dry_evaluation")
    if "—" in headline or " vs " in headline.lower():
        devices.append("contrast")
    return devices


def harvest_titles(decks_dir: Path) -> list[str]:
    """Every slide title from the author's real decks (the primary voice signal).

    python-pptx title placeholder first, first non-empty text frame as the
    documented fallback; 2-14 words keeps titles and drops body paragraphs."""
    from pptx import Presentation

    titles: list[str] = []
    for deck in sorted(decks_dir.glob("*.pptx")):
        for slide in Presentation(str(deck)).slides:
            text = ""
            try:
                if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
                    text = slide.shapes.title.text_frame.text
            except (AttributeError, KeyError):
                text = ""
            if not text.strip():
                frames = [s.text_frame.text for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]
                text = frames[0] if frames else ""
            text = " ".join(text.split())
            if 2 <= len(text.split()) <= 14:
                titles.append(text)
    return titles


def measure_voice(titles: list[str]) -> VoiceStats:
    """Quantify the house voice from real titles — the numbers the proposer targets."""
    import statistics

    if not titles:
        raise ValueError("cannot measure voice from zero titles")
    counts = sorted(len(t.split()) for t in titles)
    n = len(counts)

    def _title_case(text: str) -> bool:
        words = [w for w in text.split() if w]
        capped = sum(1 for w in words if w[:1].isupper())
        return bool(words) and capped >= max(2, len(words) * 0.6)

    return VoiceStats(
        title_count=n,
        median_words=int(statistics.median(counts)),
        p10_words=counts[n // 10],
        p90_words=counts[(9 * n) // 10],
        title_case_ratio=round(sum(_title_case(t) for t in titles) / n, 3),
        question_ratio=round(sum(t.rstrip().endswith("?") for t in titles) / n, 3),
        parenthetical_ratio=round(sum("(" in t for t in titles) / n, 3),
        copular_ratio=round(sum(bool(re.search(r"\b(is|are)\b", t, re.I)) for t in titles) / n, 3),
    )


def compile_voice_profile(corpus_dir: Path, decks_dir: Path | None = None) -> VoiceProfile:
    """Compile exemplars.yaml (+ real deck titles when available) into a
    deterministic, content-addressed profile."""
    source = corpus_dir / "references" / "exemplars.yaml"
    data = yaml.safe_load(source.read_text(encoding="utf-8"))
    exemplars: list[VoiceExemplar] = []
    gaps: list[str] = []
    for entry in data.get("exemplars", []):
        quoted = _QUOTED.findall(entry.get("why", ""))
        headline = next((q for q in quoted if len(q.split()) >= 2), None)
        if headline is None:
            gaps.append(entry["id"])
            continue
        exemplars.append(
            VoiceExemplar(
                id=entry["id"],
                role=entry["rule"],
                headline=headline,
                word_count=len(headline.split()),
                syntax=_classify_syntax(headline),
                devices=_detect_devices(headline),
                source_image=entry["image"],
            )
        )
    if not exemplars:
        raise ValueError("voice profile compiled zero exemplars — corpus unusable")
    counts = [e.word_count for e in exemplars]
    stats = measure_voice(harvest_titles(decks_dir)) if decks_dir and decks_dir.is_dir() else None
    return VoiceProfile(
        corpus_path=str(source),
        exemplars=sorted(exemplars, key=lambda e: e.id),
        coverage_gaps=sorted(gaps),
        word_count_range=(min(counts), max(counts)),
        stats=stats,
    )


class TransformMapping(StrictModel):
    """A DIRECTIONAL registered rewrite. Unapproved mappings license nothing."""

    schema_: Literal["pitchdeck.transform_mapping.v1"] = Field(
        default="pitchdeck.transform_mapping.v1", alias="schema"
    )
    id: str = Field(min_length=1)
    kind: Literal["aggregation", "generalization"]
    source_terms: list[str] = Field(min_length=1, description="Exact claim-side terms (all must appear).")
    target: str = Field(min_length=1)
    direction_note: str = Field(min_length=1)
    counterexamples: list[str] = Field(default_factory=list)
    status: Literal["candidate", "approved"] = "candidate"
    approved_by: str | None = None


class PolicyViolation(StrictModel):
    code: Literal[
        "GUARD_DROPPED",
        "GUARD_ADDED",
        "MODALITY_STRENGTHENED",
        "COORDINATED_SPAN_BROKEN",
        "UNREGISTERED_REWRITE_TERM",
    ]
    detail: str


def _coordinated_spans(claim_text: str) -> list[list[str]]:
    """Extract guarded coordinated lists: 'a, b, or c' / 'a, b, and c'."""
    spans: list[list[str]] = []
    for match in re.finditer(r"((?:[\w'-]+, )+(?:or|and) [\w'-]+)", claim_text.lower()):
        members = [m.strip() for m in re.split(r", | or | and ", match.group(1)) if m.strip() not in {"or", "and"}]
        if len(members) >= 3:
            spans.append(members)
    return spans


def check_rendering_policy(
    claim_text: str,
    rendering_text: str,
    mappings: list[TransformMapping] | None = None,
) -> list[PolicyViolation]:
    """Fail-closed policy check for a proposed rendering against its claim.

    This composes with (does not replace) planning.verify_rendering: that
    proves the transform class; this proves guard survival, modality
    direction, coordinated-span integrity, and rewrite registration."""
    violations: list[PolicyViolation] = []
    lc, lr = f" {claim_text.lower()} ", f" {rendering_text.lower()} "
    approved = {m.target.lower(): m for m in (mappings or []) if m.status == "approved" and m.approved_by}

    for guard in (*POLARITY_GUARDS, *SCOPE_GUARDS, *TIME_GUARDS):
        if f" {guard} " in lc and f" {guard} " not in lr:
            violations.append(PolicyViolation(code="GUARD_DROPPED", detail=f"claim guard '{guard}' missing from rendering"))
        if f" {guard} " in lr and f" {guard} " not in lc:
            violations.append(PolicyViolation(code="GUARD_ADDED", detail=f"rendering adds guard '{guard}' absent from claim"))
    for weak in MODALITY_GUARDS:
        if f" {weak} " in lc:
            for strong in STRENGTHENERS:
                if f" {strong} " in lr and f" {strong} " not in lc:
                    violations.append(
                        PolicyViolation(
                            code="MODALITY_STRENGTHENED",
                            detail=f"claim modality '{weak}' strengthened to '{strong}'",
                        )
                    )
    for span in _coordinated_spans(claim_text):
        present = [m for m in span if m in lr]
        if present and len(present) < len(span):
            missing = sorted(set(span) - set(present))
            violations.append(
                PolicyViolation(
                    code="COORDINATED_SPAN_BROKEN",
                    detail=f"coordinated span {span} partially rendered; missing {missing}",
                )
            )
    # Any rendering word absent from the claim must be licensed by an
    # APPROVED mapping whose source terms all appear in the claim.
    claim_words = set(re.findall(r"[\w'-]+", lc))
    for word in set(re.findall(r"[\w'-]+", lr)) - claim_words:
        mapping = approved.get(word)
        if mapping and all(t.lower() in lc for t in mapping.source_terms):
            continue
        if word in {"a", "an", "the", "be", "is", "are", "what's", "it", "its"}:
            continue  # closed-class glue introduced by legal inflection
        violations.append(
            PolicyViolation(
                code="UNREGISTERED_REWRITE_TERM",
                detail=f"rendering term '{word}' is neither claim text nor an approved mapping target",
            )
        )
    return violations
