"""Template-as-base emission: inherit the house deck instead of redrawing it.

Every AI-deck tool that gets brand fidelity right does the same thing — open the
real .pptx as the base presentation, delete its slides (which keeps the theme,
slide master, and layouts), and add slides on the author's own layouts. The
band, texture, logo mark, footer line, page number, fonts, and colors then come
from the template by construction rather than from measurement.

That matters here because measurement was demonstrably lossy: the band fill was
reproduced as #065E7C when the template's is #076889, the photographic band
texture was invented as diagonal lines before being found in the layout, and the
bottom-left mark blind judges named in every round is simply the template's own
logo. None of those are approximation problems worth solving — they are
inheritance problems.

Inputs: a house .pptx. Outputs: a stripped Presentation plus a LayoutProfile
describing each layout's placeholders (slides inherit from layouts strictly by
placeholder idx, so the idx map IS the contract). Failure modes: a template with
no usable layout raises rather than silently falling back to python-pptx's
default white theme, which would look nothing like the house.
"""

from __future__ import annotations

from pathlib import Path
from typing import Literal

from pydantic import Field

from .models import StrictModel

EMU_PER_INCH = 914400.0
_RELS_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"


class PlaceholderSlot(StrictModel):
    """One fillable slot. `idx` is the inheritance key, not a position."""

    idx: int
    name: str
    kind: str
    x: float = 0.0
    y: float = 0.0
    w: float = 0.0
    h: float = 0.0


class LayoutSlot(StrictModel):
    index: int
    name: str
    placeholders: list[PlaceholderSlot] = Field(default_factory=list)
    role: Literal["title-only", "title-content", "blank", "other"] = "other"

    @property
    def title_idx(self) -> int | None:
        for slot in self.placeholders:
            if "title" in slot.kind.lower() or "titre" in slot.name.lower() or slot.idx == 0:
                return slot.idx
        return None


class LayoutProfile(StrictModel):
    schema_: Literal["pitchdeck.layout_profile.v1"] = Field(
        default="pitchdeck.layout_profile.v1", alias="schema"
    )
    template: str
    layouts: list[LayoutSlot]

    def by_role(self, role: str) -> LayoutSlot | None:
        return next((layout for layout in self.layouts if layout.role == role), None)

    def named(self, fragment: str) -> LayoutSlot | None:
        lowered = fragment.lower()
        return next((layout for layout in self.layouts if lowered in layout.name.lower()), None)


def profile_template(template: Path) -> LayoutProfile:
    """Enumerate every layout and its placeholder idx map."""
    from pptx import Presentation

    presentation = Presentation(str(template))
    width_in = presentation.slide_width / EMU_PER_INCH
    height_in = presentation.slide_height / EMU_PER_INCH
    layouts: list[LayoutSlot] = []
    for index, layout in enumerate(presentation.slide_layouts):
        slots: list[PlaceholderSlot] = []
        for placeholder in layout.placeholders:
            try:
                x = (placeholder.left or 0) / EMU_PER_INCH / width_in
                y = (placeholder.top or 0) / EMU_PER_INCH / height_in
                w = (placeholder.width or 0) / EMU_PER_INCH / width_in
                h = (placeholder.height or 0) / EMU_PER_INCH / height_in
            except (TypeError, ZeroDivisionError):
                x = y = w = h = 0.0
            slots.append(PlaceholderSlot(
                idx=placeholder.placeholder_format.idx,
                name=placeholder.name,
                kind=str(placeholder.placeholder_format.type),
                x=round(x, 4), y=round(y, 4), w=round(w, 4), h=round(h, 4),
            ))
        titles = [s for s in slots if "TITLE" in s.kind.upper() or s.idx == 0]
        if not slots:
            role = "blank"
        elif titles and len(slots) == 1:
            role = "title-only"
        elif titles:
            role = "title-content"
        else:
            role = "other"
        layouts.append(LayoutSlot(index=index, name=layout.name, placeholders=slots, role=role))  # type: ignore[arg-type]
    if not layouts:
        raise ValueError(f"template '{template}' exposes no slide layouts")
    return LayoutProfile(template=str(template), layouts=layouts)


def open_stripped_template(template: Path):
    """The house deck with every slide removed — theme, master, and layouts kept.

    This is the canonical python-pptx idiom (issue #310): drop the relationship
    and remove the sldId. Orphaned media stays in the package, which costs file
    size but nothing else."""
    from pptx import Presentation

    presentation = Presentation(str(template))
    slide_ids = presentation.slides._sldIdLst
    for slide_id in list(slide_ids):
        relationship_id = slide_id.get(_RELS_NS)
        if relationship_id:
            presentation.part.drop_rel(relationship_id)
        slide_ids.remove(slide_id)
    return presentation


def pick_layout(profile: LayoutProfile, *, prefer_named: str | None, needs_body: bool) -> LayoutSlot:
    """Choose the author's layout that fits what this slide has to carry.

    Preference order: an explicitly named house layout, then a layout whose role
    matches, then any layout with a title. A template with none of those is a
    template we cannot honestly use."""
    if prefer_named:
        named = profile.named(prefer_named)
        if named is not None:
            return named
    wanted = "title-content" if needs_body else "title-only"
    exact = profile.by_role(wanted)
    if exact is not None:
        return exact
    fallback = next((layout for layout in profile.layouts if layout.title_idx is not None), None)
    if fallback is None:
        raise ValueError(f"template '{profile.template}' has no layout with a title placeholder")
    return fallback


def fill_title(slide, layout: LayoutSlot, text: str) -> bool:
    """Write the assertion into the layout's own title placeholder.

    Returns False when the layout has no title slot, so the caller can fall back
    to an explicit textbox rather than silently dropping the assertion."""
    title_idx = layout.title_idx
    if title_idx is None:
        return False
    for placeholder in slide.placeholders:
        if placeholder.placeholder_format.idx == title_idx:
            placeholder.text_frame.text = text
            return True
    return False


def drop_unused_placeholders(slide, keep_idxs: set[int]) -> list[int]:
    """Remove empty inherited placeholders so they do not render as prompt text.

    An unfilled placeholder shows 'Click to add text' in some viewers; leaving
    them is worse than never adding them."""
    removed: list[int] = []
    for placeholder in list(slide.placeholders):
        idx = placeholder.placeholder_format.idx
        if idx in keep_idxs:
            continue
        if placeholder.has_text_frame and placeholder.text_frame.text.strip():
            continue
        placeholder._element.getparent().remove(placeholder._element)
        removed.append(idx)
    return removed


def _rewrite_preserving_format(shape, text: str) -> None:
    """Replace a shape's text while KEEPING its run formatting.

    Assigning to ``text_frame.text`` discards every run and its properties, so
    the template's small red disclaimer type came back as default body text.
    Writing into the first existing run keeps size, colour, weight, and font;
    later runs are removed so no stale fragment survives."""
    frame = shape.text_frame
    paragraphs = frame.paragraphs
    if not paragraphs or not paragraphs[0].runs:
        frame.text = text
        return
    first = paragraphs[0].runs[0]
    first.text = text
    for run in paragraphs[0].runs[1:]:
        run._r.getparent().remove(run._r)
    for paragraph in paragraphs[1:]:
        paragraph._p.getparent().remove(paragraph._p)


def retarget_disclaimer(presentation, disclaimer: str, *, stale_markers: tuple[str, ...]) -> dict:
    """Replace the template's inherited disclaimer with this deck's own.

    A house template LOCKS its disclaimer deliberately — that is the point of
    brand governance — which means inheriting one onto a different owner's deck
    silently asserts the wrong ownership. The text lives on the slide master and
    on several layouts, so every copy must be retargeted; a surviving stale
    marker is a false legal claim, so callers should fail closed on one.

    Returns a receipt naming every shape rewritten and any marker still found."""
    rewritten: list[str] = []
    for master_index, master in enumerate(presentation.slide_masters):
        for shape in master.shapes:
            if shape.has_text_frame and any(m in shape.text_frame.text for m in stale_markers):
                _rewrite_preserving_format(shape, disclaimer)
                rewritten.append(f"master[{master_index}]:{shape.name}")
    for layout_index, layout in enumerate(presentation.slide_layouts):
        for shape in layout.shapes:
            if shape.has_text_frame and any(m in shape.text_frame.text for m in stale_markers):
                _rewrite_preserving_format(shape, disclaimer)
                rewritten.append(f"layout[{layout_index}]:{shape.name}")
    residual = scan_stale_markers(presentation, stale_markers)
    return {"rewritten": rewritten, "residual_markers": residual}


def scan_stale_markers(presentation, stale_markers: tuple[str, ...]) -> list[str]:
    """Every place a stale owner marker still appears (masters, layouts, slides)."""
    found: list[str] = []
    pools = [(f"master[{i}]", m.shapes) for i, m in enumerate(presentation.slide_masters)]
    pools += [(f"layout[{i}]", l.shapes) for i, l in enumerate(presentation.slide_layouts)]
    pools += [(f"slide[{i}]", s.shapes) for i, s in enumerate(presentation.slides)]
    for where, shapes in pools:
        for shape in shapes:
            if shape.has_text_frame and any(m in shape.text_frame.text for m in stale_markers):
                found.append(f"{where}:{shape.name}")
    return found
