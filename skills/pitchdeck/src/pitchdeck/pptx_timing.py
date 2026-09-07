"""Native PowerPoint p:timing emission shared by the canonical and legacy exporters.

Grounded in the supplied style-corpus decks (outputs/progressive-reveal/research/
provided-preset-*.xml and ReqML_GE_Presentation-slide*-timing.xml retained by
research/brief.md) and the ECMA-376 p:timing structure: a tmRoot par holds one
mainSeq p:seq; each on-click row opens a new click group (cTn with
delay="indefinite"); with-previous rows share the previous row's start;
after-previous rows start at the previous row's computed end plus their own
delay — the same arithmetic as timeline() in ui/src/animations.ts.

Playback mechanics are carried by the behaviors (p:set visibility,
p:animEffect filters, p:anim ppt_x/y/w/h, p:animScale, p:animRot, p:animClr,
p:animMotion), exactly the elements observed in the supplied decks. presetID
labels the effect in PowerPoint's animation pane; per-effect ID confidence is
documented in docs/ANIMATIONS.md. This exporter does not yet resolve nested
targets (diagram nodes/edges, grouped children); these are recorded as skipped
with a reason, never silently retargeted.
"""
from __future__ import annotations

from typing import Callable

from lxml import etree
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

# Effect -> PowerPoint preset ID. Confidence per effect lives in docs/ANIMATIONS.md:
# grounded IDs come from the supplied decks' XML; other pane labels remain
# unverified. Emitted behaviors are not proof of application playback.
_PRESET = {
    "appear": 1, "fade": 53, "fly": 2, "wipe": 22, "zoom": 23, "peek": 12,
    "split": 16, "expand": 50, "stretch": 17, "rise": 34, "grow-turn": 31,
    "blinds": 3, "box": 4, "bars": 14, "checker": 5, "strips": 18,
    "spin": 8, "grow-shrink": 6, "transparency": 9, "dim": 9, "pulse": 35,
    "font-color": 3, "fill-color": 1, "line-color": 7, "motion-line": 42,
}
_DIR_SUB = {"left": "8", "right": "2", "up": "1", "down": "4", "in": "16", "out": "32",
            "horizontal": "10", "vertical": "5"}


def _subtype(a) -> str:
    """Directional bitfield only where the pane uses one (grounded: fly 8=left)."""
    if a.phase in ("emphasis", "motion") or a.effect in ("appear", "fade", "expand", "rise", "grow-turn"):
        return "0"
    if a.effect in ("zoom", "box"):
        return "32" if a.direction == "out" else "16"
    return _DIR_SUB.get(a.direction, "0")
_NODE_TYPE = {"on-click": "clickEffect", "with-previous": "withEffect", "after-previous": "afterEffect"}


def _e(parent, tag: str, **attrs):
    el = etree.SubElement(parent, qn(tag))
    for key, value in attrs.items():
        el.set(key.rstrip("_"), str(value))
    return el


def _tgt(cbhvr, spid: int) -> None:
    _e(_e(cbhvr, "p:tgtEl"), "p:spTgt", spid=spid)


def _set_attr(ch, spid: int, attr: str, value: str, ids, *, delay: int = 0) -> None:
    node = _e(ch, "p:set")
    cb = _e(node, "p:cBhvr")
    ctn = _e(cb, "p:cTn", id=ids(), dur="1", fill="hold")
    _e(_e(ctn, "p:stCondLst"), "p:cond", delay=delay)
    _tgt(cb, spid)
    _e(_e(cb, "p:attrNameLst"), "p:attrName").text = attr
    _e(_e(node, "p:to"), "p:strVal", val=value)


def _fx(ch, spid: int, filt: str, dur: int, ids, *, out: bool = False) -> None:
    node = _e(ch, "p:animEffect", transition="out" if out else "in", filter=filt)
    cb = _e(node, "p:cBhvr")
    _e(cb, "p:cTn", id=ids(), dur=dur)
    _tgt(cb, spid)


def _anim(ch, spid: int, attr: str, tavs, dur: int, ids) -> None:
    node = _e(ch, "p:anim", calcmode="lin", valueType="num")
    cb = _e(node, "p:cBhvr", additive="base")
    _e(cb, "p:cTn", id=ids(), dur=dur, fill="hold")
    _tgt(cb, spid)
    _e(_e(cb, "p:attrNameLst"), "p:attrName").text = attr
    tav_list = _e(node, "p:tavLst")
    for tm, val in tavs:
        _e(_e(_e(tav_list, "p:tav", tm=tm), "p:val"), "p:strVal" if isinstance(val, str) else "p:fltVal", val=val)


def _anim_scale(ch, spid: int, x: int, y: int, dur: int, ids, *, auto_rev: bool = False) -> None:
    node = _e(ch, "p:animScale")
    cb = _e(node, "p:cBhvr")
    ctn = _e(cb, "p:cTn", id=ids(), dur=dur, fill="hold")
    if auto_rev:
        ctn.set("autoRev", "1")
    _tgt(cb, spid)
    _e(node, "p:to", x=x, y=y)


def _anim_rot(ch, spid: int, by_60k: int, dur: int, ids) -> None:
    node = _e(ch, "p:animRot", by=by_60k)
    cb = _e(node, "p:cBhvr")
    _e(cb, "p:cTn", id=ids(), dur=dur, fill="hold")
    _tgt(cb, spid)
    _e(_e(cb, "p:attrNameLst"), "p:attrName").text = "r"


def _anim_clr(ch, spid: int, attr: str, hex_color: str, dur: int, ids) -> None:
    node = _e(ch, "p:animClr", clrSpc="rgb", dir="cw")
    cb = _e(node, "p:cBhvr")
    _e(cb, "p:cTn", id=ids(), dur=dur, fill="hold")
    _tgt(cb, spid)
    _e(_e(cb, "p:attrNameLst"), "p:attrName").text = attr
    _e(_e(node, "p:to"), "a:srgbClr", val=hex_color.lstrip("#").upper())


def _fly_tavs(direction: str, out: bool):
    home_x, home_y = "#ppt_x", "#ppt_y"
    off = {"left": ("0-#ppt_w/2", home_y), "right": ("1+#ppt_w/2", home_y),
           "up": (home_x, "0-#ppt_h/2"), "down": (home_x, "1+#ppt_h/2")}
    off_x, off_y = off.get(direction, off["left"])
    if out:
        return [("0", home_x), ("100000", off_x)], [("0", home_y), ("100000", off_y)]
    return [("0", off_x), ("100000", home_x)], [("0", off_y), ("100000", home_y)]


def _behaviors(ch, spid: int, a, dur: int, ids) -> None:
    """Append the behaviors for one AnimationEffect row onto one shape."""
    d, out = a.direction, a.phase == "exit"
    entrance = a.phase == "entrance"
    if entrance:
        _set_attr(ch, spid, "style.visibility", "visible", ids)
    axis = "horizontal" if d in ("left", "right", "horizontal") else "vertical"
    if a.effect == "appear":
        if out:
            _set_attr(ch, spid, "style.visibility", "hidden", ids)
        return
    if a.effect == "fade":
        _fx(ch, spid, "fade", dur, ids, out=out)
    elif a.effect == "fly":
        x_tavs, y_tavs = _fly_tavs(d, out)
        _anim(ch, spid, "ppt_x", x_tavs, dur, ids)
        _anim(ch, spid, "ppt_y", y_tavs, dur, ids)
    elif a.effect == "wipe":
        _fx(ch, spid, f"wipe({d if d in ('left', 'right', 'up', 'down') else 'left'})", dur, ids, out=out)
    elif a.effect == "peek":
        edge = d if d in ("left", "right", "up", "down") else "left"
        _fx(ch, spid, f"wipe({edge})", dur, ids, out=out)
        near = {"left": "#ppt_x-#ppt_w/4", "right": "#ppt_x+#ppt_w/4"}.get(edge)
        if near:
            tavs = [("0", "#ppt_x"), ("100000", near)] if out else [("0", near), ("100000", "#ppt_x")]
            _anim(ch, spid, "ppt_x", tavs, dur, ids)
        else:
            near_y = "#ppt_y-#ppt_h/4" if edge == "up" else "#ppt_y+#ppt_h/4"
            tavs = [("0", "#ppt_y"), ("100000", near_y)] if out else [("0", near_y), ("100000", "#ppt_y")]
            _anim(ch, spid, "ppt_y", tavs, dur, ids)
    elif a.effect == "split":
        _fx(ch, spid, f"split({'out' if out else 'in'}{'Horizontal' if axis == 'horizontal' else 'Vertical'})",
            dur, ids, out=out)
    elif a.effect == "zoom":
        big = d == "out"
        start, end = ("0", "#ppt_w") if not big else ("#ppt_w*2", "#ppt_w")
        start_h, end_h = ("0", "#ppt_h") if not big else ("#ppt_h*2", "#ppt_h")
        if out:
            start, end, start_h, end_h = end, start, end_h, start_h
        _anim(ch, spid, "ppt_w", [("0", start), ("100000", end)], dur, ids)
        _anim(ch, spid, "ppt_h", [("0", start_h), ("100000", end_h)], dur, ids)
    elif a.effect == "expand":
        _fx(ch, spid, "fade", dur, ids, out=out)
        tavs = [("0", "#ppt_w"), ("100000", "0")] if out else [("0", "0"), ("100000", "#ppt_w")]
        _anim(ch, spid, "ppt_w", tavs, dur, ids)
    elif a.effect == "stretch":
        attr = "ppt_w" if axis == "horizontal" else "ppt_h"
        full = f"#{attr}"
        tavs = [("0", full), ("100000", "0")] if out else [("0", "0"), ("100000", full)]
        _anim(ch, spid, attr, tavs, dur, ids)
    elif a.effect == "rise":
        _fx(ch, spid, "fade", dur, ids, out=out)
        tavs = ([("0", "#ppt_y"), ("100000", "#ppt_y+#ppt_h/2")] if out
                else [("0", "#ppt_y+#ppt_h/2"), ("100000", "#ppt_y")])
        _anim(ch, spid, "ppt_y", tavs, dur, ids)
    elif a.effect == "grow-turn":
        # Native approximation: scale+fade without the turn (initial-rotation
        # state is not portable through p:animRot). docs/ANIMATIONS.md names it.
        _fx(ch, spid, "fade", dur, ids, out=out)
        _anim(ch, spid, "ppt_w", ([("0", "#ppt_w"), ("100000", "0")] if out
                                  else [("0", "0"), ("100000", "#ppt_w")]), dur, ids)
        _anim(ch, spid, "ppt_h", ([("0", "#ppt_h"), ("100000", "0")] if out
                                  else [("0", "0"), ("100000", "#ppt_h")]), dur, ids)
    elif a.effect == "blinds":
        _fx(ch, spid, f"blinds({axis})", dur, ids, out=out)
    elif a.effect == "box":
        _fx(ch, spid, f"box({'out' if d == 'out' else 'in'})", dur, ids, out=out)
    elif a.effect == "bars":
        _fx(ch, spid, f"randombar({axis})", dur, ids, out=out)
    elif a.effect == "checker":
        _fx(ch, spid, f"checkerboard({'across' if axis == 'horizontal' else 'down'})", dur, ids, out=out)
    elif a.effect == "strips":
        strip = {"left": "downLeft", "right": "downRight", "up": "upLeft", "down": "downRight"}.get(d, "downLeft")
        _fx(ch, spid, f"strips({strip})", dur, ids, out=out)
    elif a.effect == "spin":
        _anim_rot(ch, spid, int(21600000 * a.amount), dur, ids)
    elif a.effect == "grow-shrink":
        scale = int(100000 * a.amount)
        _anim_scale(ch, spid, scale, scale, dur, ids)
    elif a.effect in ("transparency", "dim"):
        _anim(ch, spid, "style.opacity", [("0", 1.0), ("100000", round(max(0.0, 1 - a.amount), 3))], dur, ids)
    elif a.effect == "pulse":
        scale = int(100000 * (1 + a.amount))
        _anim_scale(ch, spid, scale, scale, max(dur // 2, 1), ids, auto_rev=True)
    elif a.effect == "font-color":
        _anim_clr(ch, spid, "style.color", a.color, dur, ids)
    elif a.effect == "fill-color":
        _set_attr(ch, spid, "fill.type", "solid", ids)
        _set_attr(ch, spid, "fill.on", "true", ids)
        _anim_clr(ch, spid, "fillcolor", a.color, dur, ids)
    elif a.effect == "line-color":
        _set_attr(ch, spid, "stroke.on", "true", ids)
        _anim_clr(ch, spid, "stroke.color", a.color, dur, ids)
    elif a.effect == "motion-line":
        node = _e(ch, "p:animMotion", origin="layout", path=f"M 0 0 L {a.dx:.5f} {a.dy:.5f}",
                  pathEditMode="relative", rAng="0", ptsTypes="AA")
        cb = _e(node, "p:cBhvr")
        _e(cb, "p:cTn", id=ids(), dur=dur, fill="hold")
        _tgt(cb, spid)
        names = _e(cb, "p:attrNameLst")
        _e(names, "p:attrName").text = "ppt_x"
        _e(names, "p:attrName").text = "ppt_y"
    if out:
        _set_attr(ch, spid, "style.visibility", "hidden", ids, delay=max(dur - 1, 0))


def _click_groups(animations):
    """Mirror of the browser timeline(): (row, start_ms) grouped per click."""
    groups: list[list] = []
    first_auto = False
    start = end = 0
    for a in animations:
        if not groups and a.start != "on-click":
            first_auto = True
        if a.start == "on-click" or not groups:
            groups.append([])
        if a.start == "on-click":
            start = a.delay_ms
        elif a.start == "after-previous":
            start = end + a.delay_ms
        else:
            start = start + a.delay_ms
        end = start + a.duration_ms
        groups[-1].append((a, start))
    return groups, first_auto


def apply_slide_timing(slide, animations, resolve: Callable[[str], tuple[list[int], str | None]]) -> dict:
    """Append a p:timing tree for the slide's authored animation rows.

    resolve(target) returns (spids, skip_reason). A row with zero resolvable
    targets is skipped with its reasons in the receipt; nothing is silently
    retargeted. Returns a receipt dict; emits no p:timing when nothing applies.
    """
    receipt: dict = {"clicks": 0, "effects": 0, "skipped": []}
    counter = [0]

    def ids() -> str:
        counter[0] += 1
        return str(counter[0])

    rows: list[tuple] = []  # (effect_row, start_ms, [spids])
    for a in animations:
        spids: list[int] = []
        for target in a.targets:
            found, reason = resolve(target)
            if found:
                spids.extend(found)
            else:
                receipt["skipped"].append({"row": a.id, "target": target,
                                           "reason": reason or "no shape emitted for target"})
        rows.append((a, spids))
    if not any(spids for _, spids in rows):
        return receipt

    for stale in slide._element.findall(qn("p:timing")):
        slide._element.remove(stale)
    timing = etree.SubElement(slide._element, qn("p:timing"))
    tn_list = _e(timing, "p:tnLst")
    root_ctn = _e(_e(tn_list, "p:par"), "p:cTn", id=ids(), dur="indefinite", restart="never", nodeType="tmRoot")
    seq = _e(_e(root_ctn, "p:childTnLst"), "p:seq", concurrent="1", nextAc="seek")
    main_ctn = _e(seq, "p:cTn", id=ids(), dur="indefinite", nodeType="mainSeq")
    main_children = _e(main_ctn, "p:childTnLst")

    kept = {id(a): spids for a, spids in rows if spids}
    groups, first_auto = _click_groups([a for a, spids in rows if spids])
    grp_ids: dict[int, int] = {}
    bld_entries: list[tuple[int, int]] = []
    for index, group in enumerate(groups):
        click_ctn = _e(_e(main_children, "p:par"), "p:cTn", id=ids(), fill="hold")
        click_cond = _e(click_ctn, "p:stCondLst")
        _e(click_cond, "p:cond", delay="indefinite")
        if index == 0 and first_auto:
            # Supplied-deck pattern: a leading with/after row runs on slide begin.
            _e(_e(click_cond, "p:cond", evt="onBegin", delay="0"), "p:tn", val="2")
        inner_ctn = _e(_e(_e(click_ctn, "p:childTnLst"), "p:par"), "p:cTn", id=ids(), fill="hold")
        _e(_e(inner_ctn, "p:stCondLst"), "p:cond", delay="0")
        inner_children = _e(inner_ctn, "p:childTnLst")
        for a, start_ms in group:
            dur = max(a.duration_ms, 1)
            # One effect node per shape (PowerPoint's own serialization);
            # additional shapes of a multi-target row ride the same start as
            # withEffect nodes so one authored row stays one concept.
            for position, spid in enumerate(kept[id(a)]):
                grp = grp_ids.setdefault(spid, 0)
                grp_ids[spid] = grp + 1
                bld_entries.append((spid, grp))
                effect_ctn = _e(_e(inner_children, "p:par"), "p:cTn", id=ids(),
                                presetID=_PRESET[a.effect],
                                presetClass={"entrance": "entr", "exit": "exit", "emphasis": "emph",
                                             "motion": "path"}[a.phase],
                                presetSubtype=_subtype(a), fill="hold", grpId=grp,
                                nodeType=_NODE_TYPE[a.start] if position == 0 else "withEffect")
                _e(_e(effect_ctn, "p:stCondLst"), "p:cond", delay=start_ms)
                _behaviors(_e(effect_ctn, "p:childTnLst"), spid, a, dur, ids)
            receipt["effects"] += 1
    receipt["clicks"] = sum(1 for g in groups if any(a.start == "on-click" for a, _ in g))
    for holder, evt in (("p:prevCondLst", "onPrev"), ("p:nextCondLst", "onNext")):
        cond = _e(_e(seq, holder), "p:cond", evt=evt, delay="0")
        _e(_e(cond, "p:tgtEl"), "p:sldTgt")
    bld_list = _e(timing, "p:bldLst")
    for spid, grp in dict.fromkeys(bld_entries):
        _e(bld_list, "p:bldP", spid=spid, grpId=grp)
    return receipt


def spid_resolver(slide, *, prefix: str = "el:") -> Callable[[str], tuple[list[int], str | None]]:
    """Map browser animation targets to top-level shape IDs by emitted name.

    Canonical names: el:ID (elements/groups), el:D:node:N / el:D:edge:E
    (diagram parts, nested inside the diagram group). Legacy semantic names:
    sem:body:N / sem:visual:N. Nested shapes are reported, not retargeted:
    This exporter resolves only top-level shapes and whole groups.
    """
    top: dict[str, list[int]] = {}
    nested: set[str] = set()

    def walk(shapes, in_group: bool) -> None:
        for shape in shapes:
            name = shape.name or ""
            if name:
                if in_group:
                    nested.add(name)
                else:
                    top.setdefault(name, []).append(shape.shape_id)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                walk(shape.shapes, True)

    walk(slide.shapes, False)

    def resolve(target: str) -> tuple[list[int], str | None]:
        name = prefix + target.replace("/node/", ":node:").replace("/edge/", ":edge:")
        spids = list(top.get(name, []))
        for candidate, found in top.items():
            if candidate.startswith(name + ":"):
                spids.extend(found)
        if spids:
            return spids, None
        if name in nested or any(n.startswith(name + ":") for n in nested):
            return [], "nested-in-group: this exporter does not yet resolve individual grouped targets"
        return [], None

    return resolve
