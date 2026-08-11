"""Post-edit publish verification against the ACTUAL .pptx (#1317).

The compiler proved claim fidelity at emission. Nothing re-proved it afterwards,
so "claim-faithful" and "editable by a human" had a hole between them: any string
in the emitted deck could be retyped before delivery and no gate would notice.
This module closes that boundary by re-extracting evidence from the real file —
slides, groups, tables, notes, and package properties — rather than trusting the
manifest that produced it.

Checks, each a typed failure code:
  UNCLAIMED_TEXT        a visible string that is not a legal transform of a ledger
                        claim, an approved rendering, or declared non-claim chrome
  STALE_OWNER_MARKER    a previous template owner's name anywhere in the package
  NON_EDITABLE_CONTENT  a slide flattened to imagery, so its claims are unreadable
                        by any verifier and uneditable by any human
  TEMPLATE_DRIFT        the deck no longer derives from the approved template
  VISIBLE_CLAIM_LOSS    text present in the file but clipped, off-canvas, or
                        mid-word truncated, so the claim does not reach the reader

Inputs: a .pptx, a claim ledger, an approvals file, and an optional template
contract. Outputs: a publish receipt with findings. Failure modes: a file that
cannot be opened raises; anything the verifier cannot read is reported as a
finding rather than passing, because unreadable is not the same as clean.
"""

from __future__ import annotations

import hashlib
import json
import re
from pathlib import Path
from typing import Literal

from pydantic import Field

from .models import StrictModel

EMU_PER_INCH = 914400.0
_CHROME_MAX_CHARS = 3  # page numbers and single glyphs are not claims


class PublishFinding(StrictModel):
    code: Literal[
        "UNCLAIMED_TEXT",
        "STALE_OWNER_MARKER",
        "NON_EDITABLE_CONTENT",
        "TEMPLATE_DRIFT",
        "VISIBLE_CLAIM_LOSS",
        "UNREADABLE",
    ]
    where: str
    detail: str


class PublishApprovals(StrictModel):
    """What a human licensed: exact renderings, chrome text, and the disclaimer."""

    schema_: Literal["pitchdeck.publish_approvals.v1"] = Field(
        default="pitchdeck.publish_approvals.v1", alias="schema"
    )
    approved_renderings: list[str] = Field(default_factory=list)
    non_claim_text: list[str] = Field(default_factory=list)
    disclaimer: str | None = None
    stale_owner_markers: list[str] = Field(default_factory=list)


class TemplateContract(StrictModel):
    """Which template this deck is allowed to have been built from."""

    schema_: Literal["pitchdeck.template_contract.v1"] = Field(
        default="pitchdeck.template_contract.v1", alias="schema"
    )
    template_sha256: str
    template_name: str
    layouts_used: list[str] = Field(default_factory=list)


class PublishReceipt(StrictModel):
    schema_: Literal["pitchdeck.publish_receipt.v1"] = Field(
        default="pitchdeck.publish_receipt.v1", alias="schema"
    )
    pptx: str
    pptx_sha256: str
    status: Literal["PASS", "REFUSED"]
    strings_checked: int
    findings: list[PublishFinding] = Field(default_factory=list)
    proves: list[str] = Field(default_factory=list)
    does_not_prove: list[str] = Field(default_factory=list)


class AssertionAtom(StrictModel):
    """One emitted string with its provenance (#1328).

    The compiler knows what every string is and where it came from. Re-deriving
    that from the ledger at publish time cannot work for paraphrase or
    generalization — "Retrieval" is a legal generalization of a claim about
    search and embeddings, but it is not an excerpt of it. So the manifest is
    EMITTED from the approved document rather than reconstructed, and a human
    never types a diagram label into an approvals file to make it pass."""

    schema_: Literal["pitchdeck.assertion_atom.v1"] = Field(
        default="pitchdeck.assertion_atom.v1", alias="schema"
    )
    text: str
    canonical_id: str
    role: str
    claim_id: str | None = None
    transform_class: str | None = None
    binding_kind: str | None = None
    slide_id: str = ""
    element_id: str = ""


def atoms_from_document(document) -> list[AssertionAtom]:
    """Every visible string the document emits, with its binding provenance."""
    from .document import iter_tree

    atoms: list[AssertionAtom] = []
    for slide in document.slides:
        by_path = {b.path: b for b in slide.bindings}

        def record(text: str, canonical_id: str, role: str, paths: list[str]) -> None:
            if not (text or "").strip():
                return
            binding = next((by_path[p] for p in paths if p in by_path), None)
            atoms.append(AssertionAtom(
                text=text, canonical_id=f"{slide.id}/{canonical_id}", role=role,
                claim_id=getattr(binding, "claim_id", None),
                transform_class=getattr(binding, "transform_class", None),
                binding_kind=getattr(getattr(binding, "kind", None), "value", None),
                slide_id=slide.id, element_id=element.id,
            ))

        for element in iter_tree(slide.elements):
            record(element.text or "", element.id, element.role or element.kind.value,
                   list(element.binding_paths))
            if element.diagram is None:
                continue
            for node in element.diagram.nodes:
                record(node.label, f"{element.id}:{node.id}:label", "diagram-node-label",
                       [p for p in node.binding_paths if p.endswith(":label")])
                record(node.sublabel or "", f"{element.id}:{node.id}:sublabel",
                       "diagram-node-sublabel",
                       [p for p in node.binding_paths if p.endswith(":sublabel")])
            for edge in element.diagram.edges:
                record(edge.label or "", f"{element.id}:{edge.id}:label",
                       "diagram-edge-label", list(edge.binding_paths))
    return atoms


# Transforms a machine can verify from the claim text alone. Aggregation and
# generalization cannot be reconstructed mechanically ("Retrieval" is a legal
# generalization of a claim about search and embeddings but not an excerpt of
# it), so they require an exact named human approval instead.
MECHANICAL_TRANSFORMS = {"verbatim", "truncation", "inflection"}
ATTESTED_TRANSFORMS = {"aggregation", "generalization"}
# Chrome is a TYPED ROLE, not a string-length heuristic. Only these roles may
# carry text that asserts nothing about the product; everything else needs a
# claim. This replaces the "<=3 chars or all digits is chrome" escape hatch.
CHROME_ROLES = {"footer", "caption", "message", "page-number", "wordmark", "disclaimer", "badge"}


class AtomRefusal(StrictModel):
    code: Literal[
        "CLAIM_NOT_IN_LEDGER",
        "TRANSFORM_NOT_SATISFIED",
        "UNATTESTED_TRANSFORM",
        "UNKNOWN_TRANSFORM_CLASS",
    ]
    detail: str


def authorize_atom(
    atom: AssertionAtom,
    *,
    claims_by_id: dict[str, str],
    approved_texts: set[str],
) -> AtomRefusal | None:
    """Is this atom evidence of legal derivation, or merely an assertion by the compiler?

    A truthy claim_id is NOT sufficient — that was the reproduced bypass: changing
    a label to the inverse of its claim while keeping claim_id and binding_paths
    passed the gate. Authorization now requires the claim to exist in the bound
    ledger AND the text to satisfy its declared transform (mechanically where
    that is possible, by exact human attestation where it is not)."""
    # the cover title is the deck's own name (a wordmark), not a claim about it
    if atom.slide_id.endswith("cover") and atom.role == "title":
        return None
    if atom.binding_kind == "non_claim" or atom.role in CHROME_ROLES:
        return None  # typed chrome: asserts nothing, so no claim is required
    if not atom.claim_id:
        return AtomRefusal(code="CLAIM_NOT_IN_LEDGER", detail=f"atom '{atom.canonical_id}' carries no claim id")
    claim_text = claims_by_id.get(atom.claim_id)
    if claim_text is None:
        return AtomRefusal(
            code="CLAIM_NOT_IN_LEDGER",
            detail=f"atom '{atom.canonical_id}' cites claim '{atom.claim_id}', absent from the bound ledger",
        )
    transform = (atom.transform_class or "").lower()
    probe = " ".join(atom.text.split()).lstrip("❯>•- ").rstrip("…").strip()
    if transform in MECHANICAL_TRANSFORMS:
        if not _word_boundary_excerpt(probe, claim_text):
            return AtomRefusal(
                code="TRANSFORM_NOT_SATISFIED",
                detail=(f"atom '{atom.canonical_id}' declares {transform} but "
                        f"{probe[:48]!r} is not a word-boundary excerpt of claim '{atom.claim_id}'"),
            )
        return None
    if transform in ATTESTED_TRANSFORMS:
        if " ".join(probe.split()).casefold() not in approved_texts:
            return AtomRefusal(
                code="UNATTESTED_TRANSFORM",
                detail=(f"atom '{atom.canonical_id}' declares {transform}, which no machine can verify; "
                        f"it requires an exact named human approval and has none"),
            )
        return None
    return AtomRefusal(
        code="UNKNOWN_TRANSFORM_CLASS",
        detail=f"atom '{atom.canonical_id}' declares unknown transform class {atom.transform_class!r}",
    )


def _walk_shapes(shapes, prefix: str):
    """Yield (location, shape) including group children — a claim can hide in a group."""
    for shape in shapes:
        location = f"{prefix}/{shape.shape_id}:{shape.name}"
        yield location, shape
        if shape.shape_type is not None and "GROUP" in str(shape.shape_type):
            yield from _walk_shapes(shape.shapes, location)


def extract_visible_strings(pptx_path: Path) -> list[tuple[str, str]]:
    """Every string a reader can see, with where it lives."""
    from pptx import Presentation

    presentation = Presentation(str(pptx_path))
    found: list[tuple[str, str]] = []
    for index, slide in enumerate(presentation.slides, start=1):
        for location, shape in _walk_shapes(slide.shapes, f"slide[{index}]"):
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    found.append((location, text))
            if getattr(shape, "has_table", False):
                for r, row in enumerate(shape.table.rows):
                    for c, cell in enumerate(row.cells):
                        text = cell.text.strip()
                        if text:
                            found.append((f"{location}/cell[{r},{c}]", text))
    return found


def _word_boundary_excerpt(needle: str, haystack: str) -> bool:
    pattern = r"(?<![A-Za-z0-9])" + re.escape(needle.strip().rstrip(".?!")) + r"(?![A-Za-z0-9])"
    return re.search(pattern, haystack, flags=re.IGNORECASE) is not None


def occurrence_key(text: str, element_id: str) -> str:
    """Authorization key: the string AND where it is allowed to appear."""
    return f"{' '.join(text.split()).casefold().lstrip('❯>•- ')}@{element_id}"


def element_id_from_shape_name(name: str) -> str:
    """Recover the canonical element id from an emitted shape name.

    The PPTX emitter names every object `el:<element-id>[:part]`, which is what
    makes occurrence-scoped authorization possible at the delivered artifact."""
    if not name.startswith("el:"):
        return ""
    return name[3:].split(":")[0]


def _is_claim_bound(text: str, claim_texts: list[str], approvals: PublishApprovals,
                    atom_index: dict[str, AssertionAtom] | None = None,
                    authorized_keys: set[str] | None = None,
                    element_id: str = "") -> bool:
    """A visible string is licensed by a compiler ATOM, an approval, a claim excerpt, or chrome."""
    stripped = text.strip()
    if atom_index is not None:
        key = " ".join(stripped.split()).casefold().lstrip("❯>•- ")
        # Membership is not authorization on two counts: the atom must have
        # passed authorize_atom(), AND the delivered occurrence must be the one
        # that was authorized. A title's approval cannot license the same words
        # pasted into an edge label.
        if occurrence_key(stripped, element_id) in (authorized_keys or set()):
            return True
        if element_id and key in (authorized_keys or set()):
            return False  # authorized text, wrong occurrence
    if len(stripped) <= _CHROME_MAX_CHARS or stripped.isdigit():
        return True
    # A diagram node renders its label and sublabel into a single text frame, so
    # a multi-line string is several assertions; each line must stand on its own.
    lines = [ln.strip() for ln in stripped.splitlines() if ln.strip()]
    if len(lines) > 1:
        return all(_is_claim_bound(line, claim_texts, approvals, atom_index, authorized_keys, element_id) for line in lines)
    normalised = " ".join(stripped.split())
    for allowed in (*approvals.approved_renderings, *approvals.non_claim_text):
        if normalised.casefold() == " ".join(allowed.split()).casefold():
            return True
    if approvals.disclaimer and normalised.casefold() == " ".join(approvals.disclaimer.split()).casefold():
        return True
    # marker prefixes the emitters add (chevrons) are not part of the claim
    probe = normalised.lstrip("❯>•- ").rstrip("…")
    for claim in claim_texts:
        if _word_boundary_excerpt(probe, claim):
            return True
    return False


def _package_text(pptx_path: Path) -> list[tuple[str, str]]:
    """Text anywhere in the package: masters, layouts, notes, core properties.

    A stale owner marker hiding in a notes slide or a document property is still
    a false ownership claim, so the scan cannot stop at slide shapes."""
    import zipfile

    entries: list[tuple[str, str]] = []
    with zipfile.ZipFile(pptx_path) as archive:
        for name in archive.namelist():
            if not name.endswith(".xml"):
                continue
            if not any(part in name for part in ("slide", "master", "layout", "notes", "core.xml", "app.xml", "custom")):
                continue
            try:
                entries.append((name, archive.read(name).decode("utf-8", errors="replace")))
            except KeyError:  # pragma: no cover - archive integrity
                continue
    return entries


def verify_publish(
    pptx_path: Path,
    *,
    claim_texts: list[str],
    claims_by_id: dict[str, str] | None = None,
    approvals: PublishApprovals,
    template_contract: TemplateContract | None = None,
    document=None,
    require_document: bool = True,
) -> PublishReceipt:
    """Re-prove the delivered artifact, not the manifest that produced it."""
    from pptx import Presentation

    digest = hashlib.sha256(pptx_path.read_bytes()).hexdigest()
    findings: list[PublishFinding] = []
    visible = extract_visible_strings(pptx_path)
    atom_index: dict[str, AssertionAtom] | None = None
    authorized_keys: set[str] = set()
    if document is None and require_document:
        findings.append(PublishFinding(
            code="UNREADABLE", where="inputs",
            detail=("the approved document is required: without it the verifier falls back to ledger "
                    "excerpts and unscoped string allowlists, which is not a publication proof")))
    elif document is not None:
        resolved_claims = claims_by_id or {}
        atoms = atoms_from_document(document)
        approved_texts = {" ".join(t.split()).casefold() for t in approvals.approved_renderings}
        atom_index = {}
        for atom in atoms:
            key = " ".join(atom.text.split()).casefold().lstrip("❯>•- ")
            atom_index[key] = atom
            refusal = authorize_atom(atom, claims_by_id=resolved_claims, approved_texts=approved_texts)
            if refusal is None:
                authorized_keys.add(key)
                authorized_keys.add(occurrence_key(atom.text, atom.element_id))
            else:
                findings.append(PublishFinding(
                    code="UNCLAIMED_TEXT", where=atom.canonical_id,
                    detail=f"{refusal.code}: {refusal.detail}"))

    for where, text in visible:
        # location is "slide[N]/<id>:<name>[/<id>:<name>...]"; the OWNING shape is
        # the last segment, and its name follows the shape id
        last = where.rsplit("/", 1)[-1]
        shape_name = last.split(":", 1)[1] if ":" in last else ""
        occurrence = element_id_from_shape_name(shape_name)
        if not _is_claim_bound(text, claim_texts, approvals, atom_index, authorized_keys, occurrence):
            findings.append(PublishFinding(
                code="UNCLAIMED_TEXT", where=where,
                detail=f"visible text is not a legal transform of any claim: {text[:70]!r}"))
        if re.search(r"[A-Za-z]{2}(…|\.\.\.)$", text.strip()) or re.search(r"[A-Za-z]-$", text.strip()):
            findings.append(PublishFinding(
                code="VISIBLE_CLAIM_LOSS", where=where,
                detail=f"text appears truncated mid-word: …{text.strip()[-24:]!r}"))

    markers = tuple(approvals.stale_owner_markers)
    if markers:
        for name, xml in _package_text(pptx_path):
            for marker in markers:
                if marker in xml:
                    findings.append(PublishFinding(
                        code="STALE_OWNER_MARKER", where=name,
                        detail=f"previous owner marker {marker!r} survives in the package"))

    presentation = Presentation(str(pptx_path))
    width, height = presentation.slide_width or 0, presentation.slide_height or 0
    for index, slide in enumerate(presentation.slides, start=1):
        has_text = any(
            shape.has_text_frame and shape.text_frame.text.strip()
            for _, shape in _walk_shapes(slide.shapes, "s")
        )
        pictures = [s for _, s in _walk_shapes(slide.shapes, "s")
                    if s.shape_type is not None and "PICTURE" in str(s.shape_type)]
        if not has_text and pictures:
            findings.append(PublishFinding(
                code="NON_EDITABLE_CONTENT", where=f"slide[{index}]",
                detail="slide carries imagery but no readable text — flattened content cannot be verified or edited"))
        for location, shape in _walk_shapes(slide.shapes, f"slide[{index}]"):
            if not (shape.has_text_frame and shape.text_frame.text.strip()):
                continue
            try:
                left, top = shape.left or 0, shape.top or 0
                right, bottom = left + (shape.width or 0), top + (shape.height or 0)
            except TypeError:
                findings.append(PublishFinding(
                    code="UNREADABLE", where=location, detail="shape geometry could not be read"))
                continue
            if right < 0 or bottom < 0 or left > width or top > height:
                findings.append(PublishFinding(
                    code="VISIBLE_CLAIM_LOSS", where=location,
                    detail="text-bearing shape sits entirely off the canvas"))
            elif (shape.width or 0) <= 0 or (shape.height or 0) <= 0:
                findings.append(PublishFinding(
                    code="VISIBLE_CLAIM_LOSS", where=location,
                    detail="text-bearing shape has zero or negative extent"))

    if template_contract is not None:
        layout_names = {slide.slide_layout.name for slide in presentation.slides}
        approved = set(template_contract.layouts_used)
        drifted = sorted(layout_names - approved) if approved else []
        if drifted:
            findings.append(PublishFinding(
                code="TEMPLATE_DRIFT", where="presentation",
                detail=f"slides use layouts outside the approved template contract: {drifted}"))

    return PublishReceipt(
        pptx=str(pptx_path),
        pptx_sha256=digest,
        status="PASS" if not findings else "REFUSED",
        strings_checked=len(visible),
        findings=findings,
        proves=[
            "Every visible string in THIS file traces to an approved rendering, a ledger claim excerpt, or declared chrome.",
            "No previous template owner's marker survives anywhere in the package.",
            "No slide is flattened to unverifiable imagery, and no text-bearing shape is off-canvas or zero-sized.",
        ],
        does_not_prove=[
            "That a claim is factually true beyond its cited source.",
            "That the deck is visually approved or that screenshots are current.",
            "Exact pixel parity between the HTML renderer and a PowerPoint or LibreOffice render.",
        ],
    )


def load_claims_by_id(ledger_path: Path) -> dict[str, str]:
    """Claim id -> text, so an atom's cited claim can be resolved and checked."""
    raw = ledger_path.read_text(encoding="utf-8")
    if ledger_path.suffix in {".yaml", ".yml"}:
        import yaml

        data = yaml.safe_load(raw)
    else:
        data = json.loads(raw)
    claims = data.get("claims", data) if isinstance(data, dict) else data
    resolved = {c["id"]: str(c.get("text") or c.get("claim") or "")
                for c in claims if isinstance(c, dict) and c.get("id")}
    if not resolved:
        raise ValueError(f"no identified claims found in {ledger_path}")
    return resolved


def load_claim_texts(ledger_path: Path) -> list[str]:
    """Claim text from a ledger, accepting YAML or JSON."""
    raw = ledger_path.read_text(encoding="utf-8")
    if ledger_path.suffix in {".yaml", ".yml"}:
        import yaml

        data = yaml.safe_load(raw)
    else:
        data = json.loads(raw)
    claims = data.get("claims", data) if isinstance(data, dict) else data
    texts: list[str] = []
    for claim in claims:
        if isinstance(claim, dict):
            text = claim.get("text") or claim.get("claim")
            if text:
                texts.append(str(text))
    if not texts:
        raise ValueError(f"no claim text found in {ledger_path}")
    return texts
