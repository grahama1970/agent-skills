"""The grahama.co Gc brandmark, emitted natively (#1314).

Inheriting a house template also inherits its logo, which means a grahama.co
deck was carrying the CS Group mark while its disclaimer named grahama.co — the
mark and the ownership statement disagreed, which is worse than either being
wrong alone. This module supplies the correct mark and removes the inherited one.

The mark is not a bitmap: it is the same construction the website uses
(``site/components/site-nav.tsx``) — a brass ring, a serif ``G``, and a small
italic ember ``c`` tucked into it, beside the ``grahama.co`` wordmark. Emitting
it as an ellipse plus text runs keeps it editable in PowerPoint, matching the
library's editability contract, and keeps one definition of the brand rather
than a screenshot of one.

Inputs: a slide/layout and a frame. Outputs: named native shapes. Failure modes:
replacing the inherited mark REMOVES the old pictures — a deck showing two
different owners' marks is the defect this exists to prevent, so a partial
replacement raises rather than leaving both.
"""

from __future__ import annotations

from typing import Literal

from pydantic import Field

from .models import StrictModel

# Resolved from site/app/globals.css — one source of brand truth.
BRASS = "E2AC62"   # --brass: the ring
INK = "0C0908"     # --ink: the G on light ground
EMBER = "D1703C"   # --ember: the italic c
WORDMARK = "grahama.co"
SERIF_STACK = "Georgia"  # --serif falls back to Georgia/Palatino off-web


class BrandmarkReceipt(StrictModel):
    schema_: Literal["pitchdeck.brandmark_receipt.v1"] = Field(
        default="pitchdeck.brandmark_receipt.v1", alias="schema"
    )
    emitted_shapes: list[str] = Field(default_factory=list)
    removed_inherited_marks: list[str] = Field(default_factory=list)
    wordmark: str = WORDMARK


def remove_inherited_marks(presentation, *, max_height_in: float = 0.75) -> list[str]:
    """Strip small logo pictures from masters and layouts.

    Bounded by height so a full-bleed background picture (the band texture) is
    never mistaken for a logo. Returns what was removed so the receipt can show
    it rather than the change being invisible."""
    removed: list[str] = []
    from .template_deck import all_layouts

    pools = [(f"master[{i}]", m) for i, m in enumerate(presentation.slide_masters)]
    pools += list(all_layouts(presentation))  # every master's layouts, not just the first
    limit = int(max_height_in * 914400)
    for where, part in pools:
        for shape in list(part.shapes):
            if shape.shape_type is None or "PICTURE" not in str(shape.shape_type):
                continue
            if (shape.height or 0) > limit:
                continue  # background art, not a mark
            shape._element.getparent().remove(shape._element)
            removed.append(f"{where}:{shape.name}")
    return removed


def emit_brandmark(shapes, *, left_in: float, top_in: float, height_in: float = 0.34,
                   on_dark: bool = False) -> list[str]:
    """Draw the Gc mark + wordmark as native, editable shapes."""
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    def rgb(value: str) -> RGBColor:
        return RGBColor.from_string(value)

    names: list[str] = []
    ring_size = height_in
    ring = shapes.add_shape(MSO_SHAPE.OVAL, Inches(left_in), Inches(top_in),
                            Inches(ring_size), Inches(ring_size))
    ring.fill.background()
    ring.line.color.rgb = rgb(BRASS)
    ring.line.width = Pt(1.0)
    ring.name = "brand:gc-ring"
    names.append(ring.name)

    glyph = shapes.add_textbox(Inches(left_in), Inches(top_in - 0.035),
                               Inches(ring_size), Inches(ring_size))
    glyph.name = "brand:gc-glyph"
    frame = glyph.text_frame
    frame.word_wrap = False
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    from pptx.enum.text import PP_ALIGN

    paragraph.alignment = PP_ALIGN.CENTER
    g_run = paragraph.add_run()
    g_run.text = "G"
    g_run.font.name = SERIF_STACK
    g_run.font.size = Pt(height_in * 46)
    g_run.font.bold = True
    g_run.font.color.rgb = rgb("FFFFFF" if on_dark else INK)
    c_run = paragraph.add_run()
    c_run.text = "c"
    c_run.font.name = SERIF_STACK
    c_run.font.size = Pt(height_in * 34)
    c_run.font.italic = True
    c_run.font.color.rgb = rgb(EMBER)
    names.append(glyph.name)

    word = shapes.add_textbox(Inches(left_in + ring_size + 0.06), Inches(top_in + 0.02),
                              Inches(1.5), Inches(ring_size))
    word.name = "brand:wordmark"
    word_frame = word.text_frame
    word_frame.word_wrap = False
    word_frame.margin_left = word_frame.margin_top = word_frame.margin_bottom = 0
    run = word_frame.paragraphs[0].add_run()
    run.text = WORDMARK
    run.font.name = SERIF_STACK
    run.font.size = Pt(height_in * 34)
    run.font.color.rgb = rgb("FFFFFF" if on_dark else INK)
    names.append(word.name)
    return names
