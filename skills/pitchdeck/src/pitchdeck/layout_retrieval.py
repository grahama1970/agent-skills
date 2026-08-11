"""RESEARCH-ONLY: nearest-real-slide layout retrieval over the deck corpus.

NOT part of the compiler or the release story (#1335). It has no consumer in the
emission path, and carrying it there meant a Qdrant service dependency, index and
embedding-model drift, and operational surface for zero effect on output. It
stays available behind `index-house-slides` / `find-layout` for offline study,
and may only re-enter the product once a holdout experiment shows retrieval
improves a defined house-conformance measure without reducing determinism.

Original purpose follows.

The house corpus is consistent by construction — the same author, chrome, and
layout habits across five decks — which makes it a pattern-recognition target
rather than something to imitate by hand. This module indexes every REAL slide
(rendered PNG + title text + a measured layout signature) into Qdrant using the
same multimodal embedder ``visual_sync`` already uses, then answers: "which of
the author's real slides is closest to the slide I need?" and returns that
slide's measured geometry to compile against.

Inputs: the PPTX corpus directory and its rendered PNGs. Outputs: LayoutSignature
records (fractional zones, deterministic from the PPTX) and retrieval hits with
cosine scores. Failure modes: a missing render or an unreachable embedder raises
with the endpoint visible — retrieval never silently falls back to invented
geometry, because a made-up layout is exactly the failure this module exists to
prevent.
"""

from __future__ import annotations

import base64
import hashlib
import json
from pathlib import Path
from typing import Literal

import httpx
from pydantic import Field

from .house_spec import EMU_PER_INCH
from .models import StrictModel
from .visual_sync import EMBED_URL, HTTP_TIMEOUT, QDRANT_URL, VECTOR_SIZE

HOUSE_COLLECTION = "pitchdeck_house_slides_v1"


class LayoutBlock(StrictModel):
    kind: Literal["text", "picture", "shape"]
    x: float
    y: float
    w: float
    h: float
    words: int = 0


class LayoutSignature(StrictModel):
    """Measured structure of ONE real slide — the thing worth copying."""

    schema_: Literal["pitchdeck.layout_signature.v1"] = Field(
        default="pitchdeck.layout_signature.v1", alias="schema"
    )
    deck: str
    slide_index: int
    title: str
    blocks: list[LayoutBlock]
    picture_count: int
    text_block_count: int
    total_words: int

    @property
    def structure_key(self) -> str:
        return f"{self.text_block_count}t{self.picture_count}p"


def extract_signatures(decks_dir: Path) -> list[LayoutSignature]:
    """Measure every slide in the corpus (deterministic; no rendering needed)."""
    from pptx import Presentation

    signatures: list[LayoutSignature] = []
    for deck in sorted(decks_dir.glob("*.pptx")):
        presentation = Presentation(str(deck))
        width_in = presentation.slide_width / EMU_PER_INCH
        height_in = presentation.slide_height / EMU_PER_INCH
        for index, slide in enumerate(presentation.slides, start=1):
            blocks: list[LayoutBlock] = []
            title = ""
            for shape in slide.shapes:
                try:
                    x = shape.left / EMU_PER_INCH / width_in
                    y = shape.top / EMU_PER_INCH / height_in
                    w = shape.width / EMU_PER_INCH / width_in
                    h = shape.height / EMU_PER_INCH / height_in
                except (TypeError, AttributeError, ZeroDivisionError):
                    continue
                if not (0.0 <= x <= 1.0 and 0.0 <= y <= 1.0 and w > 0 and h > 0):
                    continue
                is_picture = shape.shape_type is not None and "PICTURE" in str(shape.shape_type)
                text = shape.text_frame.text.strip() if shape.has_text_frame else ""
                if is_picture:
                    kind = "picture"
                elif text:
                    kind = "text"
                else:
                    kind = "shape"
                if kind == "text" and not title and y < 0.25 and 2 <= len(text.split()) <= 14:
                    title = " ".join(text.split())
                blocks.append(
                    LayoutBlock(
                        kind=kind,
                        x=round(x, 4), y=round(y, 4), w=round(w, 4), h=round(h, 4),
                        words=len(text.split()),
                    )
                )
            signatures.append(
                LayoutSignature(
                    deck=deck.stem,
                    slide_index=index,
                    title=title,
                    blocks=blocks,
                    picture_count=sum(1 for b in blocks if b.kind == "picture"),
                    text_block_count=sum(1 for b in blocks if b.kind == "text"),
                    total_words=sum(b.words for b in blocks),
                )
            )
    if not signatures:
        raise ValueError(f"no slides measured under {decks_dir}")
    return signatures


def _embed(client: httpx.Client, *, text: str | None = None, image: Path | None = None) -> list[float]:
    # Service contract (embry-embedding-mm /embed, verified against its OpenAPI
    # schema 2026-08-11): the field is `image` and it accepts a data URL. The
    # previous `image_b64` field was UNKNOWN to the service and silently
    # ignored, so every "image" vector was actually the text embedding — all
    # 250 stored image vectors were identical.
    payload: dict = {}
    if text:
        payload["text"] = text
    if image is not None:
        payload["image"] = "data:image/png;base64," + base64.b64encode(image.read_bytes()).decode("ascii")
    response = client.post(EMBED_URL, json=payload, timeout=HTTP_TIMEOUT)
    response.raise_for_status()
    body = response.json()
    vector = body.get("embedding") or body.get("vector") or body.get("embeddings")
    if isinstance(vector, list) and vector and isinstance(vector[0], list):
        vector = vector[0]
    if not isinstance(vector, list) or len(vector) != VECTOR_SIZE:
        raise ValueError(f"embedder returned an unusable vector for {EMBED_URL}: keys={sorted(body)}")
    return vector


def index_house_slides(decks_dir: Path, renders_dir: Path, *, memory: bool = True) -> dict:
    """Index every real slide for retrieval AND recall.

    Three artifacts per page: a JSON record (full text + layout signature) on
    disk, a Qdrant point with REAL text and image vectors, and a /memory
    document holding metadata plus the Qdrant point id (never vectors — the
    established visual_sync pattern). Renders are matched through the
    pdftotext-derived page mapping, because pdftoppm page numbers diverge from
    pptx slide indexes whenever a deck contains hidden slides."""
    signatures = extract_signatures(decks_dir)
    mapping_path = renders_dir / "page-mapping.json"
    page_map = json.loads(mapping_path.read_text()) if mapping_path.is_file() else {}
    records_dir = renders_dir / "records"
    records_dir.mkdir(exist_ok=True)
    points: list[dict] = []
    memory_docs: list[dict] = []
    missing_renders: list[str] = []
    with httpx.Client(timeout=HTTP_TIMEOUT) as client:
        existing = client.get(f"{QDRANT_URL}/collections/{HOUSE_COLLECTION}")
        if existing.status_code != 200:
            client.put(
                f"{QDRANT_URL}/collections/{HOUSE_COLLECTION}",
                json={"vectors": {
                    "text_mm": {"size": VECTOR_SIZE, "distance": "Cosine"},
                    "image_mm": {"size": VECTOR_SIZE, "distance": "Cosine"},
                }},
            ).raise_for_status()
        for signature in signatures:
            page = page_map.get(f"{signature.deck}#{signature.slide_index}")
            if page is None:
                missing_renders.append(f"{signature.deck}#{signature.slide_index} (no page mapping)")
                continue
            image = renders_dir / f"{signature.deck}-{page:02d}.png"
            if not image.is_file():
                image = renders_dir / f"{signature.deck}-{page}.png"
            if not image.is_file():
                missing_renders.append(f"{signature.deck}#{signature.slide_index} (page {page} not rendered)")
                continue
            record = {
                **signature.model_dump(by_alias=True, mode="json"),
                "image_path": str(image),
                "image_sha256": hashlib.sha256(image.read_bytes()).hexdigest(),
                "pdf_page": page,
                "all_text": [b.model_dump(mode="json") for b in signature.blocks if b.kind == "text"],
            }
            record_path = records_dir / f"{signature.deck}-{signature.slide_index:03d}.json"
            record_path.write_text(json.dumps(record, indent=1), encoding="utf-8")
            point_id = int(hashlib.sha256(f"{signature.deck}:{signature.slide_index}".encode()).hexdigest()[:12], 16)
            vectors = {"image_mm": _embed(client, image=image)}
            if signature.title:
                vectors["text_mm"] = _embed(client, text=signature.title)
            points.append({"id": point_id, "vector": vectors,
                           "payload": {**record, "record_path": str(record_path)}})
            memory_docs.append({
                "_key": f"house_slide_{point_id}",
                "problem": f"House slide: {signature.deck} slide {signature.slide_index} ({signature.title or 'untitled'})",
                "solution": (f"Layout record at {record_path}; render at {image}; "
                             f"Qdrant point {point_id} in {HOUSE_COLLECTION} (text_mm + image_mm)."),
                "deck": signature.deck, "slide_index": signature.slide_index,
                "title": signature.title, "image_path": str(image),
                "record_path": str(record_path),
                "visual_qdrant_collection": HOUSE_COLLECTION,
                "visual_qdrant_point_id": point_id,
                "tags": ["pitchdeck", "house-slide", signature.deck],
            })
        for batch_start in range(0, len(points), 32):
            client.put(
                f"{QDRANT_URL}/collections/{HOUSE_COLLECTION}/points",
                json={"points": points[batch_start:batch_start + 32]},
            ).raise_for_status()
    memory_synced = 0
    if memory and memory_docs:
        from .visual_sync import MEMORY_URL

        with httpx.Client(timeout=HTTP_TIMEOUT) as client:
            response = client.post(f"{MEMORY_URL}/upsert",
                                   json={"collection": "pitchdeck_house_slides", "documents": memory_docs})
            response.raise_for_status()
            memory_synced = len(memory_docs)
    return {
        "collection": HOUSE_COLLECTION,
        "indexed": len(points),
        "measured": len(signatures),
        "memory_synced": memory_synced,
        "records_dir": str(records_dir),
        "missing_renders": missing_renders,
    }


def find_nearest_layout(query_text: str, *, limit: int = 3, structure: str | None = None) -> list[dict]:
    """Retrieve the author's real slides closest to a described need.

    ``structure`` (e.g. "2t1p") filters to slides with the same block census,
    so a proof slide retrieves proof-shaped layouts rather than merely
    topically similar ones."""
    with httpx.Client(timeout=HTTP_TIMEOUT) as client:
        vector = _embed(client, text=query_text)
        body: dict = {"vector": {"name": "text_mm", "vector": vector}, "limit": limit, "with_payload": True}
        response = client.post(f"{QDRANT_URL}/collections/{HOUSE_COLLECTION}/points/search", json=body)
        response.raise_for_status()
        hits = response.json()["result"]
    results = []
    for hit in hits:
        payload = hit["payload"]
        if structure and f"{payload['text_block_count']}t{payload['picture_count']}p" != structure:
            continue
        results.append({
            "score": round(hit["score"], 4),
            "deck": payload["deck"],
            "slide_index": payload["slide_index"],
            "title": payload["title"],
            "structure": f"{payload['text_block_count']}t{payload['picture_count']}p",
            "image_path": payload["image_path"],
            "blocks": payload["blocks"],
        })
    return results
