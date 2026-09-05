"""Presentation styling from the same ThemeTokens the browser consumes.

No geometry or text is changed. Raster assets keep their authored colors.
"""
from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.shapes import MSO_SHAPE


def export_font(name):
    # system-ui is a browser generic, not a font installed in Office.
    return 'Arial' if name == 'system-ui' else name


def apply_presentation_theme(presentation, tokens, heading_texts=()):
    if tokens.canvas is None:
        return
    headings = {t.strip().casefold() for t in heading_texts}
    rgb = lambda s: RGBColor.from_string(s.lstrip('#'))
    for slide in presentation.slides:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb(tokens.canvas)
        colors = {'065E7C': tokens.accent, '076889': tokens.accent, '36D6E7': tokens.accent, '22D3EE': tokens.accent, '08131F': tokens.canvas, '102437': tokens.canvas, '132D43': tokens.canvas, 'AFC2CF': tokens.muted, '595959': tokens.muted}
        for node in slide._element.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr'):
            if node.get('val', '').upper() in colors: node.set('val', colors[node.get('val').upper()].lstrip('#'))
        band = next((s for s in slide.shapes if s.name == 'chrome:band'), None)
        if band is None:
            band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, presentation.slide_width, int(presentation.slide_height * .10))
            band.name = 'chrome:band'
            # Background only, behind existing slide objects.
            tree = slide.shapes._spTree
            tree.remove(band._element); tree.insert(2, band._element)
        band.fill.solid(); band.fill.fore_color.rgb = rgb(tokens.header); band.line.fill.background()
        color = band._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        for old in list(color): color.remove(old)
        color.append(color.makeelement('{http://schemas.openxmlformats.org/drawingml/2006/main}alpha', {'val': str(round(tokens.header_opacity * 100000))}))
        # The supplied decks use this exact turbine image at 10% opacity.
        # Image alpha is independent of the solid header fill's alpha.
        texture_path = Path(__file__).resolve().parents[3] / 'best-practices-slide-design/assets/house-band-texture.png'
        texture_bytes = texture_path.read_bytes()
        header_images = []
        a = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
        r = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'
        for shape in slide.shapes:
            blip = shape._element.find(f'.//{a}blip')
            if blip is not None and shape.top < presentation.slide_height * .12 and shape.height < presentation.slide_height * .3:
                rid = blip.get(f'{r}embed')
                if rid and slide.part.related_part(rid).blob == texture_bytes:
                    header_images.append(blip)
        if not header_images:
            image = slide.shapes.add_picture(str(texture_path), int(-.09 * presentation.slide_width), int(-.01 * presentation.slide_height), int(1.18 * presentation.slide_width), int(.11 * presentation.slide_height))
            image.name = 'chrome:band-texture'
            tree = slide.shapes._spTree
            tree.remove(image._element); tree.insert(tree.index(band._element) + 1, image._element)
            header_images.append(image._element.find(f'.//{a}blip'))
        for blip in header_images:
            for old in list(blip.findall(f'{a}alphaModFix')): blip.remove(old)
            blip.append(blip.makeelement(f'{a}alphaModFix', {'amt': str(round(tokens.header_image_opacity * 100000))}))
        def walk(shapes):
            for shape in shapes:
                if shape.name == 'chrome:band-texture':
                    continue
                if hasattr(shape, 'shapes'): walk(shape.shapes)
                if not shape.has_text_frame: continue
                heading = shape.name in {'el:title', 'chrome:band-title'} or shape.text.strip().casefold() in headings
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name not in {'Consolas', 'Courier New'}:
                            run.font.name = export_font(tokens.heading_font if heading else tokens.body_font)
                        original = str(run.font.color.rgb).lower() if run.font.color.type == MSO_COLOR_TYPE.RGB else ''
                        color = tokens.muted if original == tokens.muted.lstrip('#').lower() or shape.name == 'chrome:page-number' else tokens.text
                        if shape.name == 'chrome:identity-wordmark': color = tokens.accent
                        if heading and shape.top < presentation.slide_height * .12: color = tokens.header_text
                        run.font.color.rgb = rgb(color)
        walk(slide.shapes)
