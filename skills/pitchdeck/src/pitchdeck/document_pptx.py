"""Native-editable PPTX emitter for deck_document.v1 primitives (#1271 phase 3).

Renders the scene graph as REAL PowerPoint objects — nested grpSp, preset sp,
freeform sp, cxnSp connectors, text runs, and pic — never a raster. Group
strategy (validated live against python-pptx 1.x): children are placed at
absolute canvas EMU inside the group (python-pptx recalculates the group's
off/ext/chOff/chExt from them), then the group xfrm is overridden to the
AUTHORED bbox with chOff/chExt pinned identical to off/ext. That identity
mapping means no rescaling ever occurs and a padded child_frame survives as
literal geometry (margin inside the group) rather than being swallowed by
extent recalculation.

Capability decisions are receipts, not silences: connector attach hints are
NOT written as stCxn/endCxn (python-pptx support is experimental); local
links render as link-styled runs without a jump action. Icons resolve
fail-closed via the hash-pinned library (editable shape trees required).
Every emitted object is named ``el:<element-id>`` so reimport can prove
identity. Failure modes: unknown kinds, unresolvable icons, or assets that
cannot be read raise — nothing emits partially.
"""

from __future__ import annotations

import json
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Emu, Inches, Pt

from .document import DeckDocument, DocElement, DocElementKind
from .document_html import _load_theme, house_title_case
from .io import expand_path

SLIDE_W_IN, SLIDE_H_IN = 13.333, 7.5

_PRESETS = {
    "rect": MSO_SHAPE.RECTANGLE,
    "rounded_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
    "ellipse": MSO_SHAPE.OVAL,
    "chevron": MSO_SHAPE.CHEVRON,
    "pill": MSO_SHAPE.ROUNDED_RECTANGLE,  # adjustment 0.5 -> capsule
    "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
}

_CONNECTOR_ROUTES = {
    "straight": MSO_CONNECTOR.STRAIGHT,
    "bent": MSO_CONNECTOR.ELBOW,
    "curved": MSO_CONNECTOR.CURVE,
}

_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _scrub_package_markers(pptx_path: Path, markers: tuple[str, ...], replacement: str) -> list[str]:
    """Remove stale owner strings from textual OPC parts after python-pptx save.

    python-pptx leaves docProps/app.xml from the template intact; that metadata
    can still carry prior-owner slide titles after all template slides are
    stripped. Rewrite only UTF-8 XML/rels parts and leave binary media alone.
    """
    if not markers:
        return []
    changed: list[str] = []
    with zipfile.ZipFile(pptx_path, "r") as src, tempfile.NamedTemporaryFile(
        dir=pptx_path.parent, suffix=".pptx", delete=False
    ) as tmp:
        tmp_path = Path(tmp.name)
        with zipfile.ZipFile(tmp, "w") as dst:
            for info in src.infolist():
                data = src.read(info.filename)
                if info.filename.endswith((".xml", ".rels")):
                    text = data.decode("utf-8")
                    updated = text
                    for marker in markers:
                        updated = updated.replace(marker, replacement)
                    if updated != text:
                        data = updated.encode("utf-8")
                        changed.append(info.filename)
                dst.writestr(info, data)
    tmp_path.replace(pptx_path)
    return changed


class Frame:
    """Absolute inches context: local fractions map into this rectangle."""

    def __init__(self, x: float, y: float, w: float, h: float):
        self.x, self.y, self.w, self.h = x, y, w, h

    def rect(self, bbox) -> tuple:
        return (
            Inches(self.x + bbox.x * self.w),
            Inches(self.y + bbox.y * self.h),
            Inches(bbox.w * self.w),
            Inches(bbox.h * self.h),
        )

    def point(self, x: float, y: float) -> tuple:
        return Inches(self.x + x * self.w), Inches(self.y + y * self.h)

    def sub(self, bbox, child_frame=None) -> "Frame":
        base = Frame(self.x + bbox.x * self.w, self.y + bbox.y * self.h, bbox.w * self.w, bbox.h * self.h)
        if child_frame is None:
            return base
        return Frame(
            base.x + child_frame.x * base.w,
            base.y + child_frame.y * base.h,
            child_frame.w * base.w,
            child_frame.h * base.h,
        )


def _hex(value: str) -> RGBColor:
    return RGBColor.from_string(value.lstrip("#").upper())


def _set_dash(line, dash: str) -> None:
    if dash == "solid":
        return
    from pptx.enum.dml import MSO_LINE_DASH_STYLE

    line.dash_style = MSO_LINE_DASH_STYLE.DASH if dash == "dashed" else MSO_LINE_DASH_STYLE.ROUND_DOT


def _add_arrowheads(connector, *, start: bool, end: bool) -> None:
    ln = connector.line._get_or_add_ln()
    for present, tag in ((start, "headEnd"), (end, "tailEnd")):
        if present:
            el = ln.makeelement(f"{{{_A}}}{tag}", {"type": "triangle", "w": "med", "len": "med"})
            ln.append(el)


def _pin_group_xfrm(group, bbox_rect: tuple) -> None:
    """Override off/ext AND chOff/chExt to the authored bbox (identity map):
    padding inside the group survives as literal geometry."""
    left, top, width, height = (int(v) for v in bbox_rect)
    xfrm = group._element.find(f".//{{{_A}}}xfrm")
    for tag, attrs in (
        ("off", {"x": str(left), "y": str(top)}),
        ("ext", {"cx": str(width), "cy": str(height)}),
        ("chOff", {"x": str(left), "y": str(top)}),
        ("chExt", {"cx": str(width), "cy": str(height)}),
    ):
        node = xfrm.find(f"{{{_A}}}{tag}")
        for key, value in attrs.items():
            node.set(key, value)


def _emit_icon_parts(group, library_id: str, icon_frame: Frame, tint, name_prefix: str) -> None:
    """Emit a library icon's native mapping (presets + straight freeforms) into a group."""
    from .icons import resolve_icon

    resolved = resolve_icon(library_id, require_editable=True)
    for index, part in enumerate(resolved["mapping"]["parts"]):
        if part["kind"] == "preset":
            x, y, w, h = part["bbox"]
            shape = group.shapes.add_shape(
                _PRESETS[part["preset"]],
                Inches(icon_frame.x + x * icon_frame.w), Inches(icon_frame.y + y * icon_frame.h),
                Inches(w * icon_frame.w), Inches(h * icon_frame.h),
            )
            shape.fill.background()
            shape.line.color.rgb = tint
            shape.line.width = Pt(2.0)
            shape.name = f"{name_prefix}:part{index}"
        else:
            vertices = [(Emu(Inches(icon_frame.x + vx * icon_frame.w)), Emu(Inches(icon_frame.y + vy * icon_frame.h))) for vx, vy in part["vertices"]]
            builder = group.shapes.build_freeform(vertices[0][0], vertices[0][1], scale=1)
            builder.add_line_segments(vertices[1:], close=part.get("closed", False))
            shape = builder.convert_to_shape()
            shape.fill.background()
            shape.line.color.rgb = tint
            shape.line.width = Pt(2.0)
            shape.name = f"{name_prefix}:part{index}"


def _emit_icon(container, element: DocElement, frame: Frame, palette: dict, receipt: dict) -> None:
    from .icons import resolve_icon

    resolved = resolve_icon(element.icon.library_id, require_editable=True)
    tint = _hex(palette.get(element.icon.tint_role, palette["primary"]))
    group = container.add_group_shape()
    _emit_icon_parts(group, element.icon.library_id, frame.sub(element.bbox), tint, f"el:{element.id}")
    _pin_group_xfrm(group, frame.rect(element.bbox))
    group.name = f"el:{element.id}"
    receipt["icons"].append({
        "element": element.id,
        "library_id": element.icon.library_id,
        "representation": resolved["representation"],
        "svg_sha256": resolved["svg_sha256"],
    })


def _emit_rich_text(container, element: DocElement, frame: Frame, palette: dict, scale: dict, receipt: dict) -> None:
    left, top, width, height = frame.rect(element.bbox)
    box = container.add_textbox(left, top, width, height)
    box.name = f"el:{element.id}"
    tf = box.text_frame
    tf.word_wrap = True
    for b_index, block in enumerate(element.rich_text.blocks):
        paragraph = tf.paragraphs[0] if b_index == 0 else tf.add_paragraph()
        if block.bullet_level:
            paragraph.level = block.bullet_level - 1
        for run_spec in block.runs:
            run = paragraph.add_run()
            run.text = run_spec.text
            run.font.size = Pt(scale.get(block.style_role.value, 20))
            for mark in run_spec.marks:
                if mark.type == "bold":
                    run.font.bold = True
                elif mark.type == "italic":
                    run.font.italic = True
                elif mark.type == "underline":
                    run.font.underline = True
                elif mark.type == "code":
                    run.font.name = "Consolas"
                elif mark.type == "color":
                    run.font.color.rgb = _hex(palette.get(mark.role, palette["ink"]))
                elif mark.type == "link":
                    run.font.color.rgb = _hex(palette["primary"])
                    run.font.underline = True
                    receipt["capability_decisions"].append(
                        f"{element.id}: local link to '{mark.target_slide_id}' rendered as styled run (no jump action)"
                    )


def _emit_element(container, element: DocElement, frame: Frame, *, palette: dict, scale: dict, assets: dict, asset_base: Path, receipt: dict) -> None:
    kind = element.kind
    if kind is DocElementKind.GROUP:
        group = container.add_group_shape()
        inner = frame.sub(element.bbox, element.child_frame)
        ordered = sorted(enumerate(element.children or []), key=lambda pair: (pair[1].z, pair[0]))
        for _, child in ordered:
            _emit_element(group.shapes, child, inner, palette=palette, scale=scale, assets=assets, asset_base=asset_base, receipt=receipt)
        _pin_group_xfrm(group, frame.rect(element.bbox))
        if element.rotation_deg:
            group.rotation = element.rotation_deg
        group.name = f"el:{element.id}"
        return
    if kind is DocElementKind.SHAPE:
        left, top, width, height = frame.rect(element.bbox)
        shape = container.add_shape(_PRESETS[element.shape.preset.value], left, top, width, height)
        if element.shape.preset.value == "pill":
            shape.adjustments[0] = 0.5
        if element.shape.fill_role:
            shape.fill.solid()
            shape.fill.fore_color.rgb = _hex(palette.get(element.shape.fill_role, "#FFFFFF"))
        else:
            shape.fill.background()
        if element.shape.stroke:
            shape.line.color.rgb = _hex(palette.get(element.shape.stroke.role, palette["primary"]))
            shape.line.width = Pt(element.shape.stroke.width_pt)
            _set_dash(shape.line, element.shape.stroke.dash)
        else:
            shape.line.fill.background()
        if element.rotation_deg:
            shape.rotation = element.rotation_deg
        shape.name = f"el:{element.id}"
        return
    if kind is DocElementKind.LINE:
        spec = element.line
        bx, by = frame.point(spec.start.x, spec.start.y)
        ex, ey = frame.point(spec.end.x, spec.end.y)
        connector = container.add_connector(_CONNECTOR_ROUTES[spec.route], bx, by, ex, ey)
        connector.line.color.rgb = _hex(palette["primary"])
        connector.line.width = Pt(spec.width_pt)
        _set_dash(connector.line, spec.dash)
        _add_arrowheads(connector, start=spec.arrow_start, end=spec.arrow_end)
        if spec.start_hint or spec.end_hint:
            receipt["capability_decisions"].append(
                f"{element.id}: attach hints present but stCxn/endCxn NOT written (python-pptx support experimental); geometry authoritative"
            )
        connector.name = f"el:{element.id}"
        return
    if kind is DocElementKind.ICON:
        if element.role == "badge":
            tint = _hex(palette.get(element.icon.tint_role, "#FFFFFF"))
            group = container.add_group_shape()
            badge_frame = frame.sub(element.bbox)
            # Corpus evidence (cybersummit-32/47, sbir-38, reqml-49): the band
            # badge is a BARE white line-art glyph — no enclosing ring — and it
            # fills the badge box rather than sitting inside one.
            inner = Frame(badge_frame.x + badge_frame.w * 0.06, badge_frame.y + badge_frame.h * 0.06,
                          badge_frame.w * 0.88, badge_frame.h * 0.88)
            _emit_icon_parts(group, element.icon.library_id, inner, tint, f"el:{element.id}")
            _pin_group_xfrm(group, frame.rect(element.bbox))
            group.name = f"el:{element.id}"
            return
        _emit_icon(container, element, frame, palette, receipt)
        return
    if kind is DocElementKind.RICH_TEXT:
        _emit_rich_text(container, element, frame, palette, scale, receipt)
        return
    if kind is DocElementKind.TEXT:
        left, top, width, height = frame.rect(element.bbox)
        box = container.add_textbox(left, top, width, height)
        box.name = f"el:{element.id}"
        tf = box.text_frame
        tf.word_wrap = True
        style = element.style
        text = element.text or ""
        if element.role in {"chevrons", "callout"}:
            for l_index, line in enumerate(text.split("\n")):
                if not line.strip():
                    continue
                paragraph = tf.paragraphs[0] if l_index == 0 else tf.add_paragraph()
                body = line
                if line.startswith("> "):
                    marker = paragraph.add_run()
                    marker.text = "\u276F  "
                    marker.font.bold = True
                    marker.font.size = Pt((style.size_pt if style else 20) + 2)
                    marker.font.color.rgb = _hex(palette["primary"])
                    body = line[2:]
                run = paragraph.add_run()
                run.text = body
                run.font.size = Pt(style.size_pt if style else 20)
                run.font.color.rgb = _hex("#595959")
            return
        run = tf.paragraphs[0].add_run()
        run.text = text
        run.font.size = Pt(style.size_pt if style else 20)
        run.font.bold = bool(style and style.bold)
        if style and style.color:
            run.font.color.rgb = _hex(style.color)
        if style and style.align in {"center", "right"}:
            from pptx.enum.text import PP_ALIGN as _PA
            tf.paragraphs[0].alignment = _PA.CENTER if style.align == "center" else _PA.RIGHT
        return
    if kind is DocElementKind.IMAGE:
        spec = assets[element.asset_id]
        path = expand_path(spec.local_path, base_dir=asset_base)
        left, top, width, height = frame.rect(element.bbox)
        source: object = str(path)
        if path.suffix.lower() == ".svg":
            # python-pptx cannot embed SVG; rasterize with rsvg-convert (2x for
            # legibility). The browser deck keeps the vector/animated original.
            import io
            import shutil as _sh
            import subprocess as _sp

            tool = _sh.which("rsvg-convert")
            if not tool:
                raise ValueError(f"element '{element.id}': SVG asset needs rsvg-convert to rasterize for PPTX")
            raster = _sp.run([tool, "-w", str(int(width / 914400 * 192)), str(path)], capture_output=True, timeout=60)
            if raster.returncode or not raster.stdout:
                raise ValueError(f"element '{element.id}': rsvg-convert failed: {raster.stderr.decode(errors='replace')[-300:]}")
            source = io.BytesIO(raster.stdout)
        elif path.suffix.lower() in {".webp", ".avif"}:
            # python-pptx accepts BMP/GIF/JPEG/PNG/TIFF/WMF; convert losslessly.
            import io

            from PIL import Image

            buffer = io.BytesIO()
            Image.open(path).convert("RGB").save(buffer, format="PNG")
            buffer.seek(0)
            source = buffer
        picture = container.add_picture(source, left, top, width=width, height=height)
        if element.crop:
            c = element.crop
            picture.crop_left, picture.crop_top = c.x, c.y
            picture.crop_right, picture.crop_bottom = 1.0 - (c.x + c.w), 1.0 - (c.y + c.h)
        if element.rotation_deg:
            picture.rotation = element.rotation_deg
        picture.name = f"el:{element.id}"
        return
    if kind is DocElementKind.DIAGRAM:
        _emit_diagram(container, element, frame, palette=palette, receipt=receipt)
        return
    raise ValueError(f"element '{element.id}': PPTX emitter has no handler for kind '{kind.value}'")


def _emit_diagram(container, element: DocElement, frame: Frame, *, palette: dict, receipt: dict) -> None:
    """DiagramGraph compiles to the SAME primitive contract (#1271 review):
    a native group of rounded-rect nodes, icon circles, label textboxes, and
    cxnSp connectors — separately editable, never one raster."""
    graph = element.diagram
    primary = _hex(palette["primary"])
    ink = _hex(palette["ink"])
    role_cycle = ["#065E7C", "#6F8E30", "#26558E", "#D6A300", "#065E7C"]
    unboxed = graph.recipe in {"pipeline", "scene"}
    scene = graph.recipe == "scene"  # multi-element illustration (#1315)
    group = container.add_group_shape()
    dframe = frame.sub(element.bbox)
    centers: dict[str, tuple[float, float, float, float]] = {}
    for n_index, node in enumerate(graph.nodes):
        nx = dframe.x + node.bbox.x * dframe.w
        ny = dframe.y + node.bbox.y * dframe.h
        nw = node.bbox.w * dframe.w
        nh = node.bbox.h * dframe.h
        centers[node.id] = (nx, ny, nw, nh)
        accent = _hex(role_cycle[n_index % len(role_cycle)]) if unboxed else primary
        terminal = unboxed and n_index == len(graph.nodes) - 1
        if scene:
            pass  # scene glyphs are BARE: the author never rings a scene element
        elif not unboxed:
            box = group.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(nx), Inches(ny), Inches(nw), Inches(nh))
            box.fill.background()
            box.line.color.rgb = primary
            box.line.width = Pt(2.5)
            box.name = f"el:{element.id}:node:{node.id}"
        else:
            _scale = node.scale if scene else (1.22 if terminal else (1.0 if n_index % 2 == 0 else 0.86))
            ring_size = min(nw, nh) * 0.48 * _scale
            ring = group.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(nx + nw / 2 - ring_size / 2), Inches(ny + nh * 0.28 - ring_size / 2),
                Inches(ring_size), Inches(ring_size))
            ring.fill.background()
            ring.line.color.rgb = accent
            ring.line.width = Pt(3.25 if terminal else 2.25)
            ring.name = f"el:{element.id}:node:{node.id}"
        if node.icon:
            _iscale = node.scale if scene else ((1.22 if terminal else (1.0 if n_index % 2 == 0 else 0.86)) if unboxed else 1.0)
            size = min(nw, nh) * (0.92 if scene else (0.32 if unboxed else 0.36) * _iscale)
            _emit_icon_parts(
                group, node.icon,
                Frame(nx + nw / 2 - size / 2, ny + nh * (0.5 if scene else (0.28 if unboxed else 0.32)) - size / 2, size, size),
                accent, f"el:{element.id}:node:{node.id}:icon",
            )
        if scene and not node.label.strip():
            continue  # a bare supporting glyph carries no label box
        if scene and len(node.label.split()) >= 3:
            # speech-bubble treatment (replication probe, 'LLMs Hallucinate'):
            # bordered rounded rect, body-size text — the bubble IS the content
            bw, bh = nw * 1.5, nh * 0.42
            bubble = group.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(nx + nw / 2 - bw / 2), Inches(ny + nh * 1.02),
                Inches(bw), Inches(bh))
            bubble.fill.solid(); bubble.fill.fore_color.rgb = _hex("#FFFFFF")
            bubble.line.color.rgb = accent; bubble.line.width = Pt(1.5)
            bubble.shadow.inherit = False
            bubble.name = f"el:{element.id}:node:{node.id}:bubble"
            btf = bubble.text_frame; btf.word_wrap = True
            brun = btf.paragraphs[0].add_run(); brun.text = node.label
            brun.font.size = Pt(14); brun.font.color.rgb = _hex("#292929")
            from pptx.enum.text import PP_ALIGN as _A2
            btf.paragraphs[0].alignment = _A2.CENTER
            receipt["slides"] and None
            continue
        label = group.shapes.add_textbox(Inches(nx - nw * 0.25), Inches(ny + nh * (1.06 if scene else (0.62 if unboxed else 0.52))), Inches(nw * 1.5), Inches(nh * (0.5 if scene else (0.38 if unboxed else 0.44))))
        label.name = f"el:{element.id}:node:{node.id}:label"
        tf = label.text_frame
        tf.word_wrap = True
        run = tf.paragraphs[0].add_run()
        run.text = node.label
        run.font.bold = True
        run.font.size = Pt(11 if scene else (13 if unboxed else 14))
        run.font.color.rgb = accent
        from pptx.enum.text import PP_ALIGN

        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        if node.sublabel:
            para = tf.add_paragraph()
            para.alignment = PP_ALIGN.CENTER
            sub = para.add_run()
            sub.text = node.sublabel
            sub.font.size = Pt(10)
            sub.font.color.rgb = ink
    for edge in graph.edges:
        sx, sy, sw, sh = centers[edge.source]
        tx, ty, tw, th = centers[edge.target]
        edge_frac = 0.5 if scene else (0.28 if unboxed else 0.5)
        if scene:
            # centre-to-centre: a scene's flow associates glyphs wherever they
            # sit, unlike a pipeline's strict left-to-right hand-off.
            begin = (Inches(sx + sw / 2), Inches(sy + sh / 2))
            end = (Inches(tx + tw / 2), Inches(ty + th / 2))
        else:
            begin = (Inches(sx + sw), Inches(sy + sh * edge_frac))
            end = (Inches(tx), Inches(ty + th * edge_frac))
        route = MSO_CONNECTOR.ELBOW if edge.route == "curve" else MSO_CONNECTOR.STRAIGHT
        connector = group.shapes.add_connector(route, begin[0], begin[1], end[0], end[1])
        connector.line.color.rgb = primary
        connector.line.width = Pt(2.5)
        _set_dash(connector.line, "dashed" if edge.line_style == "dashed" else ("dotted" if (edge.line_style == "dotted" or (unboxed and edge.line_style == "solid")) else "solid"))
        _add_arrowheads(connector, start=False, end=edge.arrowhead)
        connector.name = f"el:{element.id}:edge:{edge.id}"
        if edge.label:
            mid_x = (sx + sw + tx) / 2
            mid_y = (sy + sh / 2 + ty + th / 2) / 2
            caption = group.shapes.add_textbox(Inches(mid_x - 1.2), Inches(mid_y - 0.42), Inches(2.4), Inches(0.32))
            caption.name = f"el:{element.id}:edge:{edge.id}:label"
            run = caption.text_frame.paragraphs[0].add_run()
            run.text = edge.label
            run.font.italic = True
            run.font.size = Pt(11)
            run.font.color.rgb = ink
    _pin_group_xfrm(group, frame.rect(element.bbox))
    group.name = f"el:{element.id}"


def emit_document_pptx(
    document: DeckDocument,
    output_path: Path,
    *,
    asset_base: Path,
    theme_template: Path | None = None,
    house_template: Path | None = None,
    disclaimer: str | None = None,
    stale_owner_markers: tuple[str, ...] = (),
    brandmark: bool = False,
) -> dict:
    from .publish_gate import assert_publishable

    assert_publishable(document)
    theme = _load_theme(theme_template)
    palette = theme["palette"]
    scale = theme.get("type_scale_pt", {"body": 20, "support": 16, "title": 28})
    assets = {a.id: a for a in document.assets}
    # Template-as-base: opening the author's own deck inherits theme, master,
    # and layouts, so band/texture/logo/footer/page-number are correct by
    # construction instead of being measured and redrawn.
    layout_profile = None
    pending_disclaimer_receipt = None
    pending_brandmark_removed: list[str] | None = None
    if house_template is not None:
        from .template_deck import open_stripped_template, profile_template

        presentation = open_stripped_template(house_template)
        layout_profile = profile_template(house_template)
        # A house template LOCKS its disclaimer by design, so inheriting one
        # onto a different owner's deck silently asserts the wrong ownership.
        # Retarget it, then FAIL CLOSED if any stale owner marker survives —
        # a leftover marker is a false legal claim, not a cosmetic defect.
        if disclaimer:
            from .template_deck import retarget_disclaimer

            markers = stale_owner_markers or ("CSINC", "CS Communication", "SpartaAI", "Sparta AI")
            disclaimer_receipt = retarget_disclaimer(presentation, disclaimer, stale_markers=markers)
            if disclaimer_receipt["residual_markers"]:
                raise ValueError(
                    "STALE_OWNER_DISCLAIMER: the template's original owner still appears at "
                    f"{disclaimer_receipt['residual_markers']} — refusing to emit a deck that "
                    "asserts ownership by someone other than its owner"
                )
            pending_disclaimer_receipt = disclaimer_receipt
        if brandmark:
            # The inherited logo is the TEMPLATE owner's mark. Leaving it beside a
            # retargeted disclaimer makes the deck name two different owners, so
            # the old marks are removed and the grahama.co Gc mark takes its place.
            from .brandmark import remove_inherited_marks

            pending_brandmark_removed = remove_inherited_marks(presentation)
        blank = next((presentation.slide_layouts[l.index] for l in layout_profile.layouts
                      if l.role == "blank"), presentation.slide_layouts[-1])
    else:
        presentation = Presentation()
        presentation.slide_width = Inches(SLIDE_W_IN)
        presentation.slide_height = Inches(SLIDE_H_IN)
        blank = presentation.slide_layouts[6]
    receipt: dict = {
        "schema": "pitchdeck.pptx_primitive_receipt.v1",
        "capability_decisions": [],
        "icons": [],
        "slides": [],
    }
    if pending_disclaimer_receipt is not None:
        receipt["disclaimer"] = pending_disclaimer_receipt
    if pending_brandmark_removed is not None:
        receipt["brandmark"] = {"removed_inherited_marks": pending_brandmark_removed}
    root = Frame(0.0, 0.0, SLIDE_W_IN, SLIDE_H_IN)
    band_cfg = theme.get("chrome", {}).get("header_band", {})
    for slide_doc in document.slides:
        if slide_doc.hidden:
            continue
        inherited_chrome = False
        if layout_profile is not None:
            from .template_deck import drop_unused_placeholders, fill_title, pick_layout

            house_layout = pick_layout(layout_profile, prefer_named="Graham", needs_body=True)
            slide = presentation.slides.add_slide(presentation.slide_layouts[house_layout.index])
            title_el = next((e for e in slide_doc.elements if e.role == "title"), None)
            assertion = house_title_case(title_el.text) if (title_el and title_el.text) else ""
            filled = fill_title(slide, house_layout, assertion) if assertion else False
            drop_unused_placeholders(slide, {house_layout.title_idx} if filled else set())
            inherited_chrome = True
            receipt.setdefault("template", {
                "house_template": str(house_template),
                "layout": house_layout.name,
                "inherits": ["theme", "slide_master", "band", "footer", "logo", "page_number"],
            })
        else:
            slide = presentation.slides.add_slide(blank)
        # House chrome parity with the HTML renderer (render-oracle finding
        # 2026-08-07: PPTX shipped without the band — cross-target drift):
        # banded recipes get the petrol band with the white title inside it,
        # plus the footer rule; hero recipes stay banner-free.
        # Corpus correction (render-oracle, 2026-08-07): EVERY slide carries the
        # band in the house style — hero recipes put the deck KICKER in the
        # band and keep the assertion as the hero body (reqml-12 pattern).
        hero = bool(slide_doc.intent) and slide_doc.intent.recipe in {"cover-brand", "statement-thesis"}
        banded_by_template = inherited_chrome  # band/footer/page number come from the layout
        is_cover_slide = bool(slide_doc.intent) and slide_doc.intent.recipe == "cover-brand"
        banded = slide_doc.intent is not None and not banded_by_template
        skip_title = inherited_chrome  # the layout's placeholder already carries it
        if banded:
            band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_W_IN), Inches(SLIDE_H_IN * 0.10))
            band.fill.solid()
            band.fill.fore_color.rgb = _hex(band_cfg.get("fill", palette["primary"]))
            band.line.fill.background()
            band.name = "chrome:band"
            # The author's OWN band (Graham_Pittsburg layout, slideLayout22):
            # a solid petrol bar with the house turbine strip stretched over it
            # at alphaModFix amt=10000 (10% opacity). Emitted as a real picture
            # so the band is the same object the author edits, not an imitation.
            texture = (Path(__file__).resolve().parents[3] / "best-practices-slide-design"
                       / "assets" / "house-band-texture.png")
            if texture.is_file():
                pic = slide.shapes.add_picture(
                    str(texture), Inches(-1.2), Inches(-0.08),
                    Inches(SLIDE_W_IN + 2.4), Inches(SLIDE_H_IN * 0.10 + 0.08),
                )
                pic.name = "chrome:band-texture"
                blip = pic._element.find(f".//{{{_A}}}blip")
                if blip is not None:
                    alpha = blip.makeelement(f"{{{_A}}}alphaModFix", {"amt": "10000"})
                    blip.append(alpha)
            title_el = next((e for e in slide_doc.elements if e.role == "title"), None)
            is_cover = bool(slide_doc.intent) and slide_doc.intent.recipe == "cover-brand"
            tagline = document.deck.title.split("—")[-1].strip().upper() if "—" in document.deck.title else ""
            band_text = ((tagline if is_cover else document.deck.title.split("—")[0].strip().upper()) if hero else house_title_case(title_el.text if title_el else ""))
            title_box = slide.shapes.add_textbox(Inches(0.33), Inches(0.07), Inches(SLIDE_W_IN - 1.6), Inches(SLIDE_H_IN * 0.10 - 0.1))
            title_box.name = "chrome:band-title" if hero else (f"el:{title_el.id}" if title_el else "chrome:band-title")
            # wrap on + no autofit: wrap="none"+spAutoFit makes LibreOffice
            # re-fit the box symmetrically, which centers the title and
            # defeats algn="l" (webgpt re-review finding, 2026-08-07).
            title_box.text_frame.word_wrap = True
            title_box.text_frame.auto_size = None
            run = title_box.text_frame.paragraphs[0].add_run()
            run.text = band_text or ""
            run.font.size = Pt(24 if len(band_text or "") <= 58 else 20)
            run.font.bold = True
            run.font.color.rgb = _hex(band_cfg.get("title_color", "#FFFFFF"))
            from pptx.enum.text import PP_ALIGN as _ALIGN

            title_box.text_frame.paragraphs[0].alignment = _ALIGN.LEFT
            skip_title = not hero
        if inherited_chrome:
            ordered = sorted(enumerate(slide_doc.elements), key=lambda pair: (pair[1].z, pair[0]))
            for _, element in ordered:
                if element.role == "footer" or (
                        element.role == "title" and element.bbox.y < 0.15):
                    continue
                if element.role == "title":
                    # statement archetype: a mid-canvas title is DISPLAY text
                    # (huge centered teal), in addition to the band title.
                    display = slide.shapes.add_textbox(
                        Inches(element.bbox.x * SLIDE_W_IN), Inches(element.bbox.y * SLIDE_H_IN),
                        Inches(element.bbox.w * SLIDE_W_IN), Inches(element.bbox.h * SLIDE_H_IN))
                    display.name = f"el:{element.id}"
                    display.text_frame.word_wrap = True
                    d_run = display.text_frame.paragraphs[0].add_run()
                    d_run.text = element.text or ""
                    style = element.style
                    d_run.font.size = Pt(style.size_pt if style and style.size_pt else 44)
                    d_run.font.bold = bool(style.bold) if style else True
                    d_run.font.color.rgb = _hex((style.color if style and style.color else None) or "#065E7C")
                    from pptx.enum.text import PP_ALIGN as _PAL
                    display.text_frame.paragraphs[0].alignment = _PAL.CENTER
                    continue
                _emit_element(slide.shapes, element, root, palette=palette, scale=scale,
                              assets=assets, asset_base=asset_base, receipt=receipt)
            if brandmark:
                from .brandmark import emit_brandmark

                emit_brandmark(slide.shapes, left_in=0.33, top_in=SLIDE_H_IN - 0.62)
            receipt["slides"].append({"id": slide_doc.id, "elements": len(slide_doc.elements)})
            continue
        rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(SLIDE_H_IN - 0.045), Inches(SLIDE_W_IN), Inches(0.045))
        rule.fill.solid()
        rule.fill.fore_color.rgb = _hex(palette["primary"])
        rule.line.fill.background()
        rule.name = "chrome:footer-rule"
        # Footer identity strip (#1314). Every corpus slide carries one; blind
        # judges named its absence as the top authorship tell. Only ledger-
        # licensed facts print — the wordmark is self-evident from the deck
        # title, while sponsor marks and release markings stay absent until the
        # author supplies them (identity.IdentityFact refuses derivation).
        from .identity import ledger_from_document, strip_texts

        identity = strip_texts(ledger_from_document(document.deck.title))
        # the cover IS the wordmark; repeating it in the footer reads as a bug
        if identity.get("wordmark") and not is_cover_slide:
            mark_box = slide.shapes.add_textbox(Inches(0.33), Inches(SLIDE_H_IN - 0.44),
                                                Inches(4.0), Inches(0.32))
            mark_box.name = "chrome:identity-wordmark"
            mark_run = mark_box.text_frame.paragraphs[0].add_run()
            mark_run.text = identity["wordmark"]
            mark_run.font.size = Pt(11)
            mark_run.font.bold = True
            mark_run.font.color.rgb = _hex(palette["primary"])
        page_box = slide.shapes.add_textbox(Inches(SLIDE_W_IN - 0.7), Inches(SLIDE_H_IN - 0.42), Inches(0.5), Inches(0.3))
        page_box.name = "chrome:page-number"
        page_run = page_box.text_frame.paragraphs[0].add_run()
        page_run.text = str(slide_doc.order)
        page_run.font.size = Pt(11)
        page_run.font.color.rgb = _hex("#8a8a8a")
        ordered = sorted(enumerate(slide_doc.elements), key=lambda pair: (pair[1].z, pair[0]))
        for _, element in ordered:
            if skip_title and element.role == "title" and element.bbox.y < 0.15:
                continue
            _emit_element(slide.shapes, element, root, palette=palette, scale=scale, assets=assets, asset_base=asset_base, receipt=receipt)
        receipt["slides"].append({"id": slide_doc.id, "elements": sum(1 for _ in __import__("pitchdeck.document", fromlist=["iter_tree"]).iter_tree(slide_doc.elements))})
    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(output_path))
    if disclaimer:
        scrubbed = _scrub_package_markers(
            output_path,
            stale_owner_markers or ("CSINC", "CS Communication", "SpartaAI", "Sparta AI"),
            "grahama.co",
        )
        if scrubbed:
            receipt["package_marker_scrubbed"] = scrubbed
    receipt["output"] = str(output_path.resolve())
    (output_path.with_suffix(".receipt.json")).write_text(json.dumps(receipt, indent=1), encoding="utf-8")
    return receipt
