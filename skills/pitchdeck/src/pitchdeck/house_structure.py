"""Archetype-conditioned structural conformance on the DELIVERED pptx (#1381).

The pixel channels are spatially blind — a horizontally mirrored deck passed
them 14/15 (reports/house-gate-adversary-benchmark-2026-08-12.md). This module
measures what they cannot: WHERE each semantic role sits and HOW its type is
set, read from the delivered file's own shapes (el:<id> names resolve roles
through the canonical document), judged against per-archetype contracts derived
from the DESIGN_SLIDES catalog and the corpus records.

Inputs: the delivered .pptx and its deck.document.json. Outputs: typed
findings, one per violated contract. Failure modes: a slide whose declared
archetype is unknown FAILS (never silently skipped); unreadable geometry is a
finding, not a pass.
"""

from __future__ import annotations

import json
from pathlib import Path
from typing import Literal

from pydantic import BaseModel

EMU_PER_INCH = 914400.0

# recipe/notes -> archetype (DESIGN_SLIDES catalog)
RECIPE_ARCHETYPE = {
    "cover-brand": "cover",
    "statement-thesis": "statement",
    "assertion-chevrons-diagram": "assertion+art",
    "assertion-chevrons-scene": "assertion+art",
    "one-big-diagram": "assertion+art",
    "proof-screenshot-callout": "mixed",
    "roadmap-lanes": "assertion+art",
    "roadmap-gates": "assertion+art",
}
NOTE_ARCHETYPE = {
    "toc archetype": "toc",
    "section-divider interstitial (DESIGN_SLIDES archetype)": "section-divider",
    "close interstitial (DESIGN_SLIDES archetype)": "close",
    "bullets archetype (claim coverage)": "bullets",
    "mixed Q&A/proof archetype": "mixed",
    "proof archetype (one surface per page)": "mixed",
}

HOUSE_SIZE_PT = {  # per-role point-size ranges measured from the corpus
    "chevrons": (12.0, 22.0),
    "callout": (12.0, 22.0),
    "caption": (9.0, 14.0),
    "footer": (9.0, 14.0),
    "message": (18.0, 64.0),
    "title": (18.0, 64.0),
}


class StructureFinding(BaseModel):
    code: Literal[
        "UNKNOWN_ARCHETYPE",
        "ROLE_REGION_VIOLATION",
        "TYPOGRAPHY_VIOLATION",
        "VISUAL_SUBSTANCE_VIOLATION",
        "UNREADABLE",
    ]
    slide: int
    archetype: str
    detail: str


def _archetype_of(slide_doc: dict) -> str | None:
    intent = slide_doc.get("intent") or {}
    if intent.get("recipe") in RECIPE_ARCHETYPE:
        return RECIPE_ARCHETYPE[intent["recipe"]]
    notes = slide_doc.get("notes") or ""
    for key, arch in NOTE_ARCHETYPE.items():
        if notes.startswith(key.split(" (")[0]):
            return arch
    return None


def _role_map(slide_doc: dict) -> dict[str, str]:
    return {el["id"]: el.get("role", "") for el in slide_doc.get("elements", [])}


def check_structure(pptx_path: Path, document_path: Path) -> list[StructureFinding]:
    from pptx import Presentation

    document = json.loads(Path(document_path).read_text())
    doc_slides = [s for s in document["slides"] if not s.get("hidden")]
    pres = Presentation(str(pptx_path))
    width_in = pres.slide_width / EMU_PER_INCH
    height_in = pres.slide_height / EMU_PER_INCH
    findings: list[StructureFinding] = []

    for index, (slide, slide_doc) in enumerate(zip(pres.slides, doc_slides), start=1):
        archetype = _archetype_of(slide_doc)
        if archetype is None:
            findings.append(StructureFinding(code="UNKNOWN_ARCHETYPE", slide=index, archetype="?",
                detail=f"slide '{slide_doc.get('id')}' declares no recognizable archetype"))
            continue
        roles = _role_map(slide_doc)
        visual_area = 0.0
        for shape in slide.shapes:
            name = shape.name or ""
            if not name.startswith("el:"):
                continue
            element_id = name[3:]
            role = roles.get(element_id, "")
            try:
                x = shape.left / EMU_PER_INCH / width_in
                y = shape.top / EMU_PER_INCH / height_in
                w = shape.width / EMU_PER_INCH / width_in
                h = shape.height / EMU_PER_INCH / height_in
            except (TypeError, AttributeError, ZeroDivisionError):
                findings.append(StructureFinding(code="UNREADABLE", slide=index,
                    archetype=archetype, detail=f"el:{element_id} has unreadable geometry"))
                continue

            # --- role-region contracts (what the mirror mutant breaks) ---
            if role in {"chevrons", "callout"} and x > 0.34:
                findings.append(StructureFinding(code="ROLE_REGION_VIOLATION", slide=index,
                    archetype=archetype,
                    detail=f"{role} '{element_id}' anchored at x={x:.2f} — house {role} start left (x<=0.34)"))
            if role == "badge" and element_id == "house-mark" and not (x < 0.2 and y > 0.8):
                findings.append(StructureFinding(code="ROLE_REGION_VIOLATION", slide=index,
                    archetype=archetype,
                    detail=f"identity mark at ({x:.2f},{y:.2f}) — house mark is bottom-left"))
            if role == "badge" and element_id in {"divider-mark", "thesis-mark", "cover-mark"} and x < 0.5:
                findings.append(StructureFinding(code="ROLE_REGION_VIOLATION", slide=index,
                    archetype=archetype,
                    detail=f"product mark '{element_id}' at x={x:.2f} — house product marks sit right"))
            if archetype == "section-divider" and role == "message":
                center = x + w / 2
                if abs(center - 0.5) > 0.15:
                    findings.append(StructureFinding(code="ROLE_REGION_VIOLATION", slide=index,
                        archetype=archetype,
                        detail=f"divider heading centered at {center:.2f} — house dividers center (0.35-0.65)"))
            if role == "visual" and y < 0.12:
                findings.append(StructureFinding(code="ROLE_REGION_VIOLATION", slide=index,
                    archetype=archetype,
                    detail=f"visual '{element_id}' at y={y:.2f} overlaps the band region"))
            if role == "visual":
                visual_area += max(w, 0) * max(h, 0)

            # --- typography contracts (what the ransom-note mutant breaks) ---
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame and role in HOUSE_SIZE_PT:
                lo, hi = HOUSE_SIZE_PT[role]
                sizes = [run.font.size.pt for para in shape.text_frame.paragraphs
                         for run in para.runs if run.font.size is not None]
                for size in sizes:
                    if not (lo <= size <= hi):
                        findings.append(StructureFinding(code="TYPOGRAPHY_VIOLATION", slide=index,
                            archetype=archetype,
                            detail=f"{role} '{element_id}' run at {size:.0f}pt outside house range {lo}-{hi}pt"))
                        break
                if len(sizes) >= 3:
                    spread = max(sizes) - min(sizes)
                    if spread > 12:
                        findings.append(StructureFinding(code="TYPOGRAPHY_VIOLATION", slide=index,
                            archetype=archetype,
                            detail=f"{role} '{element_id}' mixes sizes across {spread:.0f}pt — no house hierarchy does"))

        # --- per-archetype visual substance (what two-tiny-visuals breaks) ---
        min_area = {"assertion+art": 0.18, "mixed": 0.12, "bullets": 0.08}.get(archetype)
        if min_area is not None and visual_area < min_area:
            findings.append(StructureFinding(code="VISUAL_SUBSTANCE_VIOLATION", slide=index,
                archetype=archetype,
                detail=(f"visual area {visual_area:.3f} < {min_area} required for {archetype} — "
                        "object count without area is not substance")))
    return findings
