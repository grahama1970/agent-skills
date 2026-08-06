from __future__ import annotations

import importlib.util
from datetime import UTC, datetime
import re
import shutil
from pathlib import Path

from pptx import Presentation

from .io import expand_path, sha256_file
from .models import (
    AssetKind,
    BindingKind,
    AssetManifest,
    AssetStatus,
    ClaimKind,
    ClaimLedger,
    ClaimStatus,
    DeckManifest,
    DeckSourcePolicy,
    Readiness,
    SourceManifest,
    ValidationIssue,
    ValidationReport,
    Visibility,
    VisualType,
)

_SUPPORTED_IMAGES = {".png", ".jpg", ".jpeg", ".webp", ".svg", ".bmp", ".gif", ".tif", ".tiff"}
_SUPPORTED_VIDEO = {".mp4", ".webm"}
_NEGATION_RE = re.compile(r"\b(not|no|never|without|unverified|cannot|can't|does not|not yet)\b", re.IGNORECASE)


def _qualified_occurrence(text: str, phrase: str) -> bool:
    lowered = text.lower()
    needle = phrase.lower()
    start = 0
    while True:
        index = lowered.find(needle, start)
        if index < 0:
            return True
        prefix = text[max(0, index - 48) : index]
        suffix = text[index + len(phrase) : index + len(phrase) + 36]
        if not (_NEGATION_RE.search(prefix) or _NEGATION_RE.search(suffix)):
            return False
        start = index + len(needle)


_DIGIT_RE = re.compile(r"\d+(?:[.,]\d+)*")  # known grammar: numeric tokens in deck strings


def _norm(text: str) -> str:
    return re.sub(r"\s+", " ", text.strip().lower())


def _stem(token: str) -> str:
    for suffix in ("ing", "ed", "es", "s"):
        if token.endswith(suffix) and len(token) > len(suffix) + 2:
            return token[: -len(suffix)]
    return token


def _check_rendering(slide_id, binding, resolved, claim, require_approved_claims) -> list:
    """Mechanical transform verification + NUMERIC_UNBOUND (roundtable #1226).

    verbatim/truncation/inflection are checkable against the claim's evidence
    spans and auto-pass; aggregation/generalization lean on claim approval
    (already publish-enforced). Every digit in a rendering must appear in a
    span or in the claim's recorded formula.
    """
    issues: list[ValidationIssue] = []
    span_norms = [_norm(s.text) for s in claim.evidence_spans]
    rendering = _norm(resolved)

    tc = binding.transform_class
    if tc == "verbatim" and not any(rendering in s or s in rendering for s in span_norms):
        issues.append(
            ValidationIssue(
                severity="error",
                code="TRANSFORM_MISMATCH",
                message=f"Binding '{binding.path}' claims transform_class=verbatim but matches no evidence span of '{claim.id}'.",
                slide_id=slide_id,
                claim_id=claim.id,
            )
        )
    elif tc == "truncation":
        clipped = rendering.rstrip(" .…")
        if not any(clipped in s for s in span_norms):
            issues.append(
                ValidationIssue(
                    severity="error",
                    code="TRANSFORM_MISMATCH",
                    message=f"Binding '{binding.path}' claims transform_class=truncation but is not a contiguous excerpt of any span of '{claim.id}'.",
                    slide_id=slide_id,
                    claim_id=claim.id,
                )
            )
    elif tc == "inflection":
        span_stems = {_stem(t) for s in span_norms for t in re.findall(r"[a-z0-9]+", s)}
        rendering_stems = {_stem(t) for t in re.findall(r"[a-z0-9]+", rendering)}
        if not rendering_stems <= span_stems:
            issues.append(
                ValidationIssue(
                    severity="error",
                    code="TRANSFORM_MISMATCH",
                    message=f"Binding '{binding.path}' claims transform_class=inflection but introduces tokens absent from every span of '{claim.id}'.",
                    slide_id=slide_id,
                    claim_id=claim.id,
                )
            )

    # NUMERIC_UNBOUND: every digit sequence must trace to a span or a formula.
    span_and_formula = " ".join([*span_norms, _norm(claim.formula or "")])
    for number in _DIGIT_RE.findall(resolved):
        if number not in span_and_formula:
            issues.append(
                ValidationIssue(
                    severity="error" if require_approved_claims else "warning",
                    code="NUMERIC_UNBOUND",
                    message=(
                        f"Binding '{binding.path}' renders number '{number}' that appears in no evidence "
                        f"span or recorded formula of claim '{claim.id}'."
                    ),
                    slide_id=slide_id,
                    claim_id=claim.id,
                )
            )
    return issues


def validate_bundle(
    deck: DeckManifest,
    ledger: ClaimLedger,
    sources: SourceManifest,
    assets: AssetManifest,
    *,
    source_manifest_dir: Path,
    asset_manifest_dir: Path,
    require_approved_claims: bool = False,
) -> ValidationReport:
    issues: list[ValidationIssue] = []
    source_map = {source.id: source for source in sources.sources}
    claim_map = {claim.id: claim for claim in ledger.claims}
    asset_map = {asset.id: asset for asset in assets.assets}
    allowed_public = set(sources.policy.public_deck_source_ids)
    bound_claim_ids: set[str] = set()

    for source in sources.sources:
        path = expand_path(source.path, base_dir=source_manifest_dir)
        if source.required and not path.exists():
            issues.append(
                ValidationIssue(
                    severity="error",
                    code="SOURCE_REQUIRED_MISSING",
                    message=f"Required source '{source.id}' is missing at {path}",
                    source_id=source.id,
                )
            )
        elif not path.exists():
            issues.append(
                ValidationIssue(
                    severity="warning",
                    code="SOURCE_OPTIONAL_MISSING",
                    message=f"Optional source '{source.id}' is missing at {path}",
                    source_id=source.id,
                )
            )
        if path.exists() and source.content_sha and source.content_sha.startswith("sha256:"):
            actual = sha256_file(path)
            expected = source.content_sha.split(":", 1)[1]
            if actual != expected:
                issues.append(
                    ValidationIssue(
                        severity="error",
                        code="SOURCE_SHA_MISMATCH",
                        message=f"Source '{source.id}' sha256 mismatch: expected {expected}, got {actual}",
                        source_id=source.id,
                    )
                )

    # Living-deck drift gate (#1229): a stale snapshot blocks publish.
    if require_approved_claims:
        from .drift import SOURCE_STATE_FILE, check_drift

        if (asset_manifest_dir / SOURCE_STATE_FILE).exists():
            drift = check_drift(asset_manifest_dir, deck, ledger, sources, source_manifest_dir)
            for source_id in [*drift["changed"], *drift["missing"]]:
                issues.append(
                    ValidationIssue(
                        severity="error",
                        code="SOURCE_DRIFT",
                        message=(
                            f"Source '{source_id}' changed or vanished since the recorded snapshot; "
                            f"affected claims: {', '.join(drift['affected_claims']) or 'none mapped'}. "
                            "Re-plan, re-review, and refresh the snapshot before publishing."
                        ),
                        source_id=source_id,
                    )
                )

    for slide in deck.slides:
        element_texts = [e.text for e in slide.elements if e.text]
        visual_texts = [*slide.visual.items, *slide.visual.callouts]
        if slide.visual.caption:
            visual_texts.append(slide.visual.caption)
        if slide.visual.source:
            # Mermaid/KaTeX source renders as visible text — scan it too.
            visual_texts.append(slide.visual.source)
        if slide.visual.type.value in {"mermaid", "math"} and not slide.visual.asset_id:
            issues.append(
                ValidationIssue(
                    code="DIAGRAM_NO_SNAPSHOT",
                    # Review condition (both seats): a client deck with a
                    # placeholder is not an alternate rendering — error at publish.
                    severity="error" if require_approved_claims else "warning",
                    slide_id=slide.id,
                    message=(
                        f"slide '{slide.id}' {slide.visual.type.value} visual has no snapshot asset; "
                        "publish requires one (PPTX would ship a placeholder panel)"
                    ),
                )
            )
        if slide.visual.type.value == "mermaid" and "%%{" in (slide.visual.source or ""):
            issues.append(
                ValidationIssue(
                    code="MERMAID_DIRECTIVE",
                    severity="error",
                    slide_id=slide.id,
                    message=(
                        f"slide '{slide.id}' mermaid source contains a %%{{...}}%% init directive; "
                        "directives can override renderer security config and are rejected"
                    ),
                )
            )
        if slide.footer:
            visual_texts.append(slide.footer)
        visible_text = "\n".join(
            [slide.title, slide.message, *slide.body, *element_texts, *visual_texts]
        )
        all_text = "\n".join([visible_text, slide.notes])

        # --- ContentIR bindings (roundtable session 2) -----------------------
        def _resolve_path(path: str) -> str | None:
            base, _, index = path.partition(":")
            if base == "title":
                return slide.title
            if base == "message":
                return slide.message
            if base == "footer":
                return slide.footer or ""
            if base == "body" and index.isdigit() and int(index) < len(slide.body):
                return slide.body[int(index)]
            if base == "element":
                for element in slide.elements:
                    if element.id == index:
                        return element.text or ""
                return None
            if base == "visual.items" and index.isdigit() and int(index) < len(slide.visual.items):
                return slide.visual.items[int(index)]
            if base == "visual.caption":
                return slide.visual.caption or ""
            return None

        bound_paths: set[str] = set()
        structural_qualifiers: dict[str, str] = {}
        for binding in slide.bindings:
            resolved = _resolve_path(binding.path)
            if resolved is None:
                issues.append(
                    ValidationIssue(
                        severity="error",
                        code="BINDING_UNKNOWN_PATH",
                        message=f"Binding path '{binding.path}' does not resolve on this slide.",
                        slide_id=slide.id,
                    )
                )
                continue
            bound_paths.add(binding.path)
            if binding.claim_id:
                bound_claim = claim_map.get(binding.claim_id)
                if bound_claim is None:
                    issues.append(
                        ValidationIssue(
                            severity="error",
                            code="BINDING_UNKNOWN_CLAIM",
                            message=f"Binding '{binding.path}' references unknown claim '{binding.claim_id}'.",
                            slide_id=slide.id,
                            claim_id=binding.claim_id,
                        )
                    )
                    continue
                if binding.kind == BindingKind.CLAIM_QUOTE and bound_claim.text.strip().lower() not in resolved.strip().lower():
                    issues.append(
                        ValidationIssue(
                            severity="error",
                            code="BINDING_QUOTE_MISMATCH",
                            message=(
                                f"claim_quote binding '{binding.path}' does not contain the ledger text "
                                f"of claim '{binding.claim_id}'."
                            ),
                            slide_id=slide.id,
                            claim_id=binding.claim_id,
                        )
                    )
                if binding.kind == BindingKind.QUALIFIER:
                    structural_qualifiers[binding.claim_id] = resolved
                # --- Span-first rendering gates (roundtable 2026-08-06, #1226) ---
                if binding.kind in {BindingKind.CLAIM_QUOTE, BindingKind.CLAIM_PARAPHRASE}:
                    if not bound_claim.evidence_spans:
                        issues.append(
                            ValidationIssue(
                                severity="error" if require_approved_claims else "warning",
                                code="RENDERING_UNBOUND",
                                message=(
                                    f"Binding '{binding.path}' renders claim '{binding.claim_id}' which has "
                                    "no evidence spans; publish requires every rendering to resolve to a span."
                                ),
                                slide_id=slide.id,
                                claim_id=binding.claim_id,
                            )
                        )
                    else:
                        issues.extend(
                            _check_rendering(
                                slide.id, binding, resolved, bound_claim, require_approved_claims
                            )
                        )

        # Structural qualifier authority: a high-risk claim's qualifier must be
        # STRUCTURALLY bound to visible text containing the required qualifier.
        # The negation-window text scan remains as defense-in-depth lint only.
        for claim_id in slide.claim_ids:
            claim = claim_map.get(claim_id)
            if claim is None or not claim.required_qualifier:
                continue
            bound_text = structural_qualifiers.get(claim_id)
            if bound_text is None:
                issues.append(
                    ValidationIssue(
                        severity="error" if require_approved_claims else "warning",
                        code="QUALIFIER_NOT_STRUCTURAL",
                        message=(
                            f"High-risk claim '{claim_id}' has no structural qualifier binding on this "
                            "slide; publish requires kind=qualifier bound to visible text."
                        ),
                        slide_id=slide.id,
                        claim_id=claim_id,
                    )
                )
            elif claim.required_qualifier.strip().lower() not in bound_text.strip().lower():
                issues.append(
                    ValidationIssue(
                        severity="error",
                        code="QUALIFIER_BINDING_TEXT_MISMATCH",
                        message=(
                            f"Qualifier binding for claim '{claim_id}' does not contain its required "
                            "qualifier text."
                        ),
                        slide_id=slide.id,
                        claim_id=claim_id,
                    )
                )

        # Publish gate: no unclassified substantive string may reach a published
        # artifact (webgpt exit criterion, applied at the gate).
        substantive: list[tuple[str, str]] = [("title", slide.title), ("message", slide.message)]
        substantive += [(f"body:{i}", line) for i, line in enumerate(slide.body)]
        substantive += [(f"element:{e.id}", e.text or "") for e in slide.elements if e.type == "text"]
        substantive += [(f"visual.items:{i}", item) for i, item in enumerate(slide.visual.items)]
        unbound = [path for path, text in substantive if text.strip() and path not in bound_paths]
        if unbound:
            issues.append(
                ValidationIssue(
                    severity="error" if require_approved_claims else "warning",
                    code="UNBOUND_TEXT",
                    message=(
                        f"{len(unbound)} substantive string(s) have no content binding "
                        f"({', '.join(unbound[:4])}{'…' if len(unbound) > 4 else ''}); publish requires full coverage."
                    ),
                    slide_id=slide.id,
                )
            )
        # ---------------------------------------------------------------------

        if not slide.source_refs:
            issues.append(
                ValidationIssue(
                    severity="error",
                    code="SLIDE_SOURCELESS",
                    message="Every slide must have at least one source reference.",
                    slide_id=slide.id,
                )
            )

        for ref in slide.source_refs:
            source = source_map.get(ref.source_id)
            if source is None:
                issues.append(
                    ValidationIssue(
                        severity="error",
                        code="SLIDE_UNKNOWN_SOURCE",
                        message=f"Slide references unknown source '{ref.source_id}'.",
                        slide_id=slide.id,
                        source_id=ref.source_id,
                    )
                )
                continue
            if deck.deck.source_policy == DeckSourcePolicy.PUBLIC_ONLY:
                if source.visibility != Visibility.PUBLIC:
                    issues.append(
                        ValidationIssue(
                            severity="error",
                            code="PUBLIC_DECK_PRIVATE_SOURCE",
                            message=f"Public slide references private source '{source.id}'.",
                            slide_id=slide.id,
                            source_id=source.id,
                        )
                    )
                if allowed_public and source.id not in allowed_public:
                    issues.append(
                        ValidationIssue(
                            severity="error",
                            code="PUBLIC_DECK_SOURCE_NOT_ALLOWLISTED",
                            message=f"Source '{source.id}' is not allowlisted for the public deck.",
                            slide_id=slide.id,
                            source_id=source.id,
                        )
                    )

        for claim_id in slide.claim_ids:
            bound_claim_ids.add(claim_id)
            claim = claim_map.get(claim_id)
            if claim is None:
                issues.append(
                    ValidationIssue(
                        severity="error",
                        code="SLIDE_UNKNOWN_CLAIM",
                        message=f"Slide references unknown claim '{claim_id}'.",
                        slide_id=slide.id,
                        claim_id=claim_id,
                    )
                )
                continue
            if claim.status == ClaimStatus.REJECTED:
                issues.append(
                    ValidationIssue(
                        severity="error",
                        code="SLIDE_REJECTED_CLAIM",
                        message=f"Slide uses rejected claim '{claim_id}'.",
                        slide_id=slide.id,
                        claim_id=claim_id,
                    )
                )
            if require_approved_claims and claim.status != ClaimStatus.APPROVED:
                issues.append(
                    ValidationIssue(
                        severity="error",
                        code="SLIDE_UNAPPROVED_CLAIM",
                        message=f"Slide uses non-approved claim '{claim_id}' with status={claim.status.value}.",
                        slide_id=slide.id,
                        claim_id=claim_id,
                    )
                )
            elif claim.status == ClaimStatus.CANDIDATE:
                issues.append(
                    ValidationIssue(
                        severity="warning",
                        code="SLIDE_CANDIDATE_CLAIM",
                        message=f"Slide uses candidate claim '{claim_id}'; human review remains open.",
                        slide_id=slide.id,
                        claim_id=claim_id,
                    )
                )
            # Claim-level source authority (roundtable item 5): resolve the
            # BOUND CLAIM's own sources, not just the slide's.
            for ref in claim.source_refs:
                claim_source = source_map.get(ref.source_id)
                if claim_source is None:
                    issues.append(
                        ValidationIssue(
                            severity="error",
                            code="CLAIM_UNKNOWN_SOURCE",
                            message=f"Claim '{claim_id}' references unknown source '{ref.source_id}'.",
                            slide_id=slide.id,
                            claim_id=claim_id,
                            source_id=ref.source_id,
                        )
                    )
                elif (
                    deck.deck.source_policy == DeckSourcePolicy.PUBLIC_ONLY
                    and claim_source.visibility != Visibility.PUBLIC
                ):
                    issues.append(
                        ValidationIssue(
                            severity="error",
                            code="CLAIM_PRIVATE_SOURCE",
                            message=f"Public deck binds claim '{claim_id}' whose source '{ref.source_id}' is private.",
                            slide_id=slide.id,
                            claim_id=claim_id,
                            source_id=ref.source_id,
                        )
                    )
            # Approval provenance + staleness (publish-authoritative).
            if claim.status == ClaimStatus.APPROVED:
                if claim.approval is None:
                    issues.append(
                        ValidationIssue(
                            severity="error" if require_approved_claims else "warning",
                            code="APPROVAL_PROVENANCE_MISSING",
                            message=f"Approved claim '{claim_id}' has no approval provenance (approved_by/approved_at).",
                            slide_id=slide.id,
                            claim_id=claim_id,
                        )
                    )
                else:
                    if claim.approval.fixture and require_approved_claims:
                        issues.append(
                            ValidationIssue(
                                severity="error",
                                code="APPROVAL_FIXTURE",
                                message=(
                                    f"Claim '{claim_id}' carries a FIXTURE approval stamp "
                                    f"({claim.approval.approved_by}); publish requires real human review."
                                ),
                                slide_id=slide.id,
                                claim_id=claim_id,
                            )
                        )
                if claim.approval is not None and claim.approval.expires_at:
                    try:
                        expiry = datetime.fromisoformat(claim.approval.expires_at.replace("Z", "+00:00"))
                        if expiry.tzinfo is None:
                            expiry = expiry.replace(tzinfo=UTC)
                        if expiry < datetime.now(UTC):
                            issues.append(
                                ValidationIssue(
                                    severity="error",
                                    code="APPROVAL_EXPIRED",
                                    message=f"Approval for claim '{claim_id}' expired {claim.approval.expires_at}; re-review required.",
                                    slide_id=slide.id,
                                    claim_id=claim_id,
                                )
                            )
                    except ValueError:
                        issues.append(
                            ValidationIssue(
                                severity="error",
                                code="APPROVAL_EXPIRY_INVALID",
                                message=f"Claim '{claim_id}' approval.expires_at is not ISO formatted.",
                                slide_id=slide.id,
                                claim_id=claim_id,
                            )
                        )
            if deck.deck.source_policy == DeckSourcePolicy.PUBLIC_ONLY and claim.visibility != Visibility.PUBLIC:
                issues.append(
                    ValidationIssue(
                        severity="error",
                        code="PUBLIC_DECK_PRIVATE_CLAIM",
                        message=f"Public slide references private claim '{claim_id}'.",
                        slide_id=slide.id,
                        claim_id=claim_id,
                    )
                )
            if claim.required_qualifier and claim.required_qualifier.lower() not in all_text.lower():
                issues.append(
                    ValidationIssue(
                        severity="error",
                        code="CLAIM_QUALIFIER_MISSING",
                        message=(
                            f"High-risk claim '{claim_id}' requires qualifier: "
                            f"{claim.required_qualifier}"
                        ),
                        slide_id=slide.id,
                        claim_id=claim_id,
                    )
                )

        required_non_claims = set(slide.claim_guard.requires_non_claim_ids)
        if not required_non_claims.issubset(set(slide.claim_ids)):
            missing = sorted(required_non_claims - set(slide.claim_ids))
            issues.append(
                ValidationIssue(
                    severity="error",
                    code="SLIDE_NON_CLAIM_NOT_BOUND",
                    message=f"Slide requires non-claims not present in claim_ids: {missing}",
                    slide_id=slide.id,
                )
            )

        forbidden = [
            *sources.policy.forbidden_unqualified_claims,
            *slide.claim_guard.forbidden_unqualified,
        ]
        for phrase in dict.fromkeys(forbidden):
            if phrase.lower() in visible_text.lower() and not _qualified_occurrence(visible_text, phrase):
                issues.append(
                    ValidationIssue(
                        severity="error",
                        code="FORBIDDEN_UNQUALIFIED_CLAIM",
                        message=f"Visible slide text contains forbidden unqualified phrase: '{phrase}'.",
                        slide_id=slide.id,
                    )
                )

        for element in slide.elements:
            if element.type == "asset" and element.asset_id not in asset_map:
                issues.append(
                    ValidationIssue(
                        severity="error",
                        code="ELEMENT_UNKNOWN_ASSET",
                        message=f"Freeform element '{element.id}' references unknown asset '{element.asset_id}'.",
                        slide_id=slide.id,
                    )
                )
        # True per-renderer body capacities read from pptx_builder slices:
        # statement/split/_render_cards slice [:4], flow [:6], roadmap buckets
        # 3 columns x 3 rows, collaboration [:3], appendix [:6].
        _PPTX_BODY_CAPACITY = {
            "statement": 4, "split": 4, "flow": 6, "roadmap": 9,
            "three_cards": 4, "proof_cards": 4, "collaboration": 3, "appendix": 6,
        }
        capacity = _PPTX_BODY_CAPACITY.get(slide.layout.value)
        if capacity is not None and len(slide.body) > capacity:
            issues.append(
                ValidationIssue(
                    severity="error",
                    code="TARGET_CONTENT_TRUNCATED",
                    message=(
                        f"Slide has {len(slide.body)} body items but the {slide.layout.value} "
                        f"PPTX renderer shows only {capacity}; content (possibly a qualifier) "
                        "would be silently dropped from export. Split the slide or shorten the body."
                    ),
                    slide_id=slide.id,
                )
            )
        if slide.layout.value == "freeform" and not slide.elements:
            issues.append(
                ValidationIssue(
                    severity="error",
                    code="FREEFORM_NO_ELEMENTS",
                    message="Freeform slides require at least one element.",
                    slide_id=slide.id,
                )
            )
        # Generated-pixel boundary (#1230): assets born from a generation brief
        # are decorative ONLY — never on claim-bearing visual surfaces.
        if slide.visual.asset_id:
            visual_asset = asset_map.get(slide.visual.asset_id)
            claim_bearing_visual = slide.visual.type in {
                VisualType.NATIVE_DIAGRAM,
                VisualType.MERMAID,
                VisualType.MATH,
            } or any(
                b.path.startswith("visual.") and b.kind.value in {"claim_quote", "claim_paraphrase"}
                for b in slide.bindings
            )
            if (
                visual_asset is not None
                and visual_asset.generation_brief
                and claim_bearing_visual
            ):
                issues.append(
                    ValidationIssue(
                        severity="error",
                        code="GENERATED_ASSET_CLAIM_SURFACE",
                        message=(
                            f"Generated asset '{visual_asset.id}' (has generation_brief) is bound to a "
                            "claim-bearing visual; generated pixels are decorative only."
                        ),
                        slide_id=slide.id,
                        asset_id=visual_asset.id,
                    )
                )
        if slide.visual.type in {VisualType.IMAGE, VisualType.SCREENSHOT}:
            asset_id = slide.visual.asset_id or ""
            asset = asset_map.get(asset_id)
            if asset is None:
                issues.append(
                    ValidationIssue(
                        severity="error",
                        code="SLIDE_UNKNOWN_ASSET",
                        message=f"Slide references unknown asset '{asset_id}'.",
                        slide_id=slide.id,
                        asset_id=asset_id,
                    )
                )
            else:
                if deck.deck.source_policy == DeckSourcePolicy.PUBLIC_ONLY and asset.visibility != Visibility.PUBLIC:
                    issues.append(
                        ValidationIssue(
                            severity="error",
                            code="PUBLIC_DECK_PRIVATE_ASSET",
                            message=f"Public slide references private asset '{asset.id}'.",
                            slide_id=slide.id,
                            asset_id=asset.id,
                        )
                    )
                local_path = (
                    expand_path(asset.local_path, base_dir=asset_manifest_dir)
                    if asset.local_path
                    else None
                )
                missing = local_path is None or not local_path.exists()
                if asset.required and missing:
                    issues.append(
                        ValidationIssue(
                            severity="error",
                            code="REQUIRED_ASSET_MISSING",
                            message=f"Required asset '{asset.id}' has no readable local file.",
                            slide_id=slide.id,
                            asset_id=asset.id,
                        )
                    )
                elif missing:
                    issues.append(
                        ValidationIssue(
                            severity="warning",
                            code="OPTIONAL_ASSET_MISSING",
                            message=f"Optional asset '{asset.id}' will render as an explicit missing-asset card.",
                            slide_id=slide.id,
                            asset_id=asset.id,
                        )
                    )
                elif local_path.suffix.lower() not in (
                    _SUPPORTED_VIDEO if asset.kind == AssetKind.VIDEO else _SUPPORTED_IMAGES
                ):
                    severity = "error" if asset.required else "warning"
                    expected = (
                        "video formats " + "/".join(sorted(_SUPPORTED_VIDEO))
                        if asset.kind == AssetKind.VIDEO
                        else "image formats"
                    )
                    issues.append(
                        ValidationIssue(
                            severity=severity,
                            code="ASSET_UNSUPPORTED_FORMAT",
                            message=(
                                f"Asset '{asset.id}' (kind {asset.kind.value}) uses unsupported format "
                                f"{local_path.suffix}; expected {expected}."
                            ),
                            slide_id=slide.id,
                            asset_id=asset.id,
                        )
                    )
                elif local_path.suffix.lower() == ".svg" and (
                    importlib.util.find_spec("cairosvg") is None
                    and shutil.which("rsvg-convert") is None
                ):
                    severity = "error" if asset.required else "warning"
                    issues.append(
                        ValidationIssue(
                            severity=severity,
                            code="SVG_CONVERTER_MISSING",
                            message=(
                                f"Asset '{asset.id}' is SVG but neither optional cairosvg nor "
                                "rsvg-convert is available."
                            ),
                            slide_id=slide.id,
                            asset_id=asset.id,
                        )
                    )
                if asset.status in {AssetStatus.STALE, AssetStatus.REGENERATE}:
                    issues.append(
                        ValidationIssue(
                            severity="warning",
                            code="ASSET_NOT_CURRENT",
                            message=f"Asset '{asset.id}' is marked {asset.status.value}.",
                            slide_id=slide.id,
                            asset_id=asset.id,
                        )
                    )

    mandatory = {claim.id for claim in ledger.claims if claim.kind == ClaimKind.NON_CLAIM}
    missing_mandatory = mandatory - bound_claim_ids
    for claim_id in sorted(missing_mandatory):
        issues.append(
            ValidationIssue(
                severity="error",
                code="MANDATORY_NON_CLAIM_UNBOUND",
                message=f"Mandatory non-claim '{claim_id}' is not bound to any slide.",
                claim_id=claim_id,
            )
        )

    errors = sum(1 for issue in issues if issue.severity == "error")
    warnings = sum(1 for issue in issues if issue.severity == "warning")
    if errors:
        readiness = Readiness.NOT_READY
    elif warnings:
        readiness = Readiness.USABLE_WITH_GAPS
    else:
        readiness = Readiness.READY
    return ValidationReport(readiness=readiness, errors=errors, warnings=warnings, issues=issues)


def validate_pptx(path: Path, expected_slides: int) -> list[ValidationIssue]:
    issues: list[ValidationIssue] = []
    if not path.exists():
        return [
            ValidationIssue(
                severity="error",
                code="PPTX_MISSING",
                message=f"PPTX does not exist: {path}",
            )
        ]
    try:
        presentation = Presentation(path)
    except Exception as exc:
        return [
            ValidationIssue(
                severity="error",
                code="PPTX_REOPEN_FAILED",
                message=f"Generated PPTX cannot be reopened: {exc}",
            )
        ]
    if len(presentation.slides) != expected_slides:
        issues.append(
            ValidationIssue(
                severity="error",
                code="PPTX_SLIDE_COUNT_MISMATCH",
                message=f"Expected {expected_slides} slides, found {len(presentation.slides)}.",
            )
        )
    for index, slide in enumerate(presentation.slides, start=1):
        visible_text = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    visible_text.append(text)
        if not visible_text:
            issues.append(
                ValidationIssue(
                    severity="error",
                    code="PPTX_BLANK_SLIDE",
                    message=f"Slide {index} contains no editable text.",
                    slide_id=str(index),
                )
            )
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if not notes:
            issues.append(
                ValidationIssue(
                    severity="warning",
                    code="PPTX_NOTES_MISSING",
                    message=f"Slide {index} has no speaker/source notes.",
                    slide_id=str(index),
                )
            )
    return issues
