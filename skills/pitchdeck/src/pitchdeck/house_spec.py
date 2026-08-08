"""Measured house layout/typography spec from the author's real decks (#1311).

The house style is not a matter of taste to be guessed at — it is measurable.
``measure_house_spec`` reads the author's PPTX corpus and reports the geometry,
type scale, and color scheme actually used: title zone, content-image zone,
footer zone, per-role point sizes, and accent colors ranked by run frequency.
``HOUSE_SPEC`` is the committed measurement (provenance in ``measured_from``)
that emitters compile against, so a drift between code and corpus is a diff in
this file rather than an invisible style regression.

Inputs: a directory of .pptx decks. Outputs: HouseSpec (deterministic —
sorted/median aggregation only, no wall-clock). Failure modes: an empty corpus
raises rather than yielding a spec with invented defaults; shapes whose
geometry cannot be read are skipped and counted, never guessed.
"""

from __future__ import annotations

import statistics
from collections import Counter
from pathlib import Path
from typing import Literal

from pydantic import Field

from .models import StrictModel

EMU_PER_INCH = 914400.0
_CONTENT_IMAGE_MIN_AREA = 0.06  # fractional area separating content images from inline glyphs


class Zone(StrictModel):
    """Fractional slide rectangle (median across the corpus)."""

    x: float = Field(ge=0.0, le=1.0)
    y: float = Field(ge=0.0, le=1.0)
    w: float = Field(gt=0.0, le=1.0)
    h: float = Field(gt=0.0, le=1.0)
    samples: int = Field(ge=1)


class HouseSpec(StrictModel):
    schema_: Literal["pitchdeck.house_spec.v1"] = Field(default="pitchdeck.house_spec.v1", alias="schema")
    measured_from: list[str]
    slide_size_in: tuple[float, float]
    title_zone: Zone
    content_image_zone: Zone
    inline_glyph_zone: Zone
    footer_zone: Zone
    title_pt: list[float] = Field(description="Most common title point sizes, most frequent first.")
    body_pt: list[float] = Field(description="Most common body point sizes, most frequent first.")
    accent_colors: list[str] = Field(description="Run colors ranked by frequency (hex, no #).")
    unreadable_shapes: int = Field(ge=0)


def _median_zone(boxes: list[tuple[float, float, float, float]]) -> Zone:
    xs, ys, ws, hs = zip(*boxes)
    return Zone(
        x=round(statistics.median(xs), 4),
        y=round(statistics.median(ys), 4),
        w=round(statistics.median(ws), 4),
        h=round(statistics.median(hs), 4),
        samples=len(boxes),
    )


def measure_house_spec(decks_dir: Path) -> HouseSpec:
    """Measure geometry, type scale, and colors across every deck in a directory."""
    from pptx import Presentation

    decks = sorted(decks_dir.glob("*.pptx"))
    if not decks:
        raise ValueError(f"no .pptx decks under {decks_dir} — cannot measure a house spec")

    titles: list[tuple[float, float, float, float]] = []
    content_images: list[tuple[float, float, float, float]] = []
    glyphs: list[tuple[float, float, float, float]] = []
    footers: list[tuple[float, float, float, float]] = []
    title_pt: Counter = Counter()
    body_pt: Counter = Counter()
    colors: Counter = Counter()
    unreadable = 0
    size_in = (13.333, 7.5)

    for deck in decks:
        presentation = Presentation(str(deck))
        width_in = presentation.slide_width / EMU_PER_INCH
        height_in = presentation.slide_height / EMU_PER_INCH
        size_in = (round(width_in, 3), round(height_in, 3))
        for slide in presentation.slides:
            for shape in slide.shapes:
                try:
                    x = shape.left / EMU_PER_INCH / width_in
                    y = shape.top / EMU_PER_INCH / height_in
                    w = shape.width / EMU_PER_INCH / width_in
                    h = shape.height / EMU_PER_INCH / height_in
                except (TypeError, AttributeError, ZeroDivisionError):
                    unreadable += 1
                    continue
                if not (0.0 <= x <= 1.0 and 0.0 <= y <= 1.0 and w > 0 and h > 0):
                    continue
                if shape.shape_type is not None and "PICTURE" in str(shape.shape_type):
                    (content_images if w * h >= _CONTENT_IMAGE_MIN_AREA else glyphs).append((x, y, w, h))
                if not shape.has_text_frame or not shape.text_frame.text.strip():
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size is not None:
                            points = round(run.font.size.pt, 2)
                            (title_pt if (y < 0.22 and points >= 20) else body_pt)[points] += 1
                        try:
                            if run.font.color is not None and run.font.color.rgb is not None:
                                colors[str(run.font.color.rgb)] += 1
                        except (AttributeError, TypeError, ValueError):
                            pass  # theme-referenced colors carry no explicit RGB
                if y < 0.22 and h < 0.3:
                    titles.append((x, y, w, h))
                if y > 0.86:
                    footers.append((x, y, w, h))

    if not (titles and content_images and glyphs and footers):
        raise ValueError("corpus lacks one of title/content-image/glyph/footer samples")
    return HouseSpec(
        measured_from=[d.name for d in decks],
        slide_size_in=size_in,
        title_zone=_median_zone(titles),
        content_image_zone=_median_zone(content_images),
        inline_glyph_zone=_median_zone(glyphs),
        footer_zone=_median_zone(footers),
        title_pt=[pt for pt, _ in title_pt.most_common(4)],
        body_pt=[pt for pt, _ in body_pt.most_common(6)],
        accent_colors=[c for c, _ in colors.most_common(10)],
        unreadable_shapes=unreadable,
    )


# Committed measurement (2026-08-07) over the five-deck corpus at
# /mnt/storage12tb/skills/pitchdeck/sources/style-corpus. Regenerate with
# `./run.sh measure-house-spec --decks <dir>`; a diff here is a style change.
HOUSE_BODY_PT = 17.0          # corpus body modes 16.0/17.33 (was 22.0 — visibly airier than the house)
HOUSE_CAPTION_PT = 12.0       # corpus caption/fine mode 12.0
HOUSE_TITLE_PT = (24.0, 20.0)  # corpus title modes 20/24/28; adaptive by length
HOUSE_HERO_PT = 64.0          # corpus statement-slide mode 64.0 (confirmed, unchanged)
HOUSE_ACCENTS = ("#065E7C", "#26558E", "#1D7694", "#D6A300", "#4A7EBB")
HOUSE_CONTENT_IMAGE_W = 0.355  # median content-image width fraction
HOUSE_TITLE_CASE_RATIO = 0.81  # 173/213 real titles are Title Case
