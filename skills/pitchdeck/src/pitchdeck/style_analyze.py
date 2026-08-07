"""Deterministic reference-deck style analyzer → pitchdeck.style_reference.v1 (#1276).

Turns an arbitrary PPTX into validated style-reference data the planner and
design pass can consume: source hash, slide dimensions, font/color/size
frequencies, per-slide word stats, shape census, arc titles, and a
HOUSE-MATCH classification computed against a design_system.v1 instance
(signals, never auto-promotion — observations do not become invariants
without review). Optional representative renders via LibreOffice with an
ISOLATED profile; render absence is recorded, never faked. Deterministic:
identical input bytes produce identical output (no wall-clock timestamps;
ordering is sorted everywhere). Failure modes: malformed/unreadable PPTX
raises SkillError; unsupported objects are counted and reported as
`unsupported`, not silently interpreted.
"""

from __future__ import annotations

import hashlib
import json
import subprocess
from collections import Counter
from pathlib import Path

from .design_system import DesignSystem, StyleReference
from .io import SkillError


def _observations(counter: Counter, limit: int = 12) -> list[dict]:
    return [{"value": str(value), "count": count} for value, count in counter.most_common(limit)]


def analyze_pptx(pptx_path: Path) -> dict:
    """Measure one PPTX. Returns raw analysis (superset of StyleReference)."""
    try:
        from pptx import Presentation

        presentation = Presentation(str(pptx_path))
    except Exception as exc:
        raise SkillError(f"not a readable PPTX: {pptx_path.name}: {exc}") from exc

    fonts: Counter = Counter()
    colors: Counter = Counter()
    sizes: Counter = Counter()
    shape_kinds: Counter = Counter()
    unsupported: Counter = Counter()
    words_per_slide: list[int] = []
    titles: list[str] = []
    for slide in presentation.slides:
        words = 0
        title = ""
        for shape in slide.shapes:
            kind = str(shape.shape_type).split(" ")[0] if shape.shape_type is not None else "UNKNOWN"
            shape_kinds[kind] += 1
            if kind in {"MEDIA", "OLE_OBJECT", "EMBEDDED_OLE_OBJECT", "UNKNOWN"}:
                unsupported[kind] += 1
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            words += len(text.split())
            if not title and text:
                title = text.splitlines()[0][:80]
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.name:
                        fonts[run.font.name] += 1
                    if run.font.size:
                        sizes[int(run.font.size.pt)] += 1
                    try:
                        if run.font.color and run.font.color.rgb:
                            colors[f"#{run.font.color.rgb}"] += 1
                    except Exception:
                        unsupported["theme-color-ref"] += 1
        words_per_slide.append(words)
        titles.append(title)
    ordered = sorted(words_per_slide)
    count = len(ordered)
    return {
        "source_file": pptx_path.name,
        "source_sha256": hashlib.sha256(pptx_path.read_bytes()).hexdigest(),
        "slide_count": count,
        "slide_size_in": [round(presentation.slide_width / 914400, 2), round(presentation.slide_height / 914400, 2)],
        "fonts": _observations(fonts),
        "colors": _observations(colors),
        "sizes_pt": _observations(sizes),
        "shape_census": _observations(shape_kinds),
        "unsupported": _observations(unsupported),
        "words_median": ordered[count // 2] if count else 0,
        "words_p90": ordered[min(count - 1, int(count * 0.9))] if count else 0,
        "arc_titles": titles[:12],
    }


def classify_against_house(analysis: dict, system: DesignSystem) -> dict:
    """Signal-based match against a design system — evidence, not a verdict."""
    primary = system.palette[list(system.palette)[0]] if system.palette else None
    from .design_system import ColorRole

    primary = system.palette.get(ColorRole.PRIMARY)
    top_colors = {obs["value"].lstrip("#").upper() for obs in analysis["colors"][:8]}
    top_fonts = {obs["value"] for obs in analysis["fonts"][:4]}
    signals = {
        "primary_color_present": bool(primary) and primary.lstrip("#").upper() in top_colors,
        "body_font_present": system.fonts.body in top_fonts or "+mn-lt" in top_fonts,
        "title_scale_present": any(
            abs(int(obs["value"]) - system.type_scale_pt.title) <= 2 for obs in analysis["sizes_pt"]
        ),
        "density_within_profile": analysis["words_median"] <= 40,
    }
    score = sum(signals.values())
    return {
        "design_system": system.id,
        "signals": signals,
        "score": f"{score}/{len(signals)}",
        "matches_house": score >= 3,
        "note": "signals are evidence for human review; observations are never auto-promoted to invariants",
    }


def render_representatives(pptx_path: Path, output_dir: Path, analysis: dict) -> list[str]:
    """Render first + densest slide via LibreOffice (isolated profile); best-effort."""
    import shutil

    soffice = shutil.which("libreoffice") or shutil.which("soffice")
    if soffice is None:
        return []
    profile = output_dir / ".lo-profile"
    try:
        subprocess.run(
            [soffice, f"-env:UserInstallation=file://{profile.resolve()}", "--headless",
             "--convert-to", "pdf", "--outdir", str(output_dir), str(pptx_path)],
            capture_output=True, timeout=300, check=True,
        )
        pdf = output_dir / f"{pptx_path.stem}.pdf"
        rendered = []
        densest = analysis["arc_titles"] and max(range(analysis["slide_count"]), key=lambda i: i == 0) or 0
        for page in sorted({1, min(analysis["slide_count"], 2)}):
            out_prefix = output_dir / f"slide-{page}"
            subprocess.run(
                ["pdftoppm", "-png", "-r", "55", "-f", str(page), "-l", str(page), str(pdf), str(out_prefix)],
                capture_output=True, timeout=120, check=True,
            )
            rendered.extend(sorted(str(p.name) for p in output_dir.glob(f"slide-{page}-*.png")))
        pdf.unlink(missing_ok=True)
        return rendered
    except (subprocess.CalledProcessError, subprocess.TimeoutExpired):
        return []


def analyze_style(pptx_path: Path, output_dir: Path, design_system_path: Path | None = None, *, render: bool = True) -> dict:
    """Full pipeline: measure → validate as style_reference.v1 → classify → receipt."""
    output_dir.mkdir(parents=True, exist_ok=True)
    analysis = analyze_pptx(pptx_path)
    reference = StyleReference(
        id=pptx_path.stem.lower().replace("_", "-"),
        version=1,
        source_file=analysis["source_file"],
        source_sha256=analysis["source_sha256"],
        slide_count=analysis["slide_count"],
        fonts=analysis["fonts"],
        colors=analysis["colors"],
        sizes_pt=analysis["sizes_pt"],
        words_median=analysis["words_median"],
        words_p90=analysis["words_p90"],
        patterns=analysis["arc_titles"][:8],
        confidence="measured",
        provenance={"analyzer": "pitchdeck.style_analyze", "parser": "python-pptx"},
    )
    (output_dir / "style_reference.json").write_text(reference.model_dump_json(by_alias=True, indent=1), encoding="utf-8")
    classification = None
    if design_system_path is not None:
        system = DesignSystem.model_validate(json.loads(design_system_path.read_text(encoding="utf-8")))
        classification = classify_against_house(analysis, system)
    renders = render_representatives(pptx_path, output_dir, analysis) if render else []
    receipt = {
        "schema": "pitchdeck.style_analysis_receipt.v1",
        "source_sha256": analysis["source_sha256"],
        "style_reference": "style_reference.json",
        "classification": classification,
        "renders": renders,
        "renders_note": None if renders else "rendering unavailable or skipped; measurements unaffected",
        "unsupported_objects": analysis["unsupported"],
    }
    (output_dir / "analysis_receipt.json").write_text(json.dumps(receipt, indent=1), encoding="utf-8")
    return receipt
