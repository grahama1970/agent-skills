"""Deterministic house-conformance gate for emitted decks (#1311).

The author's decks are not stylistically approximate — they are invariant.
Measured across all 263 corpus slides: 100% carry a header band, a bottom-left
mark, bottom-right footer text, and a title; 261/263 bands are filled #076889;
band height is 0.108-0.122 of the slide; the median slide carries 2 pictures
and 8 shapes. That level of consistency means "does this look like the house
style?" is a MEASUREMENT, not a judgement call — so this module measures it
instead of asking a model.

Inputs: an emitted .pptx. Outputs: one ConformanceFinding per violated
invariant, each naming the corpus evidence it is derived from. Failure modes:
an unreadable deck raises; a slide whose geometry cannot be read is reported as
UNREADABLE rather than silently passing — a slide that cannot be measured must
never count as conformant.
"""

from __future__ import annotations

from typing import Literal

from pydantic import Field

from .models import StrictModel

EMU_PER_INCH = 914400.0

# Invariants measured over the corpus (see module docstring for the counts).
HOUSE_BAND_FILL = "076889"
HOUSE_BAND_HEIGHT_RANGE = (0.100, 0.130)
HOUSE_MEDIAN_PICTURES = 2
HOUSE_MEDIAN_SHAPES = 8


class ConformanceFinding(StrictModel):
    code: Literal[
        "HOUSE_BAND_MISSING",
        "HOUSE_BAND_FILL",
        "HOUSE_BAND_HEIGHT",
        "HOUSE_BOTTOM_LEFT_MARK_MISSING",
        "HOUSE_BOTTOM_RIGHT_TEXT_MISSING",
        "HOUSE_TITLE_MISSING",
        "HOUSE_VISUAL_DENSITY",
        "UNREADABLE",
    ]
    slide: int
    detail: str
    corpus_evidence: str


def check_conformance(pptx_path) -> list[ConformanceFinding]:
    """Measure an emitted deck against the corpus invariants."""
    from pptx import Presentation
    from pptx.util import Emu  # noqa: F401  (documents the EMU unit contract)

    presentation = Presentation(str(pptx_path))
    width_in = presentation.slide_width / EMU_PER_INCH
    height_in = presentation.slide_height / EMU_PER_INCH
    findings: list[ConformanceFinding] = []

    for index, slide in enumerate(presentation.slides, start=1):
        band_height: float | None = None
        band_fill: str | None = None
        bottom_left = bottom_right = has_title = False
        visuals = 0
        unreadable = 0
        own_shapes = list(slide.shapes)
        # Chrome may be INHERITED: the corpus keeps band/marks/footer on the
        # slide layout, so a slide-only scan fails the author's own decks (it
        # did — 233 false findings on ACERT before this fix). Measure the
        # composed slide: its own shapes plus its layout's.
        for shape in list(slide.shapes) + list(slide.slide_layout.shapes):
            try:
                x = shape.left / EMU_PER_INCH / width_in
                y = shape.top / EMU_PER_INCH / height_in
                w = shape.width / EMU_PER_INCH / width_in
                h = shape.height / EMU_PER_INCH / height_in
            except (TypeError, AttributeError, ZeroDivisionError):
                unreadable += 1
                continue
            if shape in own_shapes and shape.shape_type is not None and any(
                token in str(shape.shape_type) for token in ("PICTURE", "GROUP", "FREEFORM")
            ):
                visuals += 1  # density counts AUTHORED visuals, not inherited chrome
            if w > 0.9 and y < 0.03 and 0.04 < h < 0.22 and band_height is None:
                band_height = h
                try:
                    if shape.fill.type == 1:
                        band_fill = str(shape.fill.fore_color.rgb)
                except (AttributeError, TypeError, ValueError):
                    band_fill = None
            # The corpus bottom-left mark is a LOGO ROW (pictures), not text —
            # requiring text failed 170 of the author's own slides. Either a
            # picture mark or a text wordmark satisfies the invariant.
            has_text = shape.has_text_frame and shape.text_frame.text.strip()
            is_mark = shape.shape_type is not None and any(
                token in str(shape.shape_type) for token in ("PICTURE", "GROUP"))
            if y > 0.85 and x < 0.35 and (has_text or is_mark):
                bottom_left = True
            if y > 0.85 and x > 0.6 and shape.has_text_frame and shape.text_frame.text.strip():
                bottom_right = True
            if y < 0.22 and shape.has_text_frame and shape.text_frame.text.strip():
                has_title = True

        if unreadable:
            findings.append(ConformanceFinding(
                code="UNREADABLE", slide=index,
                detail=f"{unreadable} shape(s) had unreadable geometry",
                corpus_evidence="a slide that cannot be measured cannot be certified conformant"))
        if band_height is None:
            findings.append(ConformanceFinding(
                code="HOUSE_BAND_MISSING", slide=index, detail="no full-width header band",
                corpus_evidence="262/263 corpus slides carry a header band"))
        else:
            if not (HOUSE_BAND_HEIGHT_RANGE[0] <= band_height <= HOUSE_BAND_HEIGHT_RANGE[1]):
                findings.append(ConformanceFinding(
                    code="HOUSE_BAND_HEIGHT", slide=index,
                    detail=f"band height {band_height:.3f} outside {HOUSE_BAND_HEIGHT_RANGE}",
                    corpus_evidence="corpus band heights are 0.108 and 0.122"))
            if band_fill is not None and band_fill.upper() != HOUSE_BAND_FILL:
                findings.append(ConformanceFinding(
                    code="HOUSE_BAND_FILL", slide=index,
                    detail=f"band fill #{band_fill} != house #{HOUSE_BAND_FILL}",
                    corpus_evidence="261/263 corpus bands are filled #076889"))
        if not bottom_left:
            findings.append(ConformanceFinding(
                code="HOUSE_BOTTOM_LEFT_MARK_MISSING", slide=index, detail="no bottom-left identity mark",
                corpus_evidence="262/263 corpus slides carry a bottom-left mark"))
        if not bottom_right:
            findings.append(ConformanceFinding(
                code="HOUSE_BOTTOM_RIGHT_TEXT_MISSING", slide=index, detail="no bottom-right footer text",
                corpus_evidence="262/263 corpus slides carry bottom-right footer text"))
        if not has_title:
            findings.append(ConformanceFinding(
                code="HOUSE_TITLE_MISSING", slide=index, detail="no title text in the top band region",
                corpus_evidence="262/263 corpus slides carry a title"))
        if visuals < HOUSE_MEDIAN_PICTURES:
            findings.append(ConformanceFinding(
                code="HOUSE_VISUAL_DENSITY", slide=index,
                detail=f"{visuals} visual object(s) < house median {HOUSE_MEDIAN_PICTURES}",
                corpus_evidence="corpus median is 2 pictures and 8 shapes per slide"))
    return findings
